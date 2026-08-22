"""
program_pptx.py

Builds the delivery program as a standalone PowerPoint, straight from a
project's program_schedule + program_week_labels -- no template file to keep
in sync (same from-scratch approach as org_chart_pptx.py /
methodology_pptx.py, for the same reasons).

FOUR STYLES, ONE MODEL
----------------------
The program is drawn in whichever of the four presentation styles the user
picked (program_render.STYLES). The shapes here are built from the SAME
program_render.build_model() object the on-screen preview and the letter
pack render from, so the deck cannot show a different program from the one
the user approved -- which is exactly what happened while each output
derived its own view of the ticked weeks.

Everything is native PowerPoint shapes (and, for the formal-table style, a
real PowerPoint table), not a pasted picture: the whole point of the
companion deck is that it can be tidied up further in PowerPoint.

Every row and every bar comes straight from program_schedule.py's output or
the user's own edits to it (see the Fee Estimate tab's "Delivery program"
editor) -- nothing here is invented; an empty schedule renders an explicit
red placeholder slide instead of a blank grid, same no-invention convention
as the rest of this tool.
"""

from __future__ import annotations

import io
import math

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from modules import export_i18n, program_render
from modules.divider_designer import THEME_COLOURS

_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
_LIGHT_GREY = RGBColor(0xD9, 0xD9, 0xD9)
_RED = RGBColor(0xC0, 0x00, 0x00)
_FONT = "Calibri"

_INK = RGBColor(0x1A, 0x22, 0x33)
_MUTED = RGBColor(0x7A, 0x85, 0x98)
_GRIDLINE = RGBColor(0xEE, 0xF1, 0xF4)
_ROW_BAND = RGBColor(0xF7, 0xF9, 0xFC)
_TRACK_GREY = RGBColor(0xED, 0xEF, 0xF2)
_MILESTONE_ORANGE = RGBColor(0xF9, 0x73, 0x16)

# A4 landscape -- same slide size as org_chart_pptx.py / methodology_pptx.py
_SLIDE_W = Inches(11.6929)
_SLIDE_H = Inches(8.2677)
_M = Inches(0.3)

# Superseded by export_i18n's "pptx_program_empty_note" key (Audit Round 2,
# Part 5, which made this note respect output_language) -- kept here as a
# fixed English constant only for backward compatibility with any external
# caller that still imports this name directly; _placeholder_slide() itself
# no longer reads it.
EMPTY_NOTE = ("[NO PROGRAM ENTERED -- build the delivery program in the Fee Estimate tab, then "
              "re-download this PowerPoint]")


def _resolve_palette(theme_name: str | None) -> dict:
    colours = THEME_COLOURS.get(theme_name, THEME_COLOURS["Corporate"])
    primary = RGBColor(*colours["primary"])
    accent = RGBColor(*colours["accent"])
    # Minimalist's "primary" is a light near-white wash by design -- same special
    # case export_docx._theme_colours() / methodology_pptx._resolve_palette() apply.
    dark_role = RGBColor(0x2A, 0x2A, 0x2A) if theme_name == "Minimalist" else primary
    return {"header_bg": dark_role, "active_bg": accent}


def _hex(value: str) -> RGBColor:
    return RGBColor.from_string(str(value).lstrip("#").upper())


def _tint(colour: RGBColor, towards_white: float) -> RGBColor:
    """`colour` mixed `towards_white` of the way to white -- the lane washes."""
    r, g, b = colour[0], colour[1], colour[2]
    mix = lambda c: int(round(c + (255 - c) * towards_white))  # noqa: E731
    return RGBColor(mix(r), mix(g), mix(b))


def _rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


# Average glyph advance for bold Calibri over mixed-case text, as a fraction
# of the point size. An estimate, because there is no renderer here to
# measure with -- unlike the preview, which measures. Deliberately on the
# generous side: a label that shrinks half a point more than it had to is
# invisible, a label that runs out past the end of its bar is not.
_AVG_CHAR_EM = 0.50


def _fit_label(label: str, inner_pt: float, size_pt: float, lines: int = 1,
               min_pt: float = 7.0) -> tuple[str, float]:
    """Shrink a label to fit its shape, then ellipsize if it still will not.

    Shrinking first keeps the whole label readable; ellipsizing is the last
    resort, and an ellipsis is visibly an ellipsis rather than text that
    silently runs off the end of a bar."""
    capacity = max(1.0, inner_pt) * max(1, lines)
    while size_pt > min_pt and len(label) * _AVG_CHAR_EM * size_pt > capacity:
        size_pt -= 0.5
    max_chars = int(capacity / max(0.1, _AVG_CHAR_EM * size_pt))
    if len(label) > max_chars:
        label = label[:max(1, max_chars - 1)].rstrip() + "…"
    return label, size_pt


def _inner_points(width_emu: float, margin_pt: float = 6.0) -> float:
    return width_emu / 12700.0 - 2 * margin_pt


def _pill(slide, x, y, w, h, fill_color, label: str = "", size_pt: float = 9.0,
          align=PP_ALIGN.CENTER):
    """A fully-rounded bar. ROUNDED_RECTANGLE with its adjustment pushed to
    the maximum, which is what makes the ends read as round rather than as a
    rounded rectangle."""
    w = max(int(w), int(h))
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Emu(int(w)), h)
    try:
        shape.adjustments[0] = 0.5
    except (IndexError, KeyError):
        pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    if label:
        # The label lives INSIDE its bar, and a PowerPoint shape happily
        # paints text straight out past its own edge.
        label, size_pt = _fit_label(label, _inner_points(w), size_pt)
        frame = shape.text_frame
        frame.word_wrap = False
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        frame.margin_left = Pt(6)
        frame.margin_right = Pt(6)
        frame.margin_top = Pt(0)
        frame.margin_bottom = Pt(0)
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = label
        run.font.size = Pt(size_pt)
        run.font.bold = True
        run.font.name = _FONT
        run.font.color.rgb = _WHITE
    return shape


def _line(slide, x, y, w, h, colour):
    """A hairline drawn as a filled rect -- a connector would be one more
    shape type for the user to fight with when they edit the slide."""
    return _rect(slide, Emu(int(x)), Emu(int(y)), Emu(max(1, int(w))), Emu(max(1, int(h))), colour)


def _text(slide, x, y, w, h, text, size_pt, color, bold=False, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.name = _FONT
    run.font.color.rgb = color
    return box


def _diamond(slide, cx, cy, size, colour):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND, Emu(int(cx - size / 2)), Emu(int(cy - size / 2)),
        Emu(int(size)), Emu(int(size)),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _no_borders(cell) -> None:
    """Strip a table cell's four borders.

    python-pptx has no border API, so this writes the a:ln* children
    directly. Without it the "formal table" style arrives in PowerPoint
    wearing the default banded-blue table style -- which is precisely the
    look the formal-table option exists to avoid."""
    tc_pr = cell._tc.get_or_add_tcPr()
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        for existing in tc_pr.findall(qn(tag)):
            tc_pr.remove(existing)
        element = tc_pr.makeelement(qn(tag), {"w": "0", "cap": "flat", "cmpd": "sng", "algn": "ctr"})
        element.append(element.makeelement(qn("a:noFill"), {}))
        tc_pr.insert(0, element)


# ---------------------------------------------------------------------------
# Shared slide furniture
# ---------------------------------------------------------------------------

def _new_deck():
    prs = Presentation()
    prs.slide_width = _SLIDE_W
    prs.slide_height = _SLIDE_H
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def _save(prs) -> bytes:
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()


def _heading(slide, model, language: str = "en") -> int:
    """Title + subtitle. Returns the y the content may start at -- a FIXED
    value (only _M and the two Inches() constants below, none of which
    depend on the slide's height), so it can be used as CONTENT_TOP_EMU
    below without calling this first."""
    title_h = Inches(0.5)
    _text(slide, _M, _M, Emu(int(_SLIDE_W - 2 * _M)), title_h,
          export_i18n.export_t("pptx_program_title", language), 20, _DARK_TEXT, bold=True)
    subtitle = (model.project_name or "").strip() or export_i18n.export_t("pptx_insert_project_name", language)
    if (model.client_name or "").strip():
        subtitle = f"{subtitle} — {model.client_name.strip()}"
    _text(slide, _M, Emu(int(_M + title_h)), Emu(int(_SLIDE_W - 2 * _M)), Inches(0.3), subtitle, 11,
          _DARK_TEXT if (model.project_name or "").strip() else _RED)
    return int(_M + title_h + Inches(0.45))


def _grow(prs, needed_height: float) -> None:
    """A defensive last resort only -- see org_chart_pptx._grow(). Every
    style below sizes itself to the fixed A4 slide via _fit_height()'s
    scale-to-fill search first; this only fires if that search's own
    MIN_SCALE floor still wasn't enough (an extreme scope-item/week count),
    growing the slide rather than leaving shapes silently positioned past
    its bottom edge, which is what PowerPoint does with them."""
    if int(needed_height) > int(prs.slide_height):
        prs.slide_height = Emu(int(needed_height))


# ---------------------------------------------------------------------------
# SCALE-TO-FILL, LIKE THE ORG CHART AND THE PNG RENDERER
# ---------------------------------------------------------------------------
# This used to draw every row/lane at a fixed size and grow the slide (via
# _grow, above) only once that fixed sizing overflowed it -- so a short
# program's deck was a small grid pasted at the top of an otherwise-empty
# A4 slide, and _grow() never fired to fix that because nothing overflowed.
# The slide is now the fixed A4 landscape page every export in this pack
# uses (_SLIDE_W / _SLIDE_H); each style instead measures its content at a
# reference scale, picks ONE scale factor from how much of that fixed page
# it needs, and derives row height, bar thickness, gridline weight,
# milestone size and every font from that one number -- the same approach
# program_render.py takes for the PNG (see that module's matching note),
# now shared by the companion deck so the two outputs read as the same
# document. _grow() above stays only as the same defensive floor
# program_render._fit_height() already has: past MIN_SCALE the page grows
# instead of shrinking text further.

CONTENT_LEFT_EMU = int(_M)
CONTENT_RIGHT_EMU = int(_SLIDE_W - _M)
AVAIL_W_EMU = CONTENT_RIGHT_EMU - CONTENT_LEFT_EMU

# Matches _heading()'s fixed return value exactly (same _M / Inches(0.5) /
# Inches(0.45) constants) -- needed before a slide exists, to size the
# scale-to-fill search below.
CONTENT_TOP_EMU = int(_M + Inches(0.5) + Inches(0.45))
CONTENT_BOTTOM_EMU = int(_SLIDE_H - _M)
AVAIL_H_EMU = CONTENT_BOTTOM_EMU - CONTENT_TOP_EMU

MIN_SCALE = 0.55
MAX_SCALE = 1.25


def _fit_height(natural_h_emu: float) -> tuple[float, int, int]:
    """(scale, total_slide_h_emu, avail_h_emu) for content whose natural
    (scale=1.0) height is `natural_h_emu` -- EMU counterpart to
    program_render._fit_height(); see that function's docstring for the
    ordinary/overflow split, identical here."""
    if natural_h_emu <= 0:
        return MAX_SCALE, int(_SLIDE_H), AVAIL_H_EMU
    scale = max(MIN_SCALE, min(MAX_SCALE, AVAIL_H_EMU / natural_h_emu))
    needed = natural_h_emu * scale
    if needed <= AVAIL_H_EMU + 1:
        return scale, int(_SLIDE_H), AVAIL_H_EMU
    overflow = needed - AVAIL_H_EMU
    return scale, int(_SLIDE_H) + int(overflow), AVAIL_H_EMU + int(overflow)


class _PVFlow:
    """EMU-space vertical flow -- see org_chart_pptx._PFlow for the shared
    rationale (fixed-size blocks; gaps that stretch to absorb leftover slide
    height, so a short program's leftover room becomes breathing space
    between sections rather than one dead band) and program_render._VFlow
    for the identical idea in matplotlib inches. The third "fixed" kind is
    for content whose height is already final and must NOT be multiplied by
    scale again -- the table style's row-height cap, the one place here
    where "how tall" isn't a straight scale multiply."""

    def __init__(self):
        self._items: list[dict] = []

    def block(self, h_emu, draw=None) -> None:
        self._items.append({"kind": "block", "h": max(0, int(h_emu)), "draw": draw})

    def gap(self, h_emu, draw=None) -> None:
        self._items.append({"kind": "gap", "h": max(0, int(h_emu)), "draw": draw})

    def fixed(self, h_emu, draw=None) -> None:
        self._items.append({"kind": "fixed", "h": max(0, int(h_emu)), "draw": draw})

    def render(self, top_emu: int, scale: float, avail_emu: int) -> int:
        n_gaps = sum(1 for it in self._items if it["kind"] == "gap")
        used = sum(it["h"] if it["kind"] == "fixed" else int(it["h"] * scale) for it in self._items)
        leftover = max(0, avail_emu - used)
        extra_per_gap = (leftover // n_gaps) if n_gaps else 0
        y = top_emu
        for it in self._items:
            h = it["h"] if it["kind"] == "fixed" else int(it["h"] * scale) + \
                (extra_per_gap if it["kind"] == "gap" else 0)
            y_bottom = y + h
            if it["draw"] is not None:
                it["draw"](y, y_bottom, scale)
            y = y_bottom
        return y


# Font sizes at scale=1.0, and the floor each is allowed to shrink to --
# same numbers as program_render.py's matplotlib renderer (see that
# module's matching constants block) so the PNG preview/letter-pack embed
# and this companion deck read as the same document at the same program
# size, not two different type scales.
_ROW_LABEL_PT_REF, _ROW_LABEL_PT_MIN = 12.5, 7.0
_BAR_LABEL_PT_REF, _BAR_LABEL_PT_MIN = 9.5, 6.0
_WEEK_PT_REF, _WEEK_PT_MIN = 11.0, 6.0
_WEEK_DATE_PT_REF, _WEEK_DATE_PT_MIN = 8.8, 5.5
_MILESTONE_PT_REF, _MILESTONE_PT_MIN = 9.5, 6.0
_LEGEND_PT_REF, _LEGEND_PT_MIN = 9.5, 6.5
_LANE_PT_REF, _LANE_PT_MIN = 10.5, 6.5

# The narrowest a "Wk NN" header can sit next to its neighbour, in inches at
# scale=1.0, before the two touch -- see program_render._HEADER_MIN_PITCH_IN,
# identical logic here.
_HEADER_MIN_PITCH_IN = 0.50


def _header_indices(model, week_col_w_emu: float, scale: float = 1.0, pt: float | None = None) -> list[int]:
    """See program_render._header_indices() for the full rationale --
    identical here except the "measure the widest label" step uses the
    same avg-glyph-width estimate _fit_label()/_legend() already use in
    this file, since there is no real renderer to measure with."""
    count = len(model.week_labels)
    if count <= 1:
        return list(range(count))
    pitch_emu = int(Inches(_HEADER_MIN_PITCH_IN) * max(scale, 1e-6))
    if pt is not None and model.week_labels:
        widest = max(model.week_labels, key=len)
        label_w_emu = int(len(widest) * _AVG_CHAR_EM * pt * 12700 * 1.15)
        pitch_emu = max(pitch_emu, label_w_emu)
    step = max(1, math.ceil(pitch_emu / max(int(week_col_w_emu), 1)))
    kept = list(range(0, count, step))
    last = count - 1
    if kept[-1] != last:
        while kept and last - kept[-1] < step:
            kept.pop()
        kept.append(last)
    return kept


def _week_header(slide, model, grid_left, y, week_col_w, scale, show_dates=True):
    pt = max(_WEEK_PT_MIN, _WEEK_PT_REF * scale)
    date_pt = max(_WEEK_DATE_PT_MIN, _WEEK_DATE_PT_REF * scale)
    label_h = Inches(0.26) * scale
    date_h = Inches(0.20) * scale
    for index in _header_indices(model, week_col_w, scale, pt=pt):
        cx = Emu(int(grid_left + index * week_col_w))
        _text(slide, cx, Emu(int(y)), Emu(int(week_col_w)), Emu(int(label_h)),
              model.week_labels[index], pt, _INK, bold=True, align=PP_ALIGN.CENTER)
        date_text = model.week_dates[index] if index < len(model.week_dates) else ""
        if show_dates and date_text:
            _text(slide, cx, Emu(int(y + label_h * 0.9)), Emu(int(week_col_w)), Emu(int(date_h)),
                  date_text, date_pt, _MUTED, align=PP_ALIGN.CENTER)


def _gridlines(slide, model, grid_left, grid_top, grid_bottom, week_col_w, scale):
    weight = Inches(0.006) * max(0.7, scale)
    for index in range(len(model.week_labels) + 1):
        _line(slide, grid_left + index * week_col_w, grid_top, weight,
              grid_bottom - grid_top, _GRIDLINE)


def _milestones(slide, model, grid_left, grid_right, grid_top, grid_bottom, label_y, label_h, scale):
    """Diamonds on a faint vertical rule at grid_bottom, labels in the
    reserved band [label_y, label_y + label_h) below -- the two are drawn
    apart (rather than immediately under the diamond) because grid_bottom
    is the end of the ROWS block and label_y is the start of the dedicated
    milestone block one gap further down; see the gantt/swimlanes/timeline
    renderers below for how those two y's are produced by the flow."""
    if not model.milestones:
        return
    size = Inches(0.16) * max(0.7, scale)
    pt = max(_MILESTONE_PT_MIN, _MILESTONE_PT_REF * scale)
    weeks = max(1, len(model.week_labels))
    week_col_w = (grid_right - grid_left) / weeks
    seen = set()
    for milestone in model.milestones:
        week = max(1, min(int(milestone.week), weeks))
        if week in seen:
            continue
        seen.add(week)
        cx = grid_left + week * week_col_w
        _line(slide, cx - Inches(0.004), grid_top, Inches(0.008), grid_bottom - grid_top,
              _tint(_MILESTONE_ORANGE, 0.45))
        _diamond(slide, cx, grid_bottom, size, _MILESTONE_ORANGE)
        # A milestone in the last week sits on the slide's right edge, so a
        # centred label runs off it -- the label swings inboard instead.
        width = Inches(1.6)
        label_x, align = cx - width / 2, PP_ALIGN.CENTER
        if label_x + width > CONTENT_RIGHT_EMU:
            label_x, align = CONTENT_RIGHT_EMU - width, PP_ALIGN.RIGHT
        elif label_x < CONTENT_LEFT_EMU:
            label_x, align = CONTENT_LEFT_EMU, PP_ALIGN.LEFT
        _text(slide, Emu(int(label_x)), Emu(int(label_y)), Emu(int(width)), Emu(int(label_h)),
              milestone.label, pt, _MILESTONE_ORANGE, bold=True, align=align)


def _activity_legend(model, accent: RGBColor, language: str = "en") -> list[tuple[RGBColor, str]]:
    """See program_render._activity_legend(): no key for a mark that isn't on
    the chart."""
    entries = [(accent, export_i18n.export_t("pptx_program_legend_scheduled", language))]
    if model.milestones:
        entries.append((_MILESTONE_ORANGE, export_i18n.export_t("pptx_program_legend_milestone", language)))
    return entries


def _legend(slide, entries: list[tuple[RGBColor, str]], y: int, x: int, scale: float) -> None:
    cursor = int(x)
    swatch = Inches(0.14) * scale
    pt = max(_LEGEND_PT_MIN, _LEGEND_PT_REF * scale)
    label_h = Inches(0.2) * scale
    for colour, label in entries:
        # A diamond for the milestone key, so it matches the marks on the
        # chart -- it shares its orange with the third stage colour, and
        # shape is the only thing telling the two apart.
        if label.lower().startswith("milestone"):
            _diamond(slide, cursor + swatch / 2, y + Inches(0.03) * scale + swatch / 2, swatch,
                     _MILESTONE_ORANGE)
        else:
            _rect(slide, Emu(cursor), Emu(int(y + Inches(0.03) * scale)), Emu(int(swatch)),
                  Emu(int(swatch)), colour)
        text_x = cursor + int(swatch * 1.5)
        # Estimated from the ACTUAL point size the label draws at, via the
        # same avg-glyph-width model _fit_label() already uses for bar
        # labels in this file -- not a flat per-character inch guess that
        # ignores the font size entirely. Legend text is floored at
        # _LEGEND_PT_MIN rather than shrinking all the way with scale, so
        # near MIN_SCALE a flat guess underestimated the true (larger,
        # floored) text width and let the next entry's swatch -- at the
        # far end, the milestone diamond -- crowd into this label; a fixed
        # guess also overshot at MAX_SCALE's smaller relative text. A small
        # +30% safety margin (generous, since the box's own left/right text
        # margins eat into it too) covers the estimate's own error, and
        # wrap=False is the backstop if it's still wrong: a legend entry
        # running a little wide reads far better than one word-wrapped
        # into "Documentaio / n" inside a single-line-tall box, which is
        # what a wrapped overflow looks like here.
        width_emu = int(len(label) * _AVG_CHAR_EM * pt * 12700 * 1.30)
        _text(slide, Emu(text_x), Emu(int(y)), Emu(width_emu), Emu(int(label_h)),
              label, pt, _INK, bold=True, wrap=False)
        cursor = text_x + width_emu + int(Inches(0.14) * scale)


def _placeholder_slide(model, language: str = "en") -> bytes:
    prs, slide = _new_deck()
    top = _heading(slide, model, language)
    _text(slide, _M, Emu(int(top + Inches(0.3))), Emu(int(_SLIDE_W - 2 * _M)), Inches(0.6),
          export_i18n.export_t("pptx_program_empty_note", language), 14, _RED, align=PP_ALIGN.CENTER)
    return _save(prs)


# ---------------------------------------------------------------------------
# A. Refined Gantt
# ---------------------------------------------------------------------------

_GANTT_LABEL_COL_REF_IN = 2.5
_GANTT_HEADER_H_REF = 0.46
_GANTT_ROW_H_REF = 0.80
_GANTT_GAP_REF = 0.14
_GANTT_MILESTONE_H_REF = 0.42
_GANTT_LEGEND_H_REF = 0.34


def _slide_gantt(model, accent: RGBColor, language: str = "en") -> bytes:
    n = len(model.items)
    has_ms = bool(model.milestones)
    natural_h_in = (_GANTT_HEADER_H_REF + _GANTT_GAP_REF
                   + n * _GANTT_ROW_H_REF + _GANTT_GAP_REF
                   + (_GANTT_MILESTONE_H_REF + _GANTT_GAP_REF if has_ms else 0.0)
                   + _GANTT_LEGEND_H_REF)
    scale, total_h_emu, avail_h_emu = _fit_height(int(Inches(natural_h_in)))

    prs, slide = _new_deck()
    top = _heading(slide, model, language)

    label_col_w = int(Inches(_GANTT_LABEL_COL_REF_IN) * scale)
    grid_left = CONTENT_LEFT_EMU + label_col_w
    grid_right = CONTENT_RIGHT_EMU
    weeks = max(1, len(model.week_labels))
    week_col_w = (grid_right - grid_left) / weeks

    bounds: dict[str, int] = {}
    flow = _PVFlow()

    def draw_header(y_top, y_bottom, scale):
        _week_header(slide, model, grid_left, y_top, week_col_w, scale)
        bounds["grid_top"] = y_bottom
    flow.block(Inches(_GANTT_HEADER_H_REF), draw_header)
    flow.gap(Inches(_GANTT_GAP_REF))

    # Row bands, then gridlines, then bars+labels are all drawn AFTER
    # flow.render() below rather than inline here, in insertion order --
    # PowerPoint paints shapes in the order they're added, and gridlines
    # added ahead of the bars they cross would rule white lines straight
    # over them; added behind the row bands they'd be buried instead. Each
    # per-row draw() closure below only records the label text and pill; the
    # band/gridline passes run once every row position is known.
    row_draws = []
    for index, item in enumerate(model.items):
        def draw_row(y_top, y_bottom, scale, index=index, item=item):
            row_draws.append((index, item, y_top, y_bottom))
            bounds["grid_bottom"] = y_bottom
        flow.block(Inches(_GANTT_ROW_H_REF), draw_row)

    flow.gap(Inches(_GANTT_GAP_REF))

    if has_ms:
        def draw_ms_anchor(y_top, y_bottom, scale):
            bounds["ms_y"] = y_top
            bounds["ms_h"] = y_bottom - y_top
        flow.block(Inches(_GANTT_MILESTONE_H_REF), draw_ms_anchor)
        flow.gap(Inches(_GANTT_GAP_REF))

    def draw_legend_anchor(y_top, y_bottom, scale):
        bounds["legend_y"] = y_top
    flow.block(Inches(_GANTT_LEGEND_H_REF), draw_legend_anchor)

    flow.render(top, scale, avail_h_emu)

    for index, item, y_top, y_bottom in row_draws:
        if index % 2 == 1:
            _rect(slide, Emu(CONTENT_LEFT_EMU), Emu(int(y_top)),
                  Emu(int(grid_right - CONTENT_LEFT_EMU)), Emu(int(y_bottom - y_top)), _ROW_BAND)
    _gridlines(slide, model, grid_left, bounds["grid_top"], bounds["grid_bottom"], week_col_w, scale)
    for index, item, y_top, y_bottom in row_draws:
        row_h = y_bottom - y_top
        pt = max(_ROW_LABEL_PT_MIN, _ROW_LABEL_PT_REF * scale)
        _text(slide, Emu(CONTENT_LEFT_EMU), Emu(int(y_top)), Emu(int(label_col_w)), Emu(int(row_h)),
              item.label or export_i18n.export_t("export_untitled_scope_item", language), pt, _INK, bold=True)
        x0 = grid_left + (item.start_week - 1) * week_col_w
        x1 = grid_left + item.end_week * week_col_w
        bar_h = int(row_h * 0.45)
        bar_pt = max(_BAR_LABEL_PT_MIN, _BAR_LABEL_PT_REF * scale)
        _pill(slide, Emu(int(x0 + Inches(0.02))), Emu(int(y_top + (row_h - bar_h) / 2)),
              Emu(int(x1 - x0 - Inches(0.04))), Emu(bar_h), accent,
              export_i18n.export_t("pptx_duration_weeks_short", language, weeks=item.weeks), bar_pt)

    if has_ms:
        _milestones(slide, model, grid_left, grid_right, bounds["grid_top"], bounds["grid_bottom"],
                    bounds["ms_y"], bounds["ms_h"], scale)
    _legend(slide, _activity_legend(model, accent, language), bounds["legend_y"], CONTENT_LEFT_EMU, scale)

    _grow(prs, total_h_emu)
    return _save(prs)


# ---------------------------------------------------------------------------
# B. Stage swimlanes
# ---------------------------------------------------------------------------

def _grouped_by_stage(model):
    grouped = []
    for stage_index in list(range(len(model.stages))) + [None]:
        members = [i for i in model.items if i.stage_index == stage_index]
        if members:
            grouped.append((stage_index, members))
    return grouped


_SWIM_LANE_HEADER_H_REF = 0.34
_SWIM_ROW_H_REF = 0.62
_SWIM_GAP_REF = 0.10
_SWIM_LABEL_COL_REF_IN = 2.5
_SWIM_HEADER_H_REF = 0.46
_SWIM_MILESTONE_H_REF = 0.42
_SWIM_LEGEND_H_REF = 0.34


def _slide_swimlanes(model, accent: RGBColor, language: str = "en") -> bytes:
    grouped = _grouped_by_stage(model)
    total_rows = sum(len(members) for _, members in grouped)
    has_ms = bool(model.milestones)
    natural_h_in = (_SWIM_HEADER_H_REF + _SWIM_GAP_REF
                   + len(grouped) * _SWIM_LANE_HEADER_H_REF + total_rows * _SWIM_ROW_H_REF
                   + max(0, len(grouped) - 1) * _SWIM_GAP_REF
                   + _SWIM_GAP_REF
                   + (_SWIM_MILESTONE_H_REF + _SWIM_GAP_REF if has_ms else 0.0)
                   + _SWIM_LEGEND_H_REF)
    scale, total_h_emu, avail_h_emu = _fit_height(int(Inches(natural_h_in)))

    prs, slide = _new_deck()
    top = _heading(slide, model, language)

    label_col_w = int(Inches(_SWIM_LABEL_COL_REF_IN) * scale)
    grid_left = CONTENT_LEFT_EMU + label_col_w
    grid_right = CONTENT_RIGHT_EMU
    weeks = max(1, len(model.week_labels))
    week_col_w = (grid_right - grid_left) / weeks

    bounds: dict[str, int] = {}
    flow = _PVFlow()

    def draw_header(y_top, y_bottom, scale):
        _week_header(slide, model, grid_left, y_top, week_col_w, scale)
        bounds["grid_top"] = y_bottom
    flow.block(Inches(_SWIM_HEADER_H_REF), draw_header)
    flow.gap(Inches(_SWIM_GAP_REF))

    # Same three-pass ordering as _slide_gantt(): lane tint washes, then
    # gridlines, then bars+labels on top -- recorded here, drawn after
    # flow.render() once every lane/row position is known.
    lane_draws, row_draws = [], []
    for lane_index, (stage_index, members) in enumerate(grouped):
        colour = (_hex(program_render.STAGE_COLOURS[stage_index % len(program_render.STAGE_COLOURS)])
                  if stage_index is not None else _MUTED)

        def draw_lane_header(y_top, y_bottom, scale, stage_index=stage_index, colour=colour,
                             members=members):
            # The tint spans the WHOLE lane (header + its member rows), so
            # it's computed here from the member count even though this
            # callback only owns the header's own band -- see
            # program_render._render_swimlanes's identical note.
            member_h = (y_bottom - y_top) + len(members) * int(Inches(_SWIM_ROW_H_REF) * scale)
            lane_draws.append((y_top, member_h, colour, stage_index))
        flow.block(Inches(_SWIM_LANE_HEADER_H_REF), draw_lane_header)

        for item in members:
            def draw_row(y_top, y_bottom, scale, item=item, colour=colour):
                row_draws.append((item, colour, y_top, y_bottom))
                bounds["grid_bottom"] = y_bottom
            flow.block(Inches(_SWIM_ROW_H_REF), draw_row)

        if lane_index < len(grouped) - 1:
            flow.gap(Inches(_SWIM_GAP_REF))

    flow.gap(Inches(_SWIM_GAP_REF))

    if has_ms:
        def draw_ms_anchor(y_top, y_bottom, scale):
            bounds["ms_y"] = y_top
            bounds["ms_h"] = y_bottom - y_top
        flow.block(Inches(_SWIM_MILESTONE_H_REF), draw_ms_anchor)
        flow.gap(Inches(_SWIM_GAP_REF))

    def draw_legend_anchor(y_top, y_bottom, scale):
        bounds["legend_y"] = y_top
    flow.block(Inches(_SWIM_LEGEND_H_REF), draw_legend_anchor)

    flow.render(top, scale, avail_h_emu)

    for y_top, member_h, colour, _stage_index in lane_draws:
        # ~5% tint: enough to group the rows, never enough to fight the bars
        # sitting on it.
        _rect(slide, Emu(CONTENT_LEFT_EMU), Emu(int(y_top)), Emu(int(grid_right - CONTENT_LEFT_EMU)),
              Emu(int(member_h)), _tint(colour, 0.95))
    _gridlines(slide, model, grid_left, bounds["grid_top"], bounds["grid_bottom"], week_col_w, scale)

    for y_top, _member_h, colour, stage_index in lane_draws:
        name = (model.stages[stage_index] if stage_index is not None
               else export_i18n.export_t("pptx_unassigned_stage_label", language)).upper()
        pt = max(_LANE_PT_MIN, _LANE_PT_REF * scale)
        header_h = int(Inches(_SWIM_LANE_HEADER_H_REF) * scale)
        _text(slide, Emu(CONTENT_LEFT_EMU), Emu(int(y_top)), Emu(int(grid_right - CONTENT_LEFT_EMU)),
              Emu(header_h), name, pt, colour, bold=True)
    for item, colour, y_top, y_bottom in row_draws:
        row_h = y_bottom - y_top
        pt = max(_ROW_LABEL_PT_MIN, _ROW_LABEL_PT_REF * 0.92 * scale)
        inset = int(Inches(0.10) * scale)
        _text(slide, Emu(CONTENT_LEFT_EMU + inset), Emu(int(y_top)), Emu(int(label_col_w - inset)),
              Emu(int(row_h)), item.label or export_i18n.export_t("export_untitled_scope_item", language),
              pt, _INK, bold=True)
        x0 = grid_left + (item.start_week - 1) * week_col_w
        x1 = grid_left + item.end_week * week_col_w
        bar_h = int(row_h * 0.48)
        bar_pt = max(_BAR_LABEL_PT_MIN, _BAR_LABEL_PT_REF * scale)
        _pill(slide, Emu(int(x0 + Inches(0.02))), Emu(int(y_top + (row_h - bar_h) / 2)),
              Emu(int(x1 - x0 - Inches(0.04))), Emu(bar_h), colour,
              export_i18n.export_t("pptx_duration_weeks_short", language, weeks=item.weeks), bar_pt)

    if has_ms:
        _milestones(slide, model, grid_left, grid_right, bounds["grid_top"], bounds["grid_bottom"],
                    bounds["ms_y"], bounds["ms_h"], scale)
    legend = [(_hex(program_render.STAGE_COLOURS[i % len(program_render.STAGE_COLOURS)]), name)
              for i, name in enumerate(model.stages)]
    if model.milestones:
        legend.append((_MILESTONE_ORANGE, export_i18n.export_t("pptx_milestone_legend", language)))
    _legend(slide, legend, bounds["legend_y"], CONTENT_LEFT_EMU, scale)

    _grow(prs, total_h_emu)
    return _save(prs)


# ---------------------------------------------------------------------------
# C. Formal table -- a real PowerPoint table, so it stays editable
# ---------------------------------------------------------------------------

_TABLE_HEADER_H_REF = 0.40
_TABLE_ROW_H_REF = 0.56
_TABLE_ROW_H_MAX_IN = 1.35   # "comfortable maximum" -- a 5-row table
                             # shouldn't grow rows past this even on an
                             # otherwise-empty slide.
_TABLE_GAP_REF = 0.16
_TABLE_LEGEND_H_REF = 0.32


def _slide_table(model, accent: RGBColor, language: str = "en") -> bytes:
    n = len(model.items)
    natural_h_in = (_TABLE_HEADER_H_REF + n * _TABLE_ROW_H_REF + _TABLE_GAP_REF
                   + _TABLE_LEGEND_H_REF)
    scale, total_h_emu, avail_h_emu = _fit_height(int(Inches(natural_h_in)))

    header_h = int(Inches(_TABLE_HEADER_H_REF) * scale)
    gap_h = int(Inches(_TABLE_GAP_REF) * scale)
    legend_h = int(Inches(_TABLE_LEGEND_H_REF) * scale)
    chrome_h = header_h + gap_h + legend_h
    remaining_for_rows = max(0, avail_h_emu - chrome_h)
    # Rows are NOT a straight scale multiply -- per the brief, "distribute
    # across the available height, up to a comfortable maximum". Computed
    # this way, a short table spends its leftover height on taller (still
    # capped) rows instead of a gap the reader reads as broken layout; only
    # once the cap is hit does genuine leftover exist, and it lands in the
    # one gap between the table and the legend rather than stranding the
    # table in the slide's top third.
    row_h = min(int(Inches(_TABLE_ROW_H_MAX_IN)), int(remaining_for_rows / n)) if n else 0

    prs, slide = _new_deck()
    top = _heading(slide, model, language)

    headers = [export_i18n.export_t("export_table_header_scope_item", language),
              export_i18n.export_t("export_table_header_commence", language),
              export_i18n.export_t("export_table_header_complete", language),
              export_i18n.export_t("export_table_header_duration", language)]
    widths = [Inches(5.0), Inches(2.1), Inches(2.1), Inches(1.89)]
    rows = n + 1

    shape = slide.shapes.add_table(
        rows, len(headers), _M, Emu(int(top)), Emu(int(sum(int(w) for w in widths))),
        Emu(int(header_h + row_h * n)),
    )
    table = shape.table
    # No banded fills, no first-row emphasis: the styling below is the whole
    # look, and PowerPoint's default table style would otherwise paint over it.
    table.first_row = False
    table.horz_banding = False
    for index, width in enumerate(widths):
        table.columns[index].width = Emu(int(width))
    table.rows[0].height = Emu(int(header_h))
    for index in range(1, rows):
        table.rows[index].height = Emu(int(row_h))

    head_pt = max(6.5, 9.5 * scale)
    cell_pt = max(6.5, 9.0 * scale)
    name_pt = max(_ROW_LABEL_PT_MIN, _ROW_LABEL_PT_REF * scale)

    def _fill(cell, colour, text, size, font_colour, bold=False):
        cell.fill.solid()
        cell.fill.fore_color.rgb = colour
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Pt(6)
        cell.margin_right = Pt(6)
        cell.margin_top = Pt(2)
        cell.margin_bottom = Pt(2)
        paragraph = cell.text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = _FONT
        run.font.color.rgb = font_colour
        _no_borders(cell)

    for index, header in enumerate(headers):
        _fill(table.cell(0, index), accent, header.upper(), head_pt, _WHITE, bold=True)

    def _week_text(week: int) -> str:
        label = (model.week_labels[week - 1] if week - 1 < len(model.week_labels)
                else export_i18n.export_t("pptx_week_number_short", language, week=week))
        date_text = model.week_dates[week - 1] if week - 1 < len(model.week_dates) else ""
        return f"{label} · {date_text}" if date_text else label

    for row_index, item in enumerate(model.items, start=1):
        band = _WHITE if row_index % 2 else _ROW_BAND
        _fill(table.cell(row_index, 0), band,
              item.label or export_i18n.export_t("export_untitled_scope_item", language), name_pt,
              _INK, bold=True)
        _fill(table.cell(row_index, 1), band, _week_text(item.start_week), cell_pt, _DARK_TEXT)
        _fill(table.cell(row_index, 2), band, _week_text(item.end_week), cell_pt, _DARK_TEXT)
        weeks_key = "pptx_duration_weeks_long_singular" if item.weeks == 1 else "pptx_duration_weeks_long_plural"
        _fill(table.cell(row_index, 3), band,
              export_i18n.export_t(weeks_key, language, weeks=item.weeks), cell_pt, _DARK_TEXT)

    bottom = int(top + header_h + row_h * n)
    legend_y = bottom + gap_h + max(0, remaining_for_rows - row_h * n)
    if model.start_date_text:
        anchored_note = export_i18n.export_t(
            "export_program_anchored_note", language, start_date=model.start_date_text)
        _text(slide, _M, Emu(int(legend_y)), Emu(int(_SLIDE_W - 2 * _M)), Emu(int(legend_h)),
              anchored_note, max(6.5, 8.0 * scale), _MUTED)
    _grow(prs, total_h_emu)
    return _save(prs)


# ---------------------------------------------------------------------------
# D. Modern timeline
# ---------------------------------------------------------------------------

_TL_HEADER_H_REF = 0.36
_TL_ROW_H_REF = 0.80
_TL_GAP_REF = 0.14
_TL_MONTH_BAND_H_REF = 0.30
_TL_MILESTONE_H_REF = 0.42
_TL_LEGEND_H_REF = 0.34


def _slide_timeline(model, accent: RGBColor, language: str = "en") -> bytes:
    n = len(model.items)
    has_ms = bool(model.milestones)
    has_bands = bool(model.month_bands)
    natural_h_in = ((_TL_MONTH_BAND_H_REF if has_bands else 0.0)
                   + _TL_HEADER_H_REF + _TL_GAP_REF
                   + n * _TL_ROW_H_REF + _TL_GAP_REF
                   + (_TL_MILESTONE_H_REF + _TL_GAP_REF if has_ms else 0.0)
                   + _TL_LEGEND_H_REF)
    scale, total_h_emu, avail_h_emu = _fit_height(int(Inches(natural_h_in)))

    prs, slide = _new_deck()
    top = _heading(slide, model, language)

    grid_left = CONTENT_LEFT_EMU
    grid_right = CONTENT_RIGHT_EMU
    weeks = max(1, len(model.week_labels))
    week_col_w = (grid_right - grid_left) / weeks

    bounds: dict[str, int] = {}
    flow = _PVFlow()

    if has_bands:
        def draw_bands(y_top, y_bottom, scale):
            band_h = int((y_bottom - y_top) * 0.72)
            band_y = y_bottom - band_h
            pt = max(6.5, 8.5 * scale)
            for name, first, last in model.month_bands:
                x0 = grid_left + (first - 1) * week_col_w
                x1 = grid_left + last * week_col_w
                _rect(slide, Emu(int(x0 + Inches(0.01))), Emu(int(band_y)),
                      Emu(int(x1 - x0 - Inches(0.02))), Emu(band_h), _tint(accent, 0.9))
                _text(slide, Emu(int(x0)), Emu(int(band_y)), Emu(int(x1 - x0)), Emu(band_h),
                      name, pt, accent, bold=True, align=PP_ALIGN.CENTER)
        flow.block(Inches(_TL_MONTH_BAND_H_REF), draw_bands)

    def draw_header(y_top, y_bottom, scale):
        _week_header(slide, model, grid_left, y_top, week_col_w, scale, show_dates=False)
        bounds["grid_top"] = y_bottom
    flow.block(Inches(_TL_HEADER_H_REF), draw_header)
    flow.gap(Inches(_TL_GAP_REF))

    row_draws = []
    for index, item in enumerate(model.items):
        def draw_row(y_top, y_bottom, scale, item=item):
            row_draws.append((item, y_top, y_bottom))
            bounds["grid_bottom"] = y_bottom
        flow.block(Inches(_TL_ROW_H_REF), draw_row)

    flow.gap(Inches(_TL_GAP_REF))

    if has_ms:
        def draw_ms_anchor(y_top, y_bottom, scale):
            bounds["ms_y"] = y_top
            bounds["ms_h"] = y_bottom - y_top
        flow.block(Inches(_TL_MILESTONE_H_REF), draw_ms_anchor)
        flow.gap(Inches(_TL_GAP_REF))

    def draw_legend_anchor(y_top, y_bottom, scale):
        bounds["legend_y"] = y_top
    flow.block(Inches(_TL_LEGEND_H_REF), draw_legend_anchor)

    flow.render(top, scale, avail_h_emu)

    _gridlines(slide, model, grid_left, bounds["grid_top"], bounds["grid_bottom"], week_col_w, scale)
    for item, y_top, y_bottom in row_draws:
        row_h = y_bottom - y_top
        x0 = grid_left + (item.start_week - 1) * week_col_w
        x1 = grid_left + item.end_week * week_col_w
        bar_h = int(row_h * 0.60)
        pt = max(_BAR_LABEL_PT_MIN, _BAR_LABEL_PT_REF * scale + 0.5)
        _pill(slide, Emu(int(x0 + Inches(0.01))), Emu(int(y_top + (row_h - bar_h) / 2)),
              Emu(int(x1 - x0 - Inches(0.02))), Emu(bar_h),
              accent, item.label or export_i18n.export_t("export_untitled_scope_item", language),
              pt, align=PP_ALIGN.LEFT)

    if has_ms:
        _milestones(slide, model, grid_left, grid_right, bounds["grid_top"], bounds["grid_bottom"],
                    bounds["ms_y"], bounds["ms_h"], scale)
    _legend(slide, _activity_legend(model, accent, language), bounds["legend_y"], CONTENT_LEFT_EMU, scale)

    _grow(prs, total_h_emu)
    return _save(prs)


_RENDERERS = {
    "gantt": _slide_gantt,
    "swimlanes": _slide_swimlanes,
    "table": _slide_table,
    "timeline": _slide_timeline,
}


def populate_program(
    program_schedule: dict[str, list[bool]],
    week_labels: list[str],
    client_name: str = "",
    project_name: str = "",
    theme_name: str | None = None,
    style: str | None = None,
    methodology_stages: list | None = None,
    start_date=None,
    analysis=None,
    output_language: str = "en",
) -> bytes:
    """
    Builds a fresh A4-landscape .pptx (returned as bytes) of the delivery
    program, in the user's chosen presentation style.

    `style` is one of program_render.STYLES; anything else (including None)
    falls back to program_render.DEFAULT_STYLE, and "swimlanes" without any
    methodology stages to group by falls back to the Gantt -- the same
    fallback the preview shows and says out loud, so the deck can never
    disagree with what the user saw.

    An empty/missing schedule renders a single placeholder slide rather than
    an empty grid.
    """
    model = program_render.build_model(
        program_schedule or {}, week_labels or [], methodology_stages or [],
        start_date, analysis, project_name or "", client_name or "",
    )
    if model.is_empty:
        return _placeholder_slide(model, output_language)
    resolved = program_render.effective_style(model, style or program_render.DEFAULT_STYLE)
    accent = _resolve_palette(theme_name)["active_bg"]
    return _RENDERERS[resolved](model, accent, output_language)
