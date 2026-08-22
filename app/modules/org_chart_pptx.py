"""
org_chart_pptx.py

Builds the org chart PowerPoint FROM SCRATCH, straight from a project's
resourcing plan -- no template file to edit or keep in sync -- in whichever
of the four approved presentation styles the user picked.

An earlier version worked by surgically editing a fixed template: filling in
[Name] placeholders, deleting boxes for disciplines not on this project,
adding ad hoc boxes for disciplines the template had no box for. That kept
breaking in new ways as real projects were tried against it -- a deleted box
left its connector dangling, a discipline outside the template's fixed
allowlist got exiled to a disconnected strip, removing boxes left the rest
asymmetrically spaced. Every fix was another special case bolted onto a
fundamentally static layout.

The version after that had no static layout, but it did have one hardcoded
look, and one hardcoded four-box management chain. This one has neither. It
builds from the shared model in modules/org_chart_render.py -- the SAME
object the on-screen preview and the pack's embedded chart render from, so
they cannot show different teams -- and dispatches to one of four style
renderers. Every row comes from the plan: 2 disciplines or 9, with or
without a Design Manager, with or without a reviewer. Columns wrap to
further rows rather than shrinking past legibility, and the slide grows
rather than storing its last row off the bottom edge (which PowerPoint
renders as simply not there).

Everything is native PowerPoint shapes, not a pasted picture: the point of
the companion deck is that it can be tidied up further in PowerPoint.

Nothing here is invented. An unassigned slot is a dashed red TBC. A role the
user removed is absent, with no TBC -- deliberate absence is not a gap. A
person's card shows exactly two lines, name and role/title -- qualifications
never appear here (see the Word pack's Key Personnel profiles for those). An
assurance element (a dedicated reviewer role) appears only where the plan
actually holds one; the separate Peer Review element -- one row per
discipline against its nominated reviewer, red TBC until one is entered --
is unconditional, appearing in every style whenever the plan has at least
one discipline, because every discipline's work needs a reviewer eventually,
not only once one has been named.
"""

from __future__ import annotations

import io

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---- palette -------------------------------------------------------------
# The fallback palette, used when no proposal theme is given. Everything
# below is now themed from divider_designer.THEME_COLOURS instead (see
# _resolve_palette) -- this chart was the last generated artefact still
# hardcoded to navy/cyan, so a Government-themed pack exported a green
# cover, green dividers, green tables and then one cyan org chart.
_NAVY = RGBColor(0x00, 0x37, 0x63)
_BLACK = RGBColor(0x00, 0x00, 0x00)
_CYAN_HEADER = RGBColor(0x00, 0xB0, 0xF0)
_CYAN_BODY = RGBColor(0xE9, 0xF9, 0xFD)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
_GREY_TEXT = RGBColor(0x59, 0x59, 0x59)
_RED_TBC = RGBColor(0xC0, 0x00, 0x00)
_RED_TBC_ON_DARK = RGBColor(0xFF, 0x8A, 0x80)
_LINE_COLOR = RGBColor(0x59, 0x59, 0x59)
_BORDER_COLOR = RGBColor(0xBF, 0xBF, 0xBF)

_FONT = "Calibri"

# A4 landscape -- the same precise dimensions modules/program_pptx.py and
# modules/methodology_pptx.py already use (both say so in their own
# comments: "same slide size as org_chart_pptx.py"), which this module had
# drifted from (it was 13.333in x 7.5in, 16:9). Matches the fix brief and
# modules/org_chart_render.py's PAGE_W_IN/PAGE_H_IN, so the companion deck is
# the same physical page as the printed/pasted chart, not a wider 16:9 slide
# with the same content floating in the top of it.
_SLIDE_W = Inches(11.6929)
_SLIDE_H = Inches(8.2677)

def _tint(rgb: RGBColor, amount: float) -> RGBColor:
    """Blend towards white. amount=0 -> unchanged, 1 -> white."""
    return RGBColor(*(int(c + (255 - c) * amount) for c in (rgb[0], rgb[1], rgb[2])))


def _text_on(bg: RGBColor) -> RGBColor:
    """Readable text colour for a given fill. A themed accent can be light
    (Minimalist) or dark (Corporate); white-on-everything was safe while the
    header was always cyan and isn't any more."""
    luminance = (0.299 * bg[0] + 0.587 * bg[1] + 0.114 * bg[2]) / 255
    return _DARK_TEXT if luminance > 0.62 else _WHITE


def _resolve_palette(theme_name: str | None) -> dict:
    """The chart's colours, from the same divider_designer.THEME_COLOURS every
    other generated graphic uses. No theme (or an unknown one) keeps the
    original navy/cyan look, so nothing changes for a project that never
    picked a theme."""
    from modules.divider_designer import THEME_COLOURS

    colours = THEME_COLOURS.get(theme_name or "")
    if not colours:
        return {
            "mgmt": _NAVY, "mgmt_text": _WHITE,
            "header": _CYAN_HEADER, "header_text": _WHITE,
            "body": _CYAN_BODY,
        }
    primary = RGBColor(*colours["primary"])
    accent = RGBColor(*colours["accent"])
    # Minimalist's "primary" is a near-white wash by design -- the same
    # special case export_docx._theme_colours() and program_pptx apply.
    mgmt = RGBColor(0x2A, 0x2A, 0x2A) if theme_name == "Minimalist" else primary
    return {
        "mgmt": mgmt, "mgmt_text": _text_on(mgmt),
        "header": accent, "header_text": _text_on(accent),
        "body": _tint(accent, 0.88),
    }


def _set_text(text_frame, lines, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    """lines: list of (text, size_pt, bold, color) tuples, one per paragraph."""
    text_frame.word_wrap = True
    text_frame.vertical_anchor = anchor
    text_frame.margin_left = Pt(4)
    text_frame.margin_right = Pt(4)
    text_frame.margin_top = Pt(2)
    text_frame.margin_bottom = Pt(2)
    for i, (text, size, bold, color) in enumerate(lines):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = _FONT
        run.font.color.rgb = color


def _rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def _connector(slide, conn_type, x1, y1, x2, y2):
    c = slide.shapes.add_connector(conn_type, x1, y1, x2, y2)
    c.line.color.rgb = _LINE_COLOR
    c.line.width = Pt(1.0)
    return c


# ---------------------------------------------------------------------------
# Four presentation styles -- see modules/org_chart_render.py, which holds the
# shared model and draws the same four styles as PNGs for the preview and the
# pack. The shapes below are built from that SAME model object, so the deck
# cannot show a different team from the one the user approved on screen.
#
# Everything is native PowerPoint shapes, not a pasted picture: the point of
# the companion deck is that it can be tidied up further in PowerPoint.
# ---------------------------------------------------------------------------

_STYLE_MARGIN = Inches(0.35)
_CARD_EDGE = RGBColor(0xE4, 0xE8, 0xEF)
_CLIENT_DARK = RGBColor(0x13, 0x1A, 0x2A)
_INK = RGBColor(0x11, 0x18, 0x27)
_MUTED = RGBColor(0x7A, 0x85, 0x98)
_TBC_FILL = RGBColor(0xFE, 0xF2, 0xF2)
_ASSURANCE_AMBER = RGBColor(0xB4, 0x53, 0x09)
_ASSURANCE_FILL = RGBColor(0xFF, 0xF7, 0xED)
_LANE_FILL = RGBColor(0xF7, 0xF9, 0xFC)


def _hex(value: str) -> RGBColor:
    return RGBColor.from_string(str(value).lstrip("#").upper())


def _round(slide, x, y, w, h, fill, line=None, dashed=False, radius=0.14):
    from pptx.enum.dml import MSO_LINE_DASH_STYLE

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    try:
        shape.adjustments[0] = radius
    except (IndexError, KeyError):
        pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.0)
        if dashed:
            shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    shape.shadow.inherit = False
    return shape


def _bar(slide, x, y, w, h, fill):
    return _rect(slide, Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)), fill)


def _circle(slide, cx, cy, d, fill, text, text_colour):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(cx - d / 2)), Emu(int(cy - d / 2)),
                                   Emu(int(d)), Emu(int(d)))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    _set_text(shape.text_frame, [(text, 9, True, text_colour)], align=PP_ALIGN.CENTER)
    # An oval's text frame insets are wide enough that two initials wrapped
    # onto two lines -- "AD" came out as "A" over "D".
    frame = shape.text_frame
    frame.word_wrap = False
    frame.margin_left = Pt(0)
    frame.margin_right = Pt(0)
    frame.margin_top = Pt(0)
    frame.margin_bottom = Pt(0)
    return shape


def _stack(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT):
    """A textbox holding a person's stacked lines, vertically centred."""
    box = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    _set_text(box.text_frame, lines, align=align)


# Reference (scale=1.0) font sizes for a person's two stacked lines. Matches
# org_chart_render.py's boosted sizes (up from the old fixed 11pt/8.5pt that
# never changed regardless of team size) so the companion deck and the
# on-screen/DOCX chart land in the same ballpark at the same team size, not
# just the same physical page size.
_NAME_PT_REF = 14.5
_ROLE_PT_REF = 10.5
_NAME_PT_MIN = 8.0
_ROLE_PT_MIN = 6.5


def _person_lines(person, role_colour: RGBColor, scale: float = 1.0):
    """The (text, size, bold, colour) rows for one person -- ALWAYS exactly
    two: name (or "TBC"), then role/title. Qualifications are deliberately
    never drawn on the chart (see org_chart_render's module docstring); a
    full CV sentence on a card overflowed the card and collided with
    whatever sat above it."""
    name_pt = max(_NAME_PT_MIN, _NAME_PT_REF * scale)
    role_pt = max(_ROLE_PT_MIN, _ROLE_PT_REF * scale)
    if person.is_tbc:
        return [("TBC", name_pt, True, _RED_TBC),
                (person.role or "", role_pt, True, _RED_TBC)]
    lines = [(person.name, name_pt, True, _INK)]
    if person.role:
        if person.role_is_placeholder:
            colour = _RED_TBC
        else:
            colour = role_colour if person.is_lead else _GREY_TEXT
        lines.append((person.role, role_pt, True, colour))
    return lines


def _title_block(slide, model, style_note: str = ""):
    box = slide.shapes.add_textbox(_STYLE_MARGIN, Inches(0.22),
                                   Emu(int(_SLIDE_W - 2 * _STYLE_MARGIN)), Inches(0.7))
    box.text_frame.word_wrap = True
    _set_text(box.text_frame, [("Project organisation", 20, True, _DARK_TEXT)],
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    if model.heading:
        right = slide.shapes.add_textbox(Emu(int(_SLIDE_W / 2)), Inches(0.24),
                                         Emu(int(_SLIDE_W / 2 - _STYLE_MARGIN)), Inches(0.32))
        _set_text(right.text_frame, [(model.heading, 11, True, _GREY_TEXT)],
                  align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.TOP)
    return int(Inches(0.95))


# ---------------------------------------------------------------------------
# Fixed A4-landscape page, scale-to-fill layout -- see org_chart_render.py's
# identical-in-spirit engine and its module note for WHY: the deck used to
# grow (via _grow, below -- now only a defensive last resort, not the normal
# path) or shrink around fixed-size shapes; now the SLIDE is fixed and the
# shapes scale to fill it.
# ---------------------------------------------------------------------------

CONTENT_TOP_EMU = int(Inches(0.95))
CONTENT_BOTTOM_EMU = int(_SLIDE_H - _STYLE_MARGIN)
AVAIL_H_EMU = CONTENT_BOTTOM_EMU - CONTENT_TOP_EMU
AVAIL_W_EMU = int(_SLIDE_W - 2 * _STYLE_MARGIN)

MIN_SCALE = 0.55
MAX_SCALE = 1.35


class _PFlow:
    """EMU-space counterpart to org_chart_render._Flow -- see that class's
    docstring for the full rationale (fixed-size blocks, stretchable gaps
    that also draw, e.g. a connector line, so leftover page height becomes
    breathing room at every seam rather than one dead band). PowerPoint's
    EMU y-axis already increases downward from the top of the slide, so
    there's no fraction/axes conversion to do here -- everything is plain
    integer EMU arithmetic."""

    def __init__(self):
        self._items = []

    def block(self, height_emu, draw) -> None:
        self._items.append({"kind": "block", "h": max(0, int(height_emu)), "draw": draw})

    def gap(self, height_emu, draw=None) -> None:
        self._items.append({"kind": "gap", "h": max(0, int(height_emu)), "draw": draw})

    def connector(self, height_emu, draw) -> None:
        self.gap(height_emu, draw=draw)

    def natural_height(self) -> int:
        return sum(it["h"] for it in self._items)

    def render(self, top_emu: int, scale: float, avail_emu: int = AVAIL_H_EMU) -> int:
        n_gaps = sum(1 for it in self._items if it["kind"] == "gap")
        used = int(self.natural_height() * scale)
        leftover = max(0, avail_emu - used)
        extra_per_gap = (leftover // n_gaps) if n_gaps else 0
        y = top_emu
        for it in self._items:
            h = int(it["h"] * scale) + (extra_per_gap if it["kind"] == "gap" else 0)
            y_bottom = y + h
            if it["draw"] is not None:
                it["draw"](y, y_bottom, scale)
            y = y_bottom
        return y


def _psolve_scale(build_flow, per_row_candidates, avail_w_emu=AVAIL_W_EMU, avail_h_emu=AVAIL_H_EMU):
    """EMU counterpart to org_chart_render._solve_scale -- same fewest-rows-
    first search, same >6% margin before a more-wrapped candidate takes
    over (see that function's docstring for why the margin matters)."""
    best = None
    for per_row in per_row_candidates:
        flow, row_w = build_flow(per_row)
        natural_h = flow.natural_height()
        height_fit = (avail_h_emu / natural_h) if natural_h > 0 else MAX_SCALE
        width_fit = (avail_w_emu / row_w) if row_w > 0 else MAX_SCALE
        scale = max(MIN_SCALE, min(MAX_SCALE, min(height_fit, width_fit)))
        if best is None or scale > best[1] * 1.06:
            best = (flow, scale, per_row)
    return best


def _grow(prs, needed: float) -> None:
    """A defensive last resort only: every style now sizes itself to fit the
    fixed A4 slide via the scale-to-fill search above, but an extreme edge
    case (a great many disciplines each with many support members, all
    already at MIN_SCALE) could still in principle need slightly more room
    than the page. Growing the slide keeps that content visible rather than
    stored off the bottom edge -- which PowerPoint renders as simply not
    there -- but this should not fire in ordinary use."""
    if int(needed) > int(prs.slide_height):
        prs.slide_height = Emu(int(needed))


def _new_deck():
    prs = Presentation()
    prs.slide_width = _SLIDE_W
    prs.slide_height = _SLIDE_H
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def _save(prs) -> bytes:
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()


def _client_label(model) -> tuple[str, RGBColor]:
    name = (model.client_name or "").strip()
    return (name or "[CLIENT NAME]", _WHITE if name else _RED_TBC_ON_DARK)


# ---------------------------------------------------------------------------
# Peer Review panel -- shared by the cards, columns and tree styles (bands
# folds the same information into its "Assurance" band instead, see below).
# Unconditional whenever the plan has at least one discipline: one row per
# discipline against its nominated reviewer, red TBC until one is entered.
# ---------------------------------------------------------------------------

_PANEL_HEADER_EMU = int(Inches(0.30))
_PANEL_ROW_EMU = int(Inches(0.28))
_PANEL_W_EMU = int(Inches(3.1))


def _draw_peer_review_panel(slide, model, scale: float, x_emu: int, top_emu: int, w_emu_ref: int) -> int:
    """Draws the panel anchored at its top-left corner and returns its
    bottom y (EMU), so a caller that needs to clear it knows how far down
    it reaches."""
    n = len(model.disciplines)
    if not n:
        return top_emu
    w = int(w_emu_ref * scale)
    header_h = int(_PANEL_HEADER_EMU * scale)
    row_h = int(_PANEL_ROW_EMU * scale)
    total_h = header_h + row_h * n
    _round(slide, Emu(int(x_emu)), Emu(int(top_emu)), Emu(w), Emu(total_h),
          _ASSURANCE_FILL, line=_ASSURANCE_AMBER, radius=0.05)
    _stack(slide, Emu(int(x_emu)), Emu(int(top_emu)), Emu(w), Emu(header_h),
          [("PEER REVIEW", max(7.0, 9.0 * scale), True, _ASSURANCE_AMBER)], align=PP_ALIGN.CENTER)
    row_y = top_emu + header_h
    pad = int(w * 0.06)
    for group in model.disciplines:
        reviewer = (group.peer_reviewer or "").strip()
        tbc = not reviewer
        _stack(slide, Emu(int(x_emu + pad)), Emu(row_y), Emu(int(w * 0.55 - pad)), Emu(row_h),
              [(group.name, max(6.5, 8.5 * scale), True, _INK)], align=PP_ALIGN.LEFT)
        _stack(slide, Emu(int(x_emu + w * 0.55)), Emu(row_y), Emu(int(w * 0.45 - pad)), Emu(row_h),
              [(reviewer or "TBC", max(6.5, 8.5 * scale), True, _RED_TBC if tbc else _ASSURANCE_AMBER)],
              align=PP_ALIGN.RIGHT)
        row_y += row_h
    return top_emu + total_h


# --- A. Executive cards ----------------------------------------------------

_CARDS_CLIENT_W = int(Inches(3.0))
_CARDS_CLIENT_H = int(Inches(0.55))
_CARDS_CARD_W = int(Inches(2.55))
_CARDS_CARD_H = int(Inches(0.66))
_CARDS_COL_GAP = int(Inches(0.26))
_CARDS_ROW_GAP = int(Inches(0.14))
_CARDS_CONNECTOR = int(Inches(0.22))
_CARDS_CAPTION = int(Inches(0.26))
_CARDS_SECTION_GAP = int(Inches(0.24))
_CARDS_ROW_TO_ROW = int(Inches(0.22))


def _avatar_card(slide, x, y, w, h, person, accent: RGBColor, scale: float, badge: str = ""):
    tbc = person.is_tbc
    _round(slide, Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)),
          _TBC_FILL if tbc else _WHITE, line=_RED_TBC if tbc else _CARD_EDGE, dashed=tbc, radius=0.12)
    bar_colour = _ASSURANCE_AMBER if badge else accent
    if not tbc and person.is_lead:
        bar_h = int(h * 0.09)
        _bar(slide, x + int(w * 0.04), y, w - int(w * 0.08), bar_h, bar_colour)
    if badge and not tbc:
        badge_w, badge_h = int(w * 0.46), int(h * 0.30)
        bx = x + w - badge_w - int(w * 0.02)
        by = y - int(badge_h * 0.45)
        _round(slide, Emu(int(bx)), Emu(int(by)), Emu(badge_w), Emu(badge_h),
              _ASSURANCE_FILL, line=_ASSURANCE_AMBER, radius=0.3)
        _stack(slide, Emu(int(bx)), Emu(int(by)), Emu(badge_w), Emu(badge_h),
              [(badge.upper(), max(5.5, 7.0 * scale), True, _ASSURANCE_AMBER)], align=PP_ALIGN.CENTER)
    avatar_d = int(h * 0.55)
    _circle(slide, x + int(h * 0.42), y + h // 2, avatar_d,
           _TBC_FILL if tbc else _tint(bar_colour, 0.86),
           person.initials, _RED_TBC if tbc else bar_colour)
    _stack(slide, Emu(int(x + h * 0.78)), Emu(int(y)), Emu(int(w - h * 0.82)), Emu(int(h)),
          _person_lines(person, bar_colour, scale))


def _render_cards_slide(prs, slide, model, accent: RGBColor):
    from modules.org_chart_render import _row_candidates, _balanced_wrap

    centre = int(_SLIDE_W / 2)
    # Same split as org_chart_render.py's cards style: the leadership
    # section sits beside the Peer Review panel and gives up width for it;
    # the discipline rows always sit below the panel (panel_extra below
    # pushes them down whenever the panel is taller than the leadership
    # section) and keep the full-width centre.
    _panel_reserve_emu = int(_PANEL_W_EMU * MAX_SCALE) if model.disciplines else 0
    lead_centre = (int(_STYLE_MARGIN) + (AVAIL_W_EMU - _panel_reserve_emu) // 2
                  if model.disciplines else centre)

    lead_people = list(model.leadership)
    top_person = lead_people[0] if lead_people else None
    rank = [(p, "") for p in lead_people[1:]] + [(p, "QA / Review") for p in model.assurance]
    rank_per_row = _balanced_wrap(len(rank), 4) if rank else 0
    rank_rows = -(-len(rank) // rank_per_row) if rank else 0
    rank_h = (rank_rows * _CARDS_CARD_H + max(0, rank_rows - 1) * _CARDS_ROW_GAP) if rank else 0

    leadership_h = _CARDS_CLIENT_H
    if top_person is not None:
        leadership_h += _CARDS_CONNECTOR + _CARDS_CARD_H
    if rank:
        leadership_h += _CARDS_CONNECTOR + rank_h
    panel_h = (_PANEL_HEADER_EMU + len(model.disciplines) * _PANEL_ROW_EMU) if model.disciplines else 0
    panel_extra = max(0, panel_h - leadership_h)

    def build_flow(per_row):
        disciplines = model.disciplines
        chunks = ([disciplines[i:i + per_row] for i in range(0, len(disciplines), per_row)]
                 if disciplines else [])
        row_infos = []
        disc_width = 0
        for chunk in chunks:
            tallest = max(len(g.people) for g in chunk)
            row_h = _CARDS_CAPTION + tallest * _CARDS_CARD_H + max(0, tallest - 1) * _CARDS_ROW_GAP
            row_infos.append((chunk, row_h))
            width = len(chunk) * _CARDS_CARD_W + (len(chunk) - 1) * _CARDS_COL_GAP
            disc_width = max(disc_width, width)
        rank_width = (rank_per_row * _CARDS_CARD_W + max(0, rank_per_row - 1) * _CARDS_COL_GAP) if rank else 0
        leadership_width = max(_CARDS_CLIENT_W, _CARDS_CARD_W, rank_width)
        # Same split fit as org_chart_render.py's cards style: the
        # leadership section's own available width excludes the panel's
        # reserved slice, the discipline rows' doesn't, so each gets its own
        # fraction and the larger (more binding) one wins.
        lead_avail = (AVAIL_W_EMU - _panel_reserve_emu) if disciplines else AVAIL_W_EMU
        lead_ratio = (leadership_width / lead_avail) if lead_avail > 0 else 1.0
        disc_ratio = (disc_width / AVAIL_W_EMU) if disc_width > 0 else 0.0
        row_width = int(max(lead_ratio, disc_ratio) * AVAIL_W_EMU)

        flow = _PFlow()

        def draw_client(y_top, y_bottom, scale):
            w = int(_CARDS_CLIENT_W * scale)
            h = y_bottom - y_top
            x = lead_centre - w // 2
            _round(slide, Emu(int(x)), Emu(int(y_top)), Emu(w), Emu(h), _CLIENT_DARK)
            label, colour = _client_label(model)
            _stack(slide, Emu(int(x)), Emu(int(y_top)), Emu(w), Emu(h),
                  [(label, max(9.0, 12.0 * scale), True, colour),
                   (model.client_role, max(6.5, 9.0 * scale), True, _GREY_TEXT)], align=PP_ALIGN.CENTER)
        flow.block(_CARDS_CLIENT_H, draw_client)

        if top_person is not None:
            def draw_conn1(y_top, y_bottom, scale):
                _connector(slide, MSO_CONNECTOR.STRAIGHT, lead_centre, y_top, lead_centre, y_bottom)
            flow.connector(_CARDS_CONNECTOR, draw_conn1)

            def draw_top(y_top, y_bottom, scale):
                w = int(_CARDS_CARD_W * scale)
                h = y_bottom - y_top
                x = lead_centre - w // 2
                _avatar_card(slide, x, y_top, w, h, top_person, accent, scale)
            flow.block(_CARDS_CARD_H, draw_top)

        if rank:
            def draw_conn2(y_top, y_bottom, scale):
                _connector(slide, MSO_CONNECTOR.STRAIGHT, lead_centre, y_top, lead_centre, y_bottom)
            flow.connector(_CARDS_CONNECTOR, draw_conn2)

            def draw_rank(y_top, y_bottom, scale):
                row_h = int(_CARDS_CARD_H * scale)
                row_gap = int(_CARDS_ROW_GAP * scale)
                col_gap = int(_CARDS_COL_GAP * scale)
                col_w = int(_CARDS_CARD_W * scale)
                row_chunks = [rank[i:i + rank_per_row] for i in range(0, len(rank), rank_per_row)]
                ry = y_top
                for row_chunk in row_chunks:
                    total = len(row_chunk) * col_w + (len(row_chunk) - 1) * col_gap
                    rx = lead_centre - total // 2
                    for person, badge in row_chunk:
                        _avatar_card(slide, rx, ry, col_w, row_h, person,
                                    _ASSURANCE_AMBER if badge else accent, scale, badge=badge)
                        rx += col_w + col_gap
                    ry += row_h + row_gap
            flow.block(rank_h, draw_rank)

        if disciplines:
            def draw_elbow(y_top, y_bottom, scale):
                # Reconciles the leadership section's lead_centre with the
                # discipline rows' full-width centre -- see the identical
                # comment in org_chart_render.py's cards style.
                if lead_centre != centre:
                    y_mid = (y_top + y_bottom) // 2
                    _connector(slide, MSO_CONNECTOR.STRAIGHT, lead_centre, y_top, lead_centre, y_mid)
                    _connector(slide, MSO_CONNECTOR.STRAIGHT, lead_centre, y_mid, centre, y_mid)
                    _connector(slide, MSO_CONNECTOR.STRAIGHT, centre, y_mid, centre, y_bottom)
                else:
                    _connector(slide, MSO_CONNECTOR.STRAIGHT, centre, y_top, centre, y_bottom)
            flow.connector(_CARDS_SECTION_GAP + panel_extra, draw_elbow)

            def draw_stub(y_top, y_bottom, scale):
                _connector(slide, MSO_CONNECTOR.STRAIGHT, centre, y_top, centre, y_bottom)
            flow.connector(_CARDS_CONNECTOR, draw_stub)

            for row_index, (chunk, row_h) in enumerate(row_infos):
                def draw_row(y_top, y_bottom, scale, chunk=chunk):
                    bus_y = y_top
                    col_w = int(_CARDS_CARD_W * scale)
                    col_gap = int(_CARDS_COL_GAP * scale)
                    total = len(chunk) * col_w + (len(chunk) - 1) * col_gap
                    x0 = centre - total // 2
                    centres = [x0 + col_w // 2 + i * (col_w + col_gap) for i in range(len(chunk))]
                    if len(centres) > 1:
                        _connector(slide, MSO_CONNECTOR.STRAIGHT, centres[0], bus_y, centres[-1], bus_y)
                    caption_h = int(_CARDS_CAPTION * scale)
                    card_h = int(_CARDS_CARD_H * scale)
                    row_gap = int(_CARDS_ROW_GAP * scale)
                    for i, group in enumerate(chunk):
                        cx = centres[i]
                        _connector(slide, MSO_CONNECTOR.STRAIGHT, cx, bus_y, cx, bus_y + int(caption_h * 0.35))
                        _stack(slide, Emu(int(x0 + i * (col_w + col_gap))), Emu(int(bus_y + caption_h * 0.35)),
                              Emu(col_w), Emu(int(caption_h * 0.5)),
                              [(group.name.upper(), max(7.0, 9.5 * scale), True, _MUTED)], align=PP_ALIGN.CENTER)
                        card_y = bus_y + caption_h
                        gx = x0 + i * (col_w + col_gap)
                        for person in group.people:
                            _avatar_card(slide, gx, card_y, col_w, card_h, person, accent, scale)
                            card_y += card_h + row_gap
                flow.block(row_h, draw_row)
                if row_index < len(row_infos) - 1:
                    flow.gap(_CARDS_ROW_TO_ROW)

        return flow, row_width

    per_row_candidates = _row_candidates(len(model.disciplines)) if model.disciplines else [1]
    flow, scale, per_row = _psolve_scale(build_flow, per_row_candidates)
    bottom = flow.render(CONTENT_TOP_EMU, scale)

    if model.disciplines:
        panel_w = int(_PANEL_W_EMU * scale)
        panel_x = int(_SLIDE_W - _STYLE_MARGIN) - panel_w
        _draw_peer_review_panel(slide, model, scale, panel_x, CONTENT_TOP_EMU, _PANEL_W_EMU)

    _grow(prs, bottom + int(Inches(0.35)))


def _slide_cards(model, accent: RGBColor) -> bytes:
    prs, slide = _new_deck()
    _title_block(slide, model)
    _render_cards_slide(prs, slide, model, accent)
    return _save(prs)


# --- B. Discipline columns -------------------------------------------------

_COLUMNS_PILL_H = int(Inches(0.55))
_COLUMNS_CLIENT_W = int(Inches(3.4))
_COLUMNS_LEAD_W1 = int(Inches(3.0))
_COLUMNS_LEAD_W2 = int(Inches(2.5))
_COLUMNS_CONNECTOR = int(Inches(0.20))
_COLUMNS_LANE_W = int(Inches(2.75))
_COLUMNS_LANE_GAP = int(Inches(0.24))
_COLUMNS_ROW_H = int(Inches(0.62))
_COLUMNS_CAPTION = int(Inches(0.32))
_COLUMNS_ROW_GAP = int(Inches(0.12))
_COLUMNS_LANE_PAD = int(Inches(0.10))
_COLUMNS_SECTION_GAP = int(Inches(0.26))
_COLUMNS_ROW_TO_ROW = int(Inches(0.22))
_COLUMNS_STRIP_H = int(Inches(0.55))
_COLUMNS_STRIP_W = int(Inches(7.5))


def _pill(slide, x, y, w, h, facecolor, label, sub, label_colour, sub_colour, scale):
    _round(slide, Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)), facecolor, radius=0.35)
    lines = [(label, max(9.0, 12.5 * scale), True, label_colour)]
    if sub:
        lines.append((sub, max(6.5, 9.5 * scale), True, sub_colour))
    _stack(slide, Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)), lines, align=PP_ALIGN.CENTER)


def _slide_columns(model, accent: RGBColor) -> bytes:
    from modules.org_chart_render import _row_candidates, DISCIPLINE_COLOURS

    prs, slide = _new_deck()
    _title_block(slide, model)
    centre = int(_SLIDE_W / 2)

    def build_flow(per_row):
        disciplines = model.disciplines
        chunks = ([disciplines[i:i + per_row] for i in range(0, len(disciplines), per_row)]
                 if disciplines else [])
        row_infos = []
        row_width = _COLUMNS_CLIENT_W
        for chunk in chunks:
            tallest = max(len(g.people) for g in chunk)
            row_h = (_COLUMNS_CAPTION + tallest * _COLUMNS_ROW_H
                    + max(0, tallest - 1) * _COLUMNS_ROW_GAP + 2 * _COLUMNS_LANE_PAD)
            row_infos.append((chunk, row_h))
            width = len(chunk) * _COLUMNS_LANE_W + (len(chunk) - 1) * _COLUMNS_LANE_GAP
            row_width = max(row_width, width)

        flow = _PFlow()

        def draw_client(y_top, y_bottom, scale):
            w = int(_COLUMNS_CLIENT_W * scale)
            h = y_bottom - y_top
            x = centre - w // 2
            label, colour = _client_label(model)
            _pill(slide, x, y_top, w, h, _CLIENT_DARK, f"{label} — Client", "", colour, colour, scale)
        flow.block(_COLUMNS_PILL_H, draw_client)

        for index, person in enumerate(model.leadership):
            def draw_conn(y_top, y_bottom, scale):
                _connector(slide, MSO_CONNECTOR.STRAIGHT, centre, y_top, centre, y_bottom)
            flow.connector(_COLUMNS_CONNECTOR, draw_conn)

            def draw_lead(y_top, y_bottom, scale, person=person, index=index):
                width_ref = _COLUMNS_LEAD_W1 if index == 0 else _COLUMNS_LEAD_W2
                w = int(width_ref * scale)
                h = y_bottom - y_top
                x = centre - w // 2
                tbc = person.is_tbc
                fill = accent if index == 0 else _tint(accent, 0.18)
                fc = _TBC_FILL if tbc else fill
                label_colour = _RED_TBC if tbc else _WHITE
                sub_colour = _RED_TBC if tbc else RGBColor(0xD7, 0xDE, 0xEA)
                _pill(slide, x, y_top, w, h, fc, person.name or "TBC", person.role,
                     label_colour, sub_colour, scale)
            flow.block(_COLUMNS_PILL_H, draw_lead)

        if disciplines:
            flow.gap(_COLUMNS_SECTION_GAP)
            for row_index, (chunk, row_h) in enumerate(row_infos):
                def draw_row(y_top, y_bottom, scale, chunk=chunk):
                    lane_w = int(_COLUMNS_LANE_W * scale)
                    lane_gap = int(_COLUMNS_LANE_GAP * scale)
                    total = len(chunk) * lane_w + (len(chunk) - 1) * lane_gap
                    x0 = centre - total // 2
                    lane_h = y_bottom - y_top
                    for i, group in enumerate(chunk):
                        colour = _hex(DISCIPLINE_COLOURS[
                            (group.people[0].group_index if group.people else disciplines.index(group))
                            % len(DISCIPLINE_COLOURS)])
                        lx = x0 + i * (lane_w + lane_gap)
                        _round(slide, Emu(int(lx)), Emu(int(y_top)), Emu(lane_w), Emu(lane_h),
                              _LANE_FILL, radius=0.06)
                        bar_h = int(lane_h * 0.06)
                        _bar(slide, lx, y_top, lane_w, bar_h, colour)
                        caption_h = int(_COLUMNS_CAPTION * scale)
                        _stack(slide, Emu(int(lx)), Emu(int(y_top + bar_h)), Emu(lane_w),
                              Emu(caption_h - bar_h),
                              [(group.name.upper(), max(7.5, 10.5 * scale), True, colour)],
                              align=PP_ALIGN.CENTER)
                        pad = int(_COLUMNS_LANE_PAD * scale)
                        row_h_in = int(_COLUMNS_ROW_H * scale)
                        row_gap_in = int(_COLUMNS_ROW_GAP * scale)
                        card_y = y_top + caption_h
                        for person in group.people:
                            tbc = person.is_tbc
                            _round(slide, Emu(int(lx + pad)), Emu(int(card_y)),
                                  Emu(int(lane_w - 2 * pad)), Emu(row_h_in),
                                  _TBC_FILL if tbc else _WHITE, line=_RED_TBC if tbc else _CARD_EDGE,
                                  dashed=tbc, radius=0.12)
                            _stack(slide, Emu(int(lx + pad)), Emu(int(card_y)),
                                  Emu(int(lane_w - 2 * pad)), Emu(row_h_in),
                                  _person_lines(person, colour, scale), align=PP_ALIGN.CENTER)
                            card_y += row_h_in + row_gap_in
                flow.block(row_h, draw_row)
                if row_index < len(row_infos) - 1:
                    flow.gap(_COLUMNS_ROW_TO_ROW)

        if model.assurance:
            flow.gap(_COLUMNS_SECTION_GAP)

            def draw_strip(y_top, y_bottom, scale):
                w = int(_COLUMNS_STRIP_W * scale)
                h = y_bottom - y_top
                x = centre - w // 2
                _round(slide, Emu(int(x)), Emu(int(y_top)), Emu(w), Emu(h), _ASSURANCE_FILL,
                      line=RGBColor(0xFB, 0xBF, 0x24), radius=0.3)
                text = " · ".join(f"{p.name or 'TBC'} — {p.role}" for p in model.assurance)
                _stack(slide, Emu(int(x)), Emu(int(y_top)), Emu(w), Emu(h),
                      [(f"★ Independent review: {text}", max(7.5, 10.5 * scale), True, _ASSURANCE_AMBER)],
                      align=PP_ALIGN.CENTER)
            flow.block(_COLUMNS_STRIP_H, draw_strip)

        if disciplines:
            flow.gap(_COLUMNS_SECTION_GAP)
            panel_h = _PANEL_HEADER_EMU + len(disciplines) * _PANEL_ROW_EMU

            def draw_panel(y_top, y_bottom, scale):
                w = int(_PANEL_W_EMU * scale)
                x = centre - w // 2
                _draw_peer_review_panel(slide, model, scale, x, y_top, _PANEL_W_EMU)
            flow.block(panel_h, draw_panel)

        return flow, row_width

    per_row_candidates = _row_candidates(len(model.disciplines)) if model.disciplines else [1]
    flow, scale, per_row = _psolve_scale(build_flow, per_row_candidates)
    bottom = flow.render(CONTENT_TOP_EMU, scale)
    _grow(prs, bottom + int(Inches(0.35)))
    return _save(prs)


def _discipline_colour_for(model, group) -> str:
    from modules.org_chart_render import DISCIPLINE_COLOURS

    index = group.people[0].group_index if group.people else model.disciplines.index(group)
    return DISCIPLINE_COLOURS[index % len(DISCIPLINE_COLOURS)]


# --- C. Governance bands ---------------------------------------------------
#
# Full slide-content-width by design -- see org_chart_render.py's identical
# rationale: no "row of N boxes" whose count trades against width, so only
# band/chip HEIGHT and font size scale. Each chip's width and which row it
# wraps onto are fixed at reference sizing, from the page width and the
# chips' own text lengths, so a scaled-up band never runs chips past margin.

_BANDS_LABEL_W = int(Inches(1.6))
_BANDS_CHIP_H = int(Inches(0.62))
_BANDS_CHIP_GAP = int(Inches(0.16))
_BANDS_PAD = int(Inches(0.14))
_BANDS_SECTION_GAP = int(Inches(0.20))


def _slide_bands(model, accent: RGBColor) -> bytes:
    from modules.org_chart_render import DISCIPLINE_COLOURS

    prs, slide = _new_deck()
    _title_block(slide, model)

    band_left = int(_STYLE_MARGIN) + _BANDS_LABEL_W
    band_right = int(_SLIDE_W - _STYLE_MARGIN)
    avail = band_right - band_left

    def chip_rows(chips):
        widths = [min(int(Inches(2.6)), max(int(Inches(1.4)),
                                            int(Inches(0.35) + Inches(0.09) * max(len(name), len(role)))))
                 for name, role, *_rest in chips]
        rows, used = [[]], 0
        for index, w in enumerate(widths):
            if rows[-1] and used + w + _BANDS_CHIP_GAP > avail - 2 * _BANDS_PAD:
                rows.append([])
                used = 0
            rows[-1].append(index)
            used += w + _BANDS_CHIP_GAP
        return rows, widths

    flow = _PFlow()

    def make_band(title, chips, fill, chip_edge=_CARD_EDGE):
        if not chips:
            return
        rows, widths = chip_rows(chips)
        band_h = 2 * _BANDS_PAD + len(rows) * _BANDS_CHIP_H + max(0, len(rows) - 1) * _BANDS_CHIP_GAP

        def draw(y_top, y_bottom, scale, chips=chips, fill=fill, chip_edge=chip_edge,
                 rows=rows, widths=widths, title=title):
            h = y_bottom - y_top
            _round(slide, Emu(band_left), Emu(int(y_top)), Emu(avail), Emu(h), fill, radius=0.06)
            _stack(slide, _STYLE_MARGIN, Emu(int(y_top)), Emu(_BANDS_LABEL_W - int(Inches(0.1))), Emu(h),
                  [(title.upper(), max(7.0, 9.0 * scale), True, _MUTED)], align=PP_ALIGN.RIGHT)
            chip_gap = int(_BANDS_CHIP_GAP * scale)
            chip_h = int(_BANDS_CHIP_H * scale)
            pad = int(_BANDS_PAD * scale)
            chip_y = y_top + pad
            for row in rows:
                x = band_left + pad
                for index in row:
                    name, role, tbc, role_colour = chips[index]
                    w = widths[index]
                    _round(slide, Emu(int(x)), Emu(int(chip_y)), Emu(int(w)), Emu(chip_h),
                          _TBC_FILL if tbc else _WHITE, line=_RED_TBC if tbc else chip_edge,
                          dashed=tbc, radius=0.14)
                    lines = [(name, max(9.0, 11.5 * scale), True, _RED_TBC if tbc else _INK),
                             (role, max(6.5, 9.0 * scale), True, _RED_TBC if tbc else role_colour)]
                    _stack(slide, Emu(int(x)), Emu(int(chip_y)), Emu(int(w)), Emu(chip_h),
                          [line for line in lines if line[0]])
                    x += w + chip_gap
                chip_y += chip_h + chip_gap

        flow.block(band_h, draw)
        flow.gap(_BANDS_SECTION_GAP)

    make_band("Client", [(_client_label(model)[0], model.client_role, False, _MUTED)],
              _CLIENT_DARK, chip_edge=_CLIENT_DARK)
    make_band("Leadership",
             [(p.name or "TBC", p.role, p.is_tbc, accent) for p in model.leadership],
             _tint(accent, 0.94))
    delivery = []
    for group in model.disciplines:
        colour = _hex(_discipline_colour_for(model, group))
        for person in group.people:
            role = (person.role if person.role.startswith(group.name)
                   else f"{person.role} · {group.name}")
            delivery.append((person.name or "TBC", role, person.is_tbc,
                            colour if person.is_lead else _GREY_TEXT))
    make_band("Delivery team", delivery, _tint(_hex(DISCIPLINE_COLOURS[1]), 0.95))
    assurance_chips = [(p.name or "TBC", p.role, p.is_tbc, _ASSURANCE_AMBER) for p in model.assurance]
    assurance_chips += [
        (group.peer_reviewer or "TBC", f"Peer review — {group.name}",
         not bool((group.peer_reviewer or "").strip()), _ASSURANCE_AMBER)
        for group in model.disciplines
    ]
    make_band("Assurance", assurance_chips, _tint(_ASSURANCE_AMBER, 0.93))

    def draw_footnote(y_top, y_bottom, scale):
        _stack(slide, _STYLE_MARGIN, Emu(int(y_top)), Emu(int(_SLIDE_W - 2 * _STYLE_MARGIN)),
              Emu(int(Inches(0.3))),
              [((("Solid reporting lines run top-down; the assurance band reviews independently "
                 "of the delivery team.") if (model.has_assurance or model.disciplines) else
                "Solid reporting lines run top-down."), max(7.0, 9.0 * scale), True, _MUTED)])
    flow.block(int(Inches(0.22)), draw_footnote)

    natural_h = flow.natural_height()
    scale = max(MIN_SCALE, min(MAX_SCALE, (AVAIL_H_EMU / natural_h) if natural_h > 0 else MAX_SCALE))
    bottom = flow.render(CONTENT_TOP_EMU, scale)
    _grow(prs, bottom + int(Inches(0.35)))
    return _save(prs)


# --- D. Classic tree -------------------------------------------------------

_TREE_BOX_W = int(Inches(2.6))
_TREE_BOX_H = int(Inches(0.62))
_TREE_COL_GAP = int(Inches(0.26))
_TREE_ROW_GAP = int(Inches(0.14))
_TREE_CONNECTOR = int(Inches(0.22))
_TREE_SECTION_GAP = int(Inches(0.22))
_TREE_ROW_TO_ROW = int(Inches(0.20))


def _tree_box(slide, x, y, w, h, person_lines, accent=None, tbc=False):
    _round(slide, Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)), _WHITE,
          line=_RED_TBC if tbc else (accent or _INK), dashed=tbc, radius=0.08)
    _stack(slide, Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)), person_lines, align=PP_ALIGN.CENTER)


def _slide_tree(model, accent: RGBColor) -> bytes:
    from modules.org_chart_render import _row_candidates, _balanced_wrap

    prs, slide = _new_deck()
    _title_block(slide, model)
    centre = int(_SLIDE_W / 2)
    # Same split as _render_cards_slide: the director level sits beside the
    # panel and gives up width for it; the discipline rows always sit below
    # the panel and keep the full-width centre.
    _panel_reserve_emu = int(_PANEL_W_EMU * MAX_SCALE) if model.disciplines else 0
    lead_centre = (int(_STYLE_MARGIN) + (AVAIL_W_EMU - _panel_reserve_emu) // 2
                  if model.disciplines else centre)

    lead_people = list(model.leadership)
    top_person = lead_people[0] if lead_people else None
    rank = lead_people[1:] + model.assurance
    rank_per_row = _balanced_wrap(len(rank), 4) if rank else 0
    rank_rows = -(-len(rank) // rank_per_row) if rank else 0
    rank_h = (rank_rows * _TREE_BOX_H + max(0, rank_rows - 1) * _TREE_ROW_GAP) if rank else 0

    director_h = _TREE_BOX_H + (_TREE_CONNECTOR if top_person is not None else 0)
    leadership_h = _TREE_BOX_H
    if top_person is not None:
        leadership_h += _TREE_CONNECTOR + _TREE_BOX_H
    if rank:
        leadership_h += _TREE_CONNECTOR + rank_h
    panel_h = (_PANEL_HEADER_EMU + len(model.disciplines) * _PANEL_ROW_EMU) if model.disciplines else 0
    panel_extra = max(0, director_h + panel_h - leadership_h)

    panel_anchor = {}

    def build_flow(per_row):
        disciplines = model.disciplines
        chunks = ([disciplines[i:i + per_row] for i in range(0, len(disciplines), per_row)]
                 if disciplines else [])
        row_infos = []
        disc_width = 0
        for chunk in chunks:
            tallest = max(len(g.people) for g in chunk)
            row_h = _TREE_CONNECTOR + tallest * _TREE_BOX_H + max(0, tallest - 1) * _TREE_ROW_GAP
            row_infos.append((chunk, row_h))
            width = len(chunk) * _TREE_BOX_W + (len(chunk) - 1) * _TREE_COL_GAP
            disc_width = max(disc_width, width)
        rank_width = (rank_per_row * _TREE_BOX_W + max(0, rank_per_row - 1) * _TREE_COL_GAP) if rank else 0
        leadership_width = max(_TREE_BOX_W, rank_width)
        lead_avail = (AVAIL_W_EMU - _panel_reserve_emu) if disciplines else AVAIL_W_EMU
        lead_ratio = (leadership_width / lead_avail) if lead_avail > 0 else 1.0
        disc_ratio = (disc_width / AVAIL_W_EMU) if disc_width > 0 else 0.0
        row_width = int(max(lead_ratio, disc_ratio) * AVAIL_W_EMU)

        flow = _PFlow()

        def draw_client(y_top, y_bottom, scale):
            w = int(_TREE_BOX_W * scale)
            h = y_bottom - y_top
            x = lead_centre - w // 2
            _tree_box(slide, x, y_top, w, h, [
                (model.client_name or "[CLIENT NAME]", max(9.0, 12.0 * scale), True,
                 _INK if model.client_name else _RED_TBC),
                (model.client_role, max(6.5, 9.5 * scale), True, _GREY_TEXT),
            ])
            if top_person is None:
                panel_anchor["y"] = y_bottom
        flow.block(_TREE_BOX_H, draw_client)

        if top_person is not None:
            def draw_conn1(y_top, y_bottom, scale):
                _connector(slide, MSO_CONNECTOR.STRAIGHT, lead_centre, y_top, lead_centre, y_bottom)
            flow.connector(_TREE_CONNECTOR, draw_conn1)

            def draw_top(y_top, y_bottom, scale):
                w = int(_TREE_BOX_W * scale)
                h = y_bottom - y_top
                x = lead_centre - w // 2
                _tree_box(slide, x, y_top, w, h, _person_lines(top_person, _INK, scale),
                         accent=accent, tbc=top_person.is_tbc)
                panel_anchor["y"] = y_top
            flow.block(_TREE_BOX_H, draw_top)

        if rank:
            def draw_conn2(y_top, y_bottom, scale):
                _connector(slide, MSO_CONNECTOR.STRAIGHT, lead_centre, y_top, lead_centre, y_bottom)
            flow.connector(_TREE_CONNECTOR, draw_conn2)

            def draw_rank(y_top, y_bottom, scale):
                row_h = int(_TREE_BOX_H * scale)
                row_gap = int(_TREE_ROW_GAP * scale)
                col_gap = int(_TREE_COL_GAP * scale)
                col_w = int(_TREE_BOX_W * scale)
                row_chunks = [rank[i:i + rank_per_row] for i in range(0, len(rank), rank_per_row)]
                ry = y_top
                for row_chunk in row_chunks:
                    total = len(row_chunk) * col_w + (len(row_chunk) - 1) * col_gap
                    rx = lead_centre - total // 2
                    for person in row_chunk:
                        _tree_box(slide, rx, ry, col_w, row_h, _person_lines(person, _INK, scale),
                                 tbc=person.is_tbc)
                        rx += col_w + col_gap
                    ry += row_h + row_gap
            flow.block(rank_h, draw_rank)

        if disciplines:
            def draw_elbow(y_top, y_bottom, scale):
                # Reconciles the director level's lead_centre with the
                # discipline rows' full-width centre -- see the identical
                # comment in _render_cards_slide.
                if lead_centre != centre:
                    y_mid = (y_top + y_bottom) // 2
                    _connector(slide, MSO_CONNECTOR.STRAIGHT, lead_centre, y_top, lead_centre, y_mid)
                    _connector(slide, MSO_CONNECTOR.STRAIGHT, lead_centre, y_mid, centre, y_mid)
                    _connector(slide, MSO_CONNECTOR.STRAIGHT, centre, y_mid, centre, y_bottom)
            flow.connector(_TREE_SECTION_GAP + panel_extra, draw_elbow)
            for row_index, (chunk, row_h) in enumerate(row_infos):
                def draw_row(y_top, y_bottom, scale, chunk=chunk):
                    col_w = int(_TREE_BOX_W * scale)
                    col_gap = int(_TREE_COL_GAP * scale)
                    total = len(chunk) * col_w + (len(chunk) - 1) * col_gap
                    x0 = centre - total // 2
                    centres = [x0 + col_w // 2 + i * (col_w + col_gap) for i in range(len(chunk))]
                    bus_y = y_top + int(_TREE_CONNECTOR * scale * 0.4)
                    _connector(slide, MSO_CONNECTOR.STRAIGHT, centre, y_top, centre, bus_y)
                    if len(centres) > 1:
                        _connector(slide, MSO_CONNECTOR.STRAIGHT, centres[0], bus_y, centres[-1], bus_y)
                    box_h = int(_TREE_BOX_H * scale)
                    row_gap = int(_TREE_ROW_GAP * scale)
                    box_top = bus_y + int(_TREE_CONNECTOR * scale * 0.6)
                    for i, group in enumerate(chunk):
                        cx = centres[i]
                        _connector(slide, MSO_CONNECTOR.STRAIGHT, cx, bus_y, cx, box_top)
                        by = box_top
                        for person in group.people:
                            _tree_box(slide, cx - col_w // 2, by, col_w, box_h,
                                     _person_lines(person, _INK, scale), tbc=person.is_tbc)
                            by += box_h + row_gap
                flow.block(row_h, draw_row)
                if row_index < len(row_infos) - 1:
                    flow.gap(_TREE_ROW_TO_ROW)

        return flow, row_width

    per_row_candidates = _row_candidates(len(model.disciplines)) if model.disciplines else [1]
    flow, scale, per_row = _psolve_scale(build_flow, per_row_candidates)
    bottom = flow.render(CONTENT_TOP_EMU, scale)

    if model.disciplines:
        panel_w = int(_PANEL_W_EMU * scale)
        panel_top = panel_anchor.get("y", CONTENT_TOP_EMU + int(director_h * scale))
        panel_x = int(_SLIDE_W - _STYLE_MARGIN) - panel_w
        _draw_peer_review_panel(slide, model, scale, panel_x, panel_top, _PANEL_W_EMU)

    _grow(prs, bottom + int(Inches(0.35)))
    return _save(prs)


_STYLE_RENDERERS = {
    "cards": _slide_cards,
    "columns": _slide_columns,
    "bands": _slide_bands,
    "tree": _slide_tree,
}


def _empty_slide(model) -> bytes:
    prs, slide = _new_deck()
    y = _title_block(slide, model)
    _stack(slide, _STYLE_MARGIN, y + Inches(0.4),
           Emu(int(_SLIDE_W - 2 * _STYLE_MARGIN)), Inches(0.6),
           [(_ORG_EMPTY_NOTE, 14, False, _RED_TBC)], align=PP_ALIGN.CENTER)
    return _save(prs)


_ORG_EMPTY_NOTE = ("[NO TEAM ASSIGNED -- add the management roles and discipline leads in "
                   "the Team & Resourcing tab, then re-download this PowerPoint]")


def populate_org_chart(resource_plan: list, client_name: str = "", project_name: str = "",
                       tender_name: str = "", theme_name: str | None = None,
                       style: str | None = None) -> bytes:
    """
    Builds a fresh .pptx (returned as bytes) of the project organisation
    chart, in the user's chosen presentation style.

    `style` is one of org_chart_render.STYLES; anything else (including None)
    falls back to org_chart_render.DEFAULT_STYLE, so the deck always matches
    the style the preview drew.

    Everything comes from `resource_plan` (a list of
    resourcing.ResourceAssignment). An unassigned slot is drawn as a dashed
    red TBC; a role the user REMOVED (see resourcing.OPTIONAL_MANAGEMENT_ROLES)
    is absent from the plan and so absent from the chart, with no TBC --
    deliberate absence is not a gap. An assurance/reviewer element appears
    only where the plan actually holds such a slot; this module never adds
    one. An empty plan renders a placeholder slide rather than a bare chart.
    """
    from modules import org_chart_render

    model = org_chart_render.build_model(resource_plan, client_name, project_name,
                                         tender_name)
    if model.is_empty:
        return _empty_slide(model)
    resolved = org_chart_render.effective_style(
        model, style or org_chart_render.DEFAULT_STYLE)
    accent = _resolve_palette(theme_name)["header"]
    return _STYLE_RENDERERS[resolved](model, accent)
