"""
methodology_pptx.py

Builds the methodology-table PowerPoint FROM SCRATCH, straight from a
project's tender analysis -- no template file to edit or keep in sync
(same approach as org_chart_pptx.py, for the same reasons: nothing to
delete, nothing to leave dangling).

FOUR PRESENTATION STYLES
-------------------------
populate_methodology(..., style=...) dispatches to one of four independent
slide builders (methodology_render.STYLES: "matrix" | "chevrons" |
"programme" | "spine" -- see that module's docstring for what each looks
like and why a separate matplotlib renderer draws the same content for the
UI's live preview). All four read the SAME column list -- one dict per
stage, built once by methodology_render.build_columns() -- so a style
picked in the UI and the exported deck can never describe different
content, only a different layout of it.

Column content, regardless of style: when the reviewed
methodology_stages.MethodologyStage grid exists, every column's name, key
tasks, engagement activities, outcome, deliverables and date chevron come
from it, real content only, with the literal "TBC" wherever the brief
didn't support a cell. Before that grid has been generated/filled in, the
table falls back to exactly what it always produced: one column of
standard Project Initiation boilerplate, one built from the brief's real
scope items (tender_analyser.ScopeItem) -- never invented, and two (or,
with a real stage grid, as many as the brief needs) of explicit
"[CONFIRM ...]" placeholders rather than guessed content.

Every style guards against overflow the same way: text shrinks first
(_fit_size), and only once it hits a legible floor does a cell drop items,
always with an honest "+N more -- see full methodology" line -- silent
truncation is never allowed to look like a complete table.

The client name (only shown by the "matrix" style's legend) and the
project name (every style's title) are the pieces of real, project-specific
data these slides need beyond the stage content -- pulled straight from
Project Setup, never invented; shown red if not entered yet, matching
every other missing-field convention in this app.

Colours: the "matrix" style is pulled from the SAME theme palette the rest
of the proposal uses (divider_designer.THEME_COLOURS), so it never looks
like a foreign template dropped into an otherwise-themed document. The
other three styles use the fixed, colourblind-validated stage-colour order
(methodology_render.STAGE_COLOURS) instead -- the same four hues the org
chart and delivery program use -- since their whole visual structure is
built around telling stages apart by colour, and stage identity has to
stay consistent with the rest of the pack, not with whichever theme this
particular proposal happens to use.
"""

from __future__ import annotations

import io
import math

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from modules.divider_designer import THEME_COLOURS

_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
_RED = RGBColor(0xC0, 0x00, 0x00)
_FONT = "Calibri"

# A4 landscape
_SLIDE_W = Inches(11.6929)
_SLIDE_H = Inches(8.2677)

# The reviewed-stages/legacy-fallback column content itself now lives in
# methodology_render.py (build_columns() and friends) so every style here
# and every live-preview PNG there read the identical data -- this is just
# the one fallback label still needed directly in this module.
_CONFIRM_DATE_RANGE = "[Date range]"


def _tint(rgb: RGBColor, amount: float) -> RGBColor:
    """Mixes `rgb` toward white by `amount` (0 = original colour, 1 = white)."""
    def mix(c):
        return round(c + (255 - c) * amount)
    return RGBColor(mix(rgb[0]), mix(rgb[1]), mix(rgb[2]))


def _resolve_palette(theme_name: str | None) -> dict:
    colours = THEME_COLOURS.get(theme_name, THEME_COLOURS["Corporate"])
    primary = RGBColor(*colours["primary"])
    accent = RGBColor(*colours["accent"])
    # Minimalist's "primary" is a light near-white wash by design (its reversed
    # light-on-dark convention elsewhere in the app) -- same special case
    # export_docx._theme_colours() applies to heading colour.
    dark_role = RGBColor(0x2A, 0x2A, 0x2A) if theme_name == "Minimalist" else primary
    return {
        "header_bg": dark_role,
        "tasks_bg": _tint(primary, 0.88),
        "eng_bg": _tint(primary, 0.78),
        "outcome_deliv_bg": _tint(accent, 0.82),
        "icon_fill": dark_role,
        "engagement_fill": accent,
        "chevron_bg": dark_role,
    }


# ---------------------------------------------------------------------------
# Icons -- rendered with PIL, entirely self-contained (no bundled asset files
# to keep in sync with this module, same reasoning as org_chart_pptx.py).
# ---------------------------------------------------------------------------

def _render_icon_png(kind: str, size: int = 240) -> bytes:
    ss = 4
    big = size * ss
    im = Image.new("RGBA", (big, big), (0, 0, 0, 0))
    d = ImageDraw.Draw(im)
    pad = int(big * 0.03)
    d.ellipse([pad, pad, big - pad, big - pad], fill=(0, 0, 0, 255))
    cx = cy = big / 2

    if kind == "hold_point":
        body_w, body_h = big * 0.375, big * 0.475
        bx0, by0 = cx - body_w / 2, cy - body_h / 2 + big * 0.03
        bx1, by1 = cx + body_w / 2, cy + body_h / 2 + big * 0.03
        d.rounded_rectangle([bx0, by0, bx1, by1], radius=big * 0.035, outline=(255, 255, 255, 255), width=int(big * 0.0225))
        clip_w, clip_h = big * 0.16, big * 0.065
        d.rounded_rectangle([cx - clip_w / 2, by0 - clip_h / 2, cx + clip_w / 2, by0 + clip_h / 2],
                             radius=big * 0.025, fill=(255, 255, 255, 255))
        sq = big * 0.035
        line_x0 = bx0 + big * 0.055
        line_x1 = bx1 - big * 0.05
        for frac in (0.24, 0.43, 0.62):
            ry = by0 + body_h * frac
            d.rectangle([line_x0, ry - sq / 2, line_x0 + sq, ry + sq / 2], outline=(255, 255, 255, 255), width=int(big * 0.0125))
            d.line([line_x0 + sq + big * 0.03, ry, line_x1, ry], fill=(255, 255, 255, 255), width=int(big * 0.0225))
    else:  # collaborative engagement -- hub and spoke
        hub_r, node_r, spoke_len = big * 0.065, big * 0.04, big * 0.23
        n_nodes = 5
        for i in range(n_nodes):
            angle = math.radians(-90 + i * (360 / n_nodes))
            nx, ny = cx + spoke_len * math.cos(angle), cy + spoke_len * math.sin(angle)
            d.line([cx, cy, nx, ny], fill=(255, 255, 255, 255), width=int(big * 0.025))
        for i in range(n_nodes):
            angle = math.radians(-90 + i * (360 / n_nodes))
            nx, ny = cx + spoke_len * math.cos(angle), cy + spoke_len * math.sin(angle)
            d.ellipse([nx - node_r, ny - node_r, nx + node_r, ny + node_r], fill=(255, 255, 255, 255))
        d.ellipse([cx - hub_r, cy - hub_r, cx + hub_r, cy + hub_r], fill=(255, 255, 255, 255))

    im = im.resize((size, size), Image.LANCZOS)
    buf = io.BytesIO()
    im.save(buf, format="PNG")
    buf.seek(0)
    return buf.read()


# ---------------------------------------------------------------------------
# Shape/text helpers (Emu-precise, mirroring org_chart_pptx.py's conventions)
# ---------------------------------------------------------------------------

def _rect(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    return box, tf


def _lines_to_paragraphs(tf, lines, size_pt, color, bullet=True, align=PP_ALIGN.LEFT):
    for i, item in enumerate(lines):
        no_bullet = False
        text = item
        if isinstance(item, tuple):
            text, no_bullet = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        prefix = "" if (no_bullet or not bullet or not text) else "–  "
        run.text = f"{prefix}{text}"
        run.font.size = Pt(size_pt)
        run.font.name = _FONT
        run.font.color.rgb = color


def _placeholder_bullets(tf, text, size_pt):
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = f"–  {text}"
    run.font.size = Pt(size_pt)
    run.font.name = _FONT
    run.font.italic = True
    run.font.color.rgb = _RED


def _centered_text(slide, x, y, w, h, text, size_pt, color, bold=True, italic=False):
    box, tf = _textbox(slide, x, y, w, h)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = _FONT
    run.font.color.rgb = color
    return box


def _rotated_label(slide, cx, cy, visual_len, visual_thick, text, size_pt, color):
    w, h = visual_len, visual_thick
    x, y = Emu(int(cx - w / 2)), Emu(int(cy - h / 2))
    box, tf = _textbox(slide, x, y, w, h)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size_pt)
        run.font.bold = True
        run.font.name = _FONT
        run.font.color.rgb = color
    box.rotation = 270
    return box


def _icon(slide, png_bytes, cx, cy, d):
    x, y = Emu(int(cx - d / 2)), Emu(int(cy - d / 2))
    slide.shapes.add_picture(io.BytesIO(png_bytes), x, y, d, d)


def _chevron(slide, x, y, w, h, fill_color, text, size_pt):
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(10)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = True
    run.font.name = _FONT
    run.font.color.rgb = _WHITE
    return shape


def _round_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_w=None, radius=0.08):
    """A rounded-rectangle panel -- used by the chevrons/programme/spine
    styles for cards, chips and bands. `fill_color`/`line_color` of None
    leaves that side unset (transparent fill, or no border)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        shape.adjustments[0] = radius
    except Exception:
        pass
    if fill_color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color is not None:
        shape.line.color.rgb = line_color
        shape.line.width = line_w or Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _chip(slide, x, y, w, h, text, fill_color, text_color, size_pt=7.0):
    """One small rounded "pill" with centred bold text -- the
    deliverable-as-chip look the chevrons/programme/spine styles use."""
    shape = _round_rect(slide, x, y, w, h, fill_color=fill_color, radius=0.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(3)
    tf.margin_right = Pt(3)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = True
    run.font.name = _FONT
    run.font.color.rgb = text_color
    return shape


def _diamond(slide, cx, cy, r, fill_color, text=None, text_size=6.0):
    """A small filled diamond -- the Programme-matched style's hold-point
    gate marker, with an optional label centred beneath it."""
    shape = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Emu(int(cx - r)), Emu(int(cy - r)), Emu(int(2 * r)), Emu(int(2 * r)))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = _WHITE
    shape.line.width = Pt(1)
    shape.shadow.inherit = False
    if text:
        _centered_text(slide, Emu(int(cx - r * 1.8)), Emu(int(cy + r * 1.1)), Emu(int(r * 3.6)), Emu(int(r * 1.8)),
                        text, text_size, fill_color, bold=True)
    return shape


def _stage_colour_rgb(index: int) -> "RGBColor":
    return _hex_to_rgb(stage_colour(index))


def _hex_to_rgb(hex_colour: str) -> "RGBColor":
    hex_colour = hex_colour.lstrip("#")
    return RGBColor(int(hex_colour[0:2], 16), int(hex_colour[2:4], 16), int(hex_colour[4:6], 16))


# ---------------------------------------------------------------------------
# Content assembly -- the shared column model lives in methodology_render.py
# now (build_columns() etc.), so this module's four style builders and that
# module's four PNG previews can never describe a different table. See that
# module's docstring for the full "why one model, two renderers" reasoning.
# ---------------------------------------------------------------------------

from modules.methodology_render import (  # noqa: E402
    DEFAULT_STYLE,
    STYLES,
    _cap_items,
    build_columns,
    is_placeholder as _is_placeholder,
    stage_carries_hold_point,
    stage_colour,
)


def _natural_lines(texts: list[str], width_emu, size_pt: float) -> int:
    """How many wrapped lines `texts` need at `size_pt` in a box `width_emu`
    wide, with NO shrinking or dropping -- the "how tall does this content
    actually want to be" half of the two-pass layout _row_heights() does."""
    chars_per_line = max(8, int((width_emu / 914400) * (1.85 / (size_pt / 72))))
    return sum(max(1, -(-len(str(t)) // chars_per_line)) for t in texts)


def _fit_size(lines, width_emu, height_emu, start_pt: float, min_pt: float = 4.2):
    """Shrink text until it fits the box, then cap what still doesn't.

    Returns (size_pt, lines). The old code rendered every list at a fixed
    5.6pt inside a fixed-height text box, so a brief with ten scope items
    simply spilled its last items out of the bottom of the box -- invisibly,
    because PowerPoint clips nothing and the shapes below just overlapped
    them. Shrinking first preserves the content; capping with an explicit
    "+N more" line at the floor is honest about the rest, which silent
    overflow never was."""
    texts = [t if not isinstance(t, tuple) else t[0] for t in lines]
    size = start_pt
    while size > min_pt:
        line_h = (size * 1.28) / 72 * 914400
        needed = _natural_lines(texts, width_emu, size)
        if needed * line_h <= height_emu:
            return size, lines
        size -= 0.2

    chars_per_line = max(8, int((width_emu / 914400) * (1.85 / (min_pt / 72))))
    line_h = (min_pt * 1.28) / 72 * 914400
    budget = max(1, int(height_emu // line_h)) - 1  # leave a row for the note
    kept, used = [], 0
    for item, text in zip(lines, texts):
        cost = max(1, -(-len(str(text)) // chars_per_line))
        if used + cost > budget:
            break
        kept.append(item)
        used += cost
    dropped = len(lines) - len(kept)
    if dropped:
        kept.append((f"+{dropped} more — see full methodology", True))
    return min_pt, kept


def _fit_single(text: str, width_emu, height_emu, start_pt: float, min_pt: float = 4.5) -> float:
    """Same shrink-to-fit as _fit_size, for a single centred text cell
    (OUTCOME) rather than a bulleted list -- there is nothing sensible to
    drop from one sentence, so this only ever shrinks, never truncates."""
    size = start_pt
    while size > min_pt:
        line_h = (size * 1.28) / 72 * 914400
        if _natural_lines([text], width_emu, size) * line_h <= height_emu:
            return size
        size -= 0.2
    return min_pt


def _render_cell_lines(slide, tf, lines, size_pt):
    """One list cell. Placeholder/TBC entries render red italic; real
    content renders as ordinary bullets, so a mixed cell shows at a glance
    which half still needs work."""
    for i, item in enumerate(lines):
        no_bullet = False
        text = item
        if isinstance(item, tuple):
            text, no_bullet = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        prefix = "" if (no_bullet or not text) else "\u2013  "
        run.text = f"{prefix}{text}"
        run.font.size = Pt(size_pt)
        run.font.name = _FONT
        placeholder = _is_placeholder(text)
        run.font.italic = placeholder
        run.font.color.rgb = _RED if placeholder else _DARK_TEXT


WVR_STATEMENT = "All design deliverables will be issued with completed Work Verification Records (WVRs)"
WVR_CONFIRM_PLACEHOLDER = "[CONFIRM WVR / QA STATEMENT FOR THIS FIRM]"


def _slide_matrix(slide, columns_data, from_stages, P, hold_icon, eng_icon,
                  project_name, client_name, wvr_confirmed) -> None:
    """Boardroom matrix -- the original/default style: navy stage headers,
    LEFT row labels (KEY TASKS / ENGAGEMENT / OUTCOME / DELIVERABLES), one
    grid row per cell type. Row heights are CONTENT-sized: each row's real
    height is measured from what its busiest column actually needs, and
    only rescaled (proportionally, never below a legible floor) if all four
    rows together would otherwise overflow the slide -- so a thin
    ENGAGEMENT row no longer reserves the same fixed height as a busy
    DELIVERABLES row, and a busy row is never silently clipped."""
    n = len(columns_data)

    M = Inches(0.14)
    row_label_x, row_label_w = M, Inches(0.17)
    content_x = Emu(int(row_label_x + row_label_w + Inches(0.04)))
    right_edge = Emu(int(_SLIDE_W - M))
    gap = Inches(0.20)

    if from_stages:
        # Every column is a real stage now, so they share the width evenly
        # rather than reserving a narrow first column for boilerplate.
        col_w = Emu(int((right_edge - content_x - (n - 1) * gap) / max(n, 1)))
        cols = [(Emu(int(content_x + i * (col_w + gap))), col_w) for i in range(n)]
    else:
        col1_w = Inches(1.55)
        col2_x = Emu(int(content_x + col1_w + Inches(0.05)))
        stage_w = Emu(int((right_edge - col2_x - 2 * gap) / 3))
        cols = [(content_x, col1_w)] + [
            (Emu(int(col2_x + i * (stage_w + gap))), stage_w) for i in range(3)
        ]
        cols = cols[:n]
    col3_x = cols[2][0] if n > 2 else right_edge
    col4_x = cols[3][0] if n > 3 else right_edge

    y_top = M
    # Slightly taller than it used to be: the title block now carries the
    # project name on a second line beneath the heading. Total column height
    # still lands inside the A4 landscape slide.
    h_title = Inches(0.40)
    y_header = Emu(int(y_top + h_title + Inches(0.04)))
    h_header = Inches(0.46)
    h_timeline = Inches(0.30)

    # ---- content-sized row heights -----------------------------------
    # Measure what each row's busiest column actually needs (at the row's
    # preferred font size, nothing dropped yet), then either use those
    # natural heights directly -- if all four fit in the space a fixed
    # layout used to reserve -- or scale every row down proportionally
    # (never below a legible floor) so the table still ends inside the
    # slide. _fit_size (called per-cell below, with the ALLOCATED height)
    # is what guarantees no individual cell overflows even after scaling.
    note_h = Inches(0.2)  # the WVR confirm/assert line under deliverables
    available_h = Inches(3.35 + 0.70 + 0.50 + 2.30)  # same total budget the old fixed layout used
    row_specs = [("tasks", 5.6, Inches(0.08)), ("engagement", 5.4, Inches(0.06)),
                ("outcome", 5.6, Inches(0.04)), ("deliverables", 5.0, Inches(0.03) + note_h + Inches(0.06))]
    floors = {"tasks": Inches(0.55), "engagement": Inches(0.35), "outcome": Inches(0.32), "deliverables": Inches(0.75)}

    natural = {}
    for key, start_pt, pad in row_specs:
        tallest = 0
        for (cx, cw), column in zip(cols, columns_data):
            box_w = Emu(int(cw - Inches(0.1)))
            value = column[key]
            texts = [str(value)] if key == "outcome" else [
                t if not isinstance(t, tuple) else t[0] for t in value]
            line_h = (start_pt * 1.28) / 72 * 914400
            tallest = max(tallest, _natural_lines(texts, box_w, start_pt) * line_h)
        natural[key] = int(tallest + pad)

    # The floor isn't just "don't get illegibly short" -- each row's own
    # rotated side-label (e.g. two-line "KEY ENGAGEMENT\nACTIVITIES" at
    # 5.2pt) needs a minimum height to render inside without spilling
    # outside its box once rotated. A row with very thin cell content (a
    # single short engagement line) previously fell straight through to
    # that tiny natural height on the "everything fits" path below, which
    # only applied floors on the scale-down branch -- so the label text
    # overflowed sideways into the next row's column. Applying the floor
    # here, before summing, means both branches respect it.
    natural = {k: max(floors[k], v) for k, v in natural.items()}

    total_natural = sum(natural.values())
    if total_natural <= available_h:
        row_h = natural
    else:
        scale = available_h / total_natural
        row_h = {k: max(floors[k], int(v * scale)) for k, v in natural.items()}

    h_tasks, h_eng, h_outcome, h_deliv = (
        Emu(row_h["tasks"]), Emu(row_h["engagement"]), Emu(row_h["outcome"]), Emu(row_h["deliverables"]))
    y_tasks = Emu(int(y_header + h_header))
    y_eng = Emu(int(y_tasks + h_tasks))
    y_outcome = Emu(int(y_eng + h_eng))
    y_deliv = Emu(int(y_outcome + h_outcome))
    y_timeline = Emu(int(y_deliv + h_deliv + Inches(0.03)))

    # ---- title + KEY legend --------------------------------------------
    # Heading plus the project name, in ONE box so the two lines can't
    # drift into the stage headers below. project_name was accepted by this
    # function and never rendered anywhere, so every exported methodology
    # table was anonymous -- indistinguishable from another bid's once saved
    # to disk.
    title_box, title_tf = _textbox(slide, M, y_top, Inches(7.5), h_title)
    title_tf.vertical_anchor = MSO_ANCHOR.TOP
    tp = title_tf.paragraphs[0]
    tp.alignment = PP_ALIGN.LEFT
    tr = tp.add_run()
    tr.text = "Our proposed methodology"
    tr.font.size = Pt(15)
    tr.font.bold = True
    tr.font.name = _FONT
    tr.font.color.rgb = _DARK_TEXT
    sp = title_tf.add_paragraph()
    sp.alignment = PP_ALIGN.LEFT
    sr = sp.add_run()
    sr.text = (project_name or "").strip() or "[Insert project name]"
    sr.font.size = Pt(7.5)
    sr.font.name = _FONT
    sr.font.color.rgb = _DARK_TEXT if (project_name or "").strip() else _RED

    key_box, key_tf = _textbox(slide, Emu(int(right_edge - Inches(2.7))), Emu(int(y_top + Inches(0.02))), Inches(0.55), Emu(int(h_title * 0.5)))
    key_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    kr = key_tf.paragraphs[0].add_run()
    kr.text = "KEY"
    kr.font.size = Pt(8)
    kr.font.bold = True
    kr.font.name = _FONT
    kr.font.color.rgb = _DARK_TEXT

    _icon(slide, hold_icon, Emu(int(right_edge - Inches(2.1))), Emu(int(y_top + Inches(0.06) + Inches(0.085))), Inches(0.17))
    label_box, label_tf = _textbox(slide, Emu(int(right_edge - Inches(1.93))), y_top, Inches(1.93), Emu(int(h_title * 0.5)))
    label_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = label_tf.paragraphs[0]
    client_display = (client_name or "").strip() or "[Insert client name]"
    r1 = p1.add_run()
    r1.text = client_display
    r1.font.size = Pt(6.5)
    r1.font.name = _FONT
    r1.font.color.rgb = _DARK_TEXT if (client_name or "").strip() else _RED
    r2 = p1.add_run()
    r2.text = " hold point"
    r2.font.size = Pt(6.5)
    r2.font.name = _FONT
    r2.font.color.rgb = _DARK_TEXT

    _icon(slide, eng_icon, Emu(int(right_edge - Inches(2.1))), Emu(int(y_top + Inches(0.24) + Inches(0.085))), Inches(0.17))
    _centered_text(slide, Emu(int(right_edge - Inches(1.93))), Emu(int(y_top + Inches(0.18))), Inches(1.9), Emu(int(h_title * 0.5)),
                    "Collaborative engagement", 6.5, _DARK_TEXT, bold=False)
    slide.shapes[-1].text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # ---- stage headers ----------------------------------------------------
    for (cx, cw), column in zip(cols, columns_data):
        _rect(slide, cx, y_header, cw, h_header, P["header_bg"])
        _centered_text(slide, cx, y_header, cw, h_header, column["name"], 9.5,
                       _RED if _is_placeholder(column["name"]) else _WHITE)

    # ---- KEY TASKS ---------------------------------------------------------
    _rect(slide, row_label_x, y_tasks, row_label_w, h_tasks, P["tasks_bg"])
    for (cx, cw), column in zip(cols, columns_data):
        _rect(slide, cx, y_tasks, cw, h_tasks, P["tasks_bg"])
        box_w = Emu(int(cw - Inches(0.1)))
        box_h = Emu(int(h_tasks - Inches(0.08)))
        box, tf = _textbox(slide, Emu(int(cx + Inches(0.05))), Emu(int(y_tasks + Inches(0.04))), box_w, box_h)
        size, lines = _fit_size(column["tasks"], box_w, box_h, 5.6)
        _render_cell_lines(slide, tf, lines, size)
    _rotated_label(slide, Emu(int(row_label_x + row_label_w / 2)), Emu(int(y_tasks + h_tasks / 2)), Emu(int(h_tasks - Inches(0.1))), row_label_w, "KEY TASKS", 6.5, _DARK_TEXT)
    if n > 2:
        _icon(slide, hold_icon, Emu(int(col3_x - gap / 2)), Emu(int(y_tasks + h_tasks - Inches(0.13))), Inches(0.15))
    if n > 3:
        _icon(slide, hold_icon, Emu(int(col4_x - gap / 2)), Emu(int(y_tasks + h_tasks - Inches(0.13))), Inches(0.15))

    # ---- KEY ENGAGEMENT ACTIVITIES -----------------------------------------
    _rect(slide, row_label_x, y_eng, row_label_w, h_eng, P["eng_bg"])
    for (cx, cw), column in zip(cols, columns_data):
        _rect(slide, cx, y_eng, cw, h_eng, P["eng_bg"])
        box_w = Emu(int(cw - Inches(0.2)))
        box_h = Emu(int(h_eng - Inches(0.06)))
        box, tf = _textbox(slide, Emu(int(cx + Inches(0.05))), Emu(int(y_eng + Inches(0.03))), box_w, box_h)
        size, lines = _fit_size(column["engagement"], box_w, box_h, 5.4)
        _render_cell_lines(slide, tf, lines, size)
        _icon(slide, eng_icon, Emu(int(cx + cw - Inches(0.11))), Emu(int(y_eng + h_eng / 2)), Inches(0.16))
    _rotated_label(slide, Emu(int(row_label_x + row_label_w / 2)), Emu(int(y_eng + h_eng / 2)), Emu(int(h_eng - Inches(0.06))), row_label_w, "KEY ENGAGEMENT\nACTIVITIES", 5.2, _DARK_TEXT)

    # ---- OUTCOME ------------------------------------------------------------
    _rect(slide, row_label_x, y_outcome, row_label_w, h_outcome, P["outcome_deliv_bg"])
    for (cx, cw), column in zip(cols, columns_data):
        _rect(slide, cx, y_outcome, cw, h_outcome, P["outcome_deliv_bg"])
        text = column["outcome"]
        placeholder = _is_placeholder(text)
        box_w = Emu(int(cw - Inches(0.1)))
        box_h = Emu(int(h_outcome - Inches(0.04)))
        fitted_size = _fit_single(str(text), box_w, box_h, 5.6)
        _centered_text(slide, Emu(int(cx + Inches(0.05))), Emu(int(y_outcome + Inches(0.02))), box_w, box_h,
                        text, fitted_size,
                        _RED if placeholder else _DARK_TEXT, bold=not placeholder, italic=placeholder)
    _rotated_label(slide, Emu(int(row_label_x + row_label_w / 2)), Emu(int(y_outcome + h_outcome / 2)), Emu(int(h_outcome - Inches(0.05))), row_label_w, "OUTCOME", 5.6, _DARK_TEXT)

    # ---- DELIVERABLES ---------------------------------------------------------
    _rect(slide, row_label_x, y_deliv, row_label_w, h_deliv, P["outcome_deliv_bg"])
    list_h = Emu(int(h_deliv - note_h - Inches(0.06)))
    for (cx, cw), column in zip(cols, columns_data):
        _rect(slide, cx, y_deliv, cw, h_deliv, P["outcome_deliv_bg"])
        box_w = Emu(int(cw - Inches(0.1)))
        box, tf = _textbox(slide, Emu(int(cx + Inches(0.05))), Emu(int(y_deliv + Inches(0.03))), box_w, list_h)
        size, lines = _fit_size(column["deliverables"], box_w, list_h, 5.0)
        _render_cell_lines(slide, tf, lines, size)
        # The WVR line asserts a QA practice about the BIDDER. The app was
        # never told whether this firm issues Work Verification Records, so
        # printing it as fact in every export was the one no-invention
        # breach in this module. It now only appears once confirmed.
        note_box, note_tf = _textbox(slide, Emu(int(cx + Inches(0.05))), Emu(int(y_deliv + list_h + Inches(0.02))), box_w, note_h)
        nr = note_tf.paragraphs[0].add_run()
        nr.text = WVR_STATEMENT if wvr_confirmed else WVR_CONFIRM_PLACEHOLDER
        nr.font.size = Pt(5.0)
        nr.font.italic = True
        nr.font.name = _FONT
        nr.font.color.rgb = _DARK_TEXT if wvr_confirmed else _RED
    _rotated_label(slide, Emu(int(row_label_x + row_label_w / 2)), Emu(int(y_deliv + h_deliv / 2)), Emu(int(h_deliv - Inches(0.1))), row_label_w, "DELIVERABLES", 6.5, _DARK_TEXT)
    if n > 2:
        _icon(slide, hold_icon, Emu(int(col3_x - gap / 2)), Emu(int(y_deliv + h_deliv * 0.12)), Inches(0.15))
    if n > 3:
        _icon(slide, hold_icon, Emu(int(col4_x - gap / 2)), Emu(int(y_deliv + h_deliv * 0.12)), Inches(0.15))

    # ---- TIMELINE ---------------------------------------------------------
    for (cx, cw), column in zip(cols, columns_data):
        chevron_text = column["chevron"]
        _chevron(slide, cx, y_timeline, cw, h_timeline, P["chevron_bg"],
                 chevron_text if chevron_text else _CONFIRM_DATE_RANGE, 5.4)
    if n > 2:
        _icon(slide, hold_icon, Emu(int(col3_x - gap / 2)), Emu(int(y_timeline + h_timeline / 2)), Inches(0.16))
    if n > 3:
        _icon(slide, hold_icon, Emu(int(col4_x - gap / 2)), Emu(int(y_timeline + h_timeline / 2)), Inches(0.16))


def _simple_title(slide, project_name: str) -> None:
    """The compact title the three new styles use -- just the heading and
    project name, no KEY legend (the matrix style's hold-point/engagement
    icon legend doesn't apply to these layouts, which show hold points and
    engagement inline instead)."""
    M = Inches(0.14)
    box, tf = _textbox(slide, M, M, Inches(9.0), Inches(0.32))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "Our proposed methodology"
    if (project_name or "").strip():
        r.text += f" — {project_name.strip()}"
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.name = _FONT
    r.font.color.rgb = _DARK_TEXT


def _deliverable_chips(slide, x, w, y, avail_h, items, fill, text_fill, size_pt=6.5) -> None:
    """Deliverables as flowing rounded chips within a fixed box, capped
    (with an honest "+N more" line) rather than silently overflowing --
    same convention _fit_size uses for bulleted cells."""
    row_h = Inches(0.22)
    max_rows = max(1, int(avail_h / row_h))
    chip_w_estimate = Inches(1.1)
    chips_per_row = max(1, int(w / (chip_w_estimate + Inches(0.06))))
    max_items = max(1, max_rows * chips_per_row)
    kept, dropped = _cap_items(items, max_items)

    cx, row_top, used_rows = x, y, 0
    for item in kept:
        chip_w = Emu(int(min(w, Inches(0.12) + Pt(size_pt) * len(str(item)) * 9525 * 0.62)))
        if cx + chip_w > x + w and cx > x:
            cx = x
            row_top = Emu(int(row_top + row_h + Inches(0.03)))
            used_rows += 1
            if used_rows >= max_rows:
                break
        _chip(slide, Emu(int(cx)), Emu(int(row_top)), chip_w, row_h, str(item), fill, text_fill, size_pt)
        cx = Emu(int(cx + chip_w + Inches(0.06)))
    if dropped:
        note_y = Emu(int(row_top + row_h + Inches(0.02)))
        box, tf = _textbox(slide, x, note_y, w, Inches(0.16))
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"+{dropped} more — see full methodology"
        r.font.size = Pt(5.6)
        r.font.italic = True
        r.font.name = _FONT
        r.font.color.rgb = RGBColor(0x5B, 0x64, 0x72)


def _slide_chevrons(slide, columns_data, from_stages, P, hold_icon, eng_icon,
                    project_name, client_name, wvr_confirmed) -> None:
    """Stage chevrons -- an arrowed banner per stage (stage colour, name +
    weeks), a bordered card below with coloured section labels and
    deliverables as tinted chips. See the module docstring and
    methodology_render.py's for the shared reasoning."""
    n = max(len(columns_data), 1)
    _simple_title(slide, project_name)

    M = Inches(0.14)
    gap = Inches(0.12)
    content_top = Emu(int(M + Inches(0.4)))
    col_w = Emu(int((_SLIDE_W - 2 * M - (n - 1) * gap) / n))

    chevron_h = Inches(0.55)
    card_top = Emu(int(content_top + chevron_h + Inches(0.06)))
    card_h = Emu(int(_SLIDE_H - card_top - Inches(0.12)))

    # Key tasks is nearly always the busiest cell (the brief's own scope
    # items land there), so it gets the majority share; the rest split what
    # is normally much shorter content.
    tasks_h = Emu(int(card_h * 0.58))
    eng_h = Emu(int(card_h * 0.14))
    outcome_h = Emu(int(card_h * 0.10))
    deliv_h = Emu(int(card_h - tasks_h - eng_h - outcome_h - Inches(0.16)))

    for i, column in enumerate(columns_data):
        cx = Emu(int(M + i * (col_w + gap)))
        colour = _hex_to_rgb(stage_colour(i))

        label = column["name"]
        if column.get("chevron") and not _is_placeholder(column["chevron"]):
            label = f"{label}   ({column['chevron']})"
        _chevron(slide, cx, content_top, col_w, chevron_h, colour, label, 7.5)

        _round_rect(slide, cx, card_top, col_w, card_h, fill_color=None, line_color=colour, line_w=Pt(1))

        pad = Inches(0.06)
        y = Emu(int(card_top + Inches(0.05)))
        box_w = Emu(int(col_w - 2 * pad))

        for label_text, key, h in (
            ("KEY TASKS", "tasks", tasks_h), ("ENGAGEMENT", "engagement", eng_h),
        ):
            lb, lbtf = _textbox(slide, Emu(int(cx + pad)), y, box_w, Inches(0.14))
            lr = lbtf.paragraphs[0].add_run()
            lr.text = label_text
            lr.font.size = Pt(6.2)
            lr.font.bold = True
            lr.font.name = _FONT
            lr.font.color.rgb = colour
            list_top = Emu(int(y + Inches(0.15)))
            list_h = Emu(int(h - Inches(0.15)))
            box, tf = _textbox(slide, Emu(int(cx + pad)), list_top, box_w, list_h)
            # No pre-cap: _fit_size shrinks first and only drops items (with
            # an honest "+N more" line) if the ALLOCATED box genuinely can't
            # hold them all -- same as the matrix style, so a stage with 8
            # tasks isn't truncated to a smaller fixed number when it would
            # actually have fit.
            size, lines = _fit_size(list(column[key]), box_w, list_h, 6.0, min_pt=4.4)
            _render_cell_lines(slide, tf, lines, size)
            y = Emu(int(y + h))

        # OUTCOME -- single sentence, shrink-only.
        lb, lbtf = _textbox(slide, Emu(int(cx + pad)), y, box_w, Inches(0.14))
        lr = lbtf.paragraphs[0].add_run()
        lr.text = "OUTCOME"
        lr.font.size = Pt(6.2)
        lr.font.bold = True
        lr.font.name = _FONT
        lr.font.color.rgb = colour
        outcome_top = Emu(int(y + Inches(0.15)))
        outcome_box_h = Emu(int(outcome_h - Inches(0.15)))
        placeholder = _is_placeholder(column["outcome"])
        fitted = _fit_single(str(column["outcome"]), box_w, outcome_box_h, 6.0, min_pt=4.5)
        obox, otf = _textbox(slide, Emu(int(cx + pad)), outcome_top, box_w, outcome_box_h)
        op = otf.paragraphs[0]
        orun = op.add_run()
        orun.text = column["outcome"]
        orun.font.size = Pt(fitted)
        orun.font.italic = placeholder
        orun.font.name = _FONT
        orun.font.color.rgb = _RED if placeholder else _DARK_TEXT
        y = Emu(int(y + outcome_h))

        # DELIVERABLES -- chips.
        lb, lbtf = _textbox(slide, Emu(int(cx + pad)), y, box_w, Inches(0.14))
        lr = lbtf.paragraphs[0].add_run()
        lr.text = "DELIVERABLES"
        lr.font.size = Pt(6.2)
        lr.font.bold = True
        lr.font.name = _FONT
        lr.font.color.rgb = colour
        chips_top = Emu(int(y + Inches(0.16)))
        chips_h = Emu(int(card_top + card_h - chips_top - Inches(0.04)))
        deliverables = column["deliverables"]
        placeholder_deliv = any(_is_placeholder(d) for d in deliverables)
        _deliverable_chips(slide, Emu(int(cx + pad)), box_w, chips_top, chips_h, deliverables,
                           RGBColor(0xFC, 0xE8, 0xE8) if placeholder_deliv else colour,
                           _RED if placeholder_deliv else _WHITE)


def _slide_programme(slide, columns_data, from_stages, P, hold_icon, eng_icon,
                     project_name, client_name, wvr_confirmed) -> None:
    """Programme-matched columns -- stage-coloured column cards (WHAT WE DO
    / WITH YOU / YOU RECEIVE, deliverables as white chips), with orange
    HOLD POINT diamonds between columns wherever a stage actually carries
    one (see methodology_render.stage_carries_hold_point -- derived from
    the stage's own content, never asserted)."""
    n = max(len(columns_data), 1)
    _simple_title(slide, project_name)

    M = Inches(0.14)
    gap = Inches(0.42)  # room for hold-point diamonds between columns
    content_top = Emu(int(M + Inches(0.4)))
    col_w = Emu(int((_SLIDE_W - 2 * M - (n - 1) * gap) / n))

    header_h = Inches(0.55)
    card_h = Emu(int(_SLIDE_H - content_top - header_h - Inches(0.14)))
    do_h = Emu(int(card_h * 0.60))
    with_h = Emu(int(card_h * 0.16))
    recv_h = Emu(int(card_h - do_h - with_h - Inches(0.1)))

    for i, column in enumerate(columns_data):
        cx = Emu(int(M + i * (col_w + gap)))
        colour = _hex_to_rgb(stage_colour(i))

        _rect(slide, cx, content_top, col_w, header_h, colour)
        header_text = column["name"]
        _centered_text(slide, cx, content_top, col_w, Emu(int(header_h * 0.62)), header_text, 8.5,
                       RGBColor(0xFF, 0xE8, 0xE8) if _is_placeholder(header_text) else _WHITE)
        if column.get("chevron") and not _is_placeholder(column["chevron"]):
            _centered_text(slide, cx, Emu(int(content_top + header_h * 0.6)), col_w, Emu(int(header_h * 0.4)),
                           column["chevron"], 6.5, _WHITE, bold=False)

        card_top = Emu(int(content_top + header_h + Inches(0.05)))
        _round_rect(slide, cx, card_top, col_w, card_h, fill_color=None, line_color=RGBColor(0xE4, 0xE8, 0xEE), line_w=Pt(0.75))

        pad = Inches(0.06)
        y = Emu(int(card_top + Inches(0.05)))
        box_w = Emu(int(col_w - 2 * pad))
        for label_text, key, h in (("WHAT WE DO", "tasks", do_h), ("WITH YOU", "engagement", with_h)):
            lb, lbtf = _textbox(slide, Emu(int(cx + pad)), y, box_w, Inches(0.14))
            lr = lbtf.paragraphs[0].add_run()
            lr.text = label_text
            lr.font.size = Pt(6.0)
            lr.font.bold = True
            lr.font.name = _FONT
            lr.font.color.rgb = colour
            list_top = Emu(int(y + Inches(0.15)))
            list_h = Emu(int(h - Inches(0.15)))
            box, tf = _textbox(slide, Emu(int(cx + pad)), list_top, box_w, list_h)
            # No pre-cap -- see _slide_chevrons's identical comment.
            size, lines = _fit_size(list(column[key]), box_w, list_h, 5.8, min_pt=4.4)
            _render_cell_lines(slide, tf, lines, size)
            y = Emu(int(y + h))

        lb, lbtf = _textbox(slide, Emu(int(cx + pad)), y, box_w, Inches(0.14))
        lr = lbtf.paragraphs[0].add_run()
        lr.text = "YOU RECEIVE"
        lr.font.size = Pt(6.0)
        lr.font.bold = True
        lr.font.name = _FONT
        lr.font.color.rgb = colour
        chips_top = Emu(int(y + Inches(0.16)))
        chips_h = Emu(int(card_top + card_h - chips_top - Inches(0.04)))
        deliverables = column["deliverables"]
        placeholder_deliv = any(_is_placeholder(d) for d in deliverables)
        _deliverable_chips(slide, Emu(int(cx + pad)), box_w, chips_top, chips_h, deliverables,
                           RGBColor(0xFC, 0xE8, 0xE8) if placeholder_deliv else _tint(colour, 0.15),
                           _RED if placeholder_deliv else _WHITE)

        # Hold-point diamond after this stage, only if it actually carries
        # one -- positioned clear of the header band, in the gap column.
        if i < n - 1 and stage_carries_hold_point(column):
            gx = Emu(int(cx + col_w + gap / 2))
            gy = Emu(int(content_top + header_h + Inches(0.28)))
            _diamond(slide, gx, gy, Emu(int(Inches(0.11))), RGBColor(0xF9, 0x73, 0x16),
                    text="HOLD\nPOINT", text_size=5.2)


def _slide_spine(slide, columns_data, from_stages, P, hold_icon, eng_icon,
                 project_name, client_name, wvr_confirmed) -> None:
    """Timeline spine -- a coloured node per stage on a vertical spine at
    the left, a full-width band per stage (tinted left cell with
    name/weeks/outcome, then What we do / What you receive columns)."""
    n = max(len(columns_data), 1)
    _simple_title(slide, project_name)

    M = Inches(0.14)
    spine_x = Emu(int(M + Inches(0.08)))
    band_x = Emu(int(M + Inches(0.32)))
    band_right = Emu(int(_SLIDE_W - M))
    content_top = Emu(int(M + Inches(0.4)))
    available_h = Emu(int(_SLIDE_H - content_top - Inches(0.1)))
    band_gap = Inches(0.06)
    band_h = Emu(int((available_h - (n - 1) * band_gap) / max(n, 1)))

    label_w = Emu(int((band_right - band_x) * 0.22))
    tasks_x = Emu(int(band_x + label_w + Inches(0.08)))
    tasks_w = Emu(int((band_right - tasks_x) * 0.5))
    deliv_x = Emu(int(tasks_x + tasks_w + Inches(0.1)))
    deliv_w = Emu(int(band_right - deliv_x))

    # The spine line is drawn first so every stage's node (drawn later, in
    # the loop below) naturally lands on top of it -- no z-order surgery
    # needed the way drawing it last, over already-placed nodes, would.
    spine_bottom = Emu(int(content_top + n * band_h + max(0, n - 1) * band_gap))
    if n > 1:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, spine_x, content_top, spine_x, spine_bottom)
        line.line.color.rgb = RGBColor(0xE4, 0xE8, 0xEE)
        line.line.width = Pt(1.5)
        line.shadow.inherit = False

    y = content_top
    for i, column in enumerate(columns_data):
        colour = _hex_to_rgb(stage_colour(i))

        _rect(slide, band_x, y, label_w, band_h, _tint(colour, 0.85))
        _round_rect(slide, band_x, y, band_right - band_x, band_h, fill_color=None,
                   line_color=RGBColor(0xE4, 0xE8, 0xEE), line_w=Pt(0.75), radius=0.02)

        pad = Inches(0.05)
        placeholder_name = _is_placeholder(column["name"])
        lbox, ltf = _textbox(slide, Emu(int(band_x + pad)), Emu(int(y + Inches(0.03))), Emu(int(label_w - 2 * pad)), Emu(int(band_h - Inches(0.06))))
        lp = ltf.paragraphs[0]
        lr = lp.add_run()
        lr.text = column["name"]
        lr.font.size = Pt(8.5)
        lr.font.bold = True
        lr.font.name = _FONT
        lr.font.color.rgb = _RED if placeholder_name else _DARK_TEXT
        if column.get("chevron") and not _is_placeholder(column["chevron"]):
            wp = ltf.add_paragraph()
            wr = wp.add_run()
            wr.text = column["chevron"]
            wr.font.size = Pt(6.5)
            wr.font.name = _FONT
            wr.font.color.rgb = RGBColor(0x5B, 0x64, 0x72)
        placeholder_outcome = _is_placeholder(column["outcome"])
        op = ltf.add_paragraph()
        orun = op.add_run()
        orun.text = column["outcome"]
        orun.font.size = Pt(6.3)
        orun.font.italic = placeholder_outcome
        orun.font.name = _FONT
        orun.font.color.rgb = _RED if placeholder_outcome else RGBColor(0x5B, 0x64, 0x72)

        # WHAT WE DO
        tlb, tltf = _textbox(slide, tasks_x, Emu(int(y + Inches(0.03))), tasks_w, Inches(0.14))
        tlr = tltf.paragraphs[0].add_run()
        tlr.text = "WHAT WE DO"
        tlr.font.size = Pt(6.0)
        tlr.font.bold = True
        tlr.font.name = _FONT
        tlr.font.color.rgb = colour
        tbody_top = Emu(int(y + Inches(0.18)))
        tbody_h = Emu(int(band_h - Inches(0.21)))
        tbox, ttf = _textbox(slide, tasks_x, tbody_top, tasks_w, tbody_h)
        # No pre-cap -- see _slide_chevrons's identical comment.
        size, tlines = _fit_size(list(column["tasks"]), tasks_w, tbody_h, 6.0, min_pt=4.4)
        _render_cell_lines(slide, ttf, tlines, size)

        # WHAT YOU RECEIVE -- chips.
        dlb, dltf = _textbox(slide, deliv_x, Emu(int(y + Inches(0.03))), deliv_w, Inches(0.14))
        dlr = dltf.paragraphs[0].add_run()
        dlr.text = "WHAT YOU RECEIVE"
        dlr.font.size = Pt(6.0)
        dlr.font.bold = True
        dlr.font.name = _FONT
        dlr.font.color.rgb = colour
        deliverables = column["deliverables"]
        placeholder_deliv = any(_is_placeholder(d) for d in deliverables)
        chips_top = Emu(int(y + Inches(0.19)))
        chips_h = Emu(int(band_h - Inches(0.22)))
        _deliverable_chips(slide, deliv_x, deliv_w, chips_top, chips_h, deliverables,
                           RGBColor(0xFC, 0xE8, 0xE8) if placeholder_deliv else _tint(colour, 0.8),
                           _RED if placeholder_deliv else _DARK_TEXT, size_pt=6.0)

        node_r = Emu(int(Inches(0.045)))
        node_cy = Emu(int(y + band_h / 2))
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(spine_x - node_r)), Emu(int(node_cy - node_r)),
                                       Emu(int(2 * node_r)), Emu(int(2 * node_r)))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colour
        shape.line.color.rgb = _WHITE
        shape.line.width = Pt(1)
        shape.shadow.inherit = False

        y = Emu(int(y + band_h + band_gap))


_RENDERERS = {
    "matrix": _slide_matrix,
    "chevrons": _slide_chevrons,
    "programme": _slide_programme,
    "spine": _slide_spine,
}


def populate_methodology(
    analysis, client_name: str = "", project_name: str = "", theme_name: str | None = None,
    stages: list | None = None, week_labels: list | None = None,
    wvr_confirmed: bool = False, style: str | None = None,
) -> bytes:
    """
    Builds a fresh A4-landscape .pptx (returned as bytes) of the delivery
    methodology, in the requested presentation `style` -- one of
    methodology_render.STYLES ("matrix" | "chevrons" | "programme" |
    "spine"); anything else, including None, falls back to
    methodology_render.DEFAULT_STYLE, same convention org_chart_pptx.py and
    program_pptx.py use for their own style pickers.

    `stages`: the reviewed methodology_stages.MethodologyStage list from the
    Draft Responses tab. When supplied, every column -- name, key tasks,
    engagement activities, outcome, deliverables and the date chevron --
    comes from it, and cells the brief didn't support render as red TBC.
    When it is empty (the stage drafter hasn't been run), the table falls
    back to exactly what it produced before: one column of standard
    initiation boilerplate, one built from the brief's scope items, and two
    of placeholders. See methodology_render.build_columns().

    `week_labels`: the delivery program's own week labels, used for the date
    chevrons -- so they read "Wk 1 - Wk 3" normally and "6 Oct - 20 Oct"
    once a program start date is set, with no regeneration needed.

    `wvr_confirmed`: whether the user has confirmed their firm actually
    issues Work Verification Records. Only the "matrix" style prints this
    line (it always has); the three new styles have no equivalent fixed
    slot for it and do not assert it either way.

    `client_name` fills the "matrix" style's legend hold-point label and
    `project_name` every style's title -- both shown in red if not yet
    entered, same convention as every other missing field here.
    """
    resolved = style if style in STYLES else DEFAULT_STYLE

    P = _resolve_palette(theme_name)
    hold_icon = _render_icon_png("hold_point")
    eng_icon = _render_icon_png("engagement")

    columns_data = build_columns(analysis, stages, week_labels)
    from_stages = bool(stages)

    prs = Presentation()
    prs.slide_width = _SLIDE_W
    prs.slide_height = _SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _RENDERERS[resolved](slide, columns_data, from_stages, P, hold_icon, eng_icon,
                         project_name, client_name, wvr_confirmed)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()
