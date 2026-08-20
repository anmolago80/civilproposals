"""
test_batch1_wiring.py -- regression tests for the wiring fixes (Batches 1+).

Each test here pins one behaviour that used to be silently wrong: data the
app already collected but never got into the exported document. They are
deliberately plain assertions over the module functions rather than UI
tests, because the bugs were all in the wiring, not the widgets.

Run from this directory:

    python test_batch1_wiring.py
"""

from __future__ import annotations

import os
import sys


class _State(dict):
    """Stands in for st.session_state: .get() plus attribute-free access is
    all the modules under test use."""


def test_schedule_filler_reads_real_quals(failures: list[str]) -> None:
    """1(a): qualifications/years came from a session key nothing ever wrote,
    so a fully-filled Team & Resourcing tab still produced an empty
    Qualifications column in the client's personnel schedule."""
    from modules import resourcing, returnable_schedules

    state = _State(
        resource_plan=[
            resourcing.ResourceAssignment(
                slot="Project Director",
                slot_kind="management",
                person_name="Jane Smith",
                qualification="BEng (Civil) (Hons), UQ, 2003",
                rpeq_status="RPEQ 12345",
                years_experience="18 years",
            ),
            resourcing.ResourceAssignment(
                slot="Structural",
                slot_kind="discipline",
                person_name="Raj Patel",
                qualification="MEng (Structural)",
                # no RPEQ entered -- must not produce a trailing comma
                years_experience="9 years",
            ),
        ],
        team_members=[],
    )
    data = returnable_schedules.build_fill_data(state)
    people = {p["name"]: p for p in data["personnel"]}

    if "Jane Smith" not in people or "Raj Patel" not in people:
        failures.append("[1a] the resourcing plan's people didn't reach the fill data")
        return
    if people["Jane Smith"]["quals"] != "BEng (Civil) (Hons), UQ, 2003, RPEQ 12345":
        failures.append(f"[1a] qualification+RPEQ not joined: {people['Jane Smith']['quals']!r}")
    if people["Jane Smith"]["years"] != "18 years":
        failures.append(f"[1a] years_experience not carried: {people['Jane Smith']['years']!r}")
    if people["Raj Patel"]["quals"] != "MEng (Structural)":
        failures.append(f"[1a] blank RPEQ left punctuation behind: {people['Raj Patel']['quals']!r}")
    if people["Jane Smith"]["role"] != "Project Director":
        failures.append("[1a] role no longer resolves from slot/custom_title")


def test_blank_person_still_placeholders(failures: list[str]) -> None:
    """The no-invention rule: a person with nothing entered must come through
    with empty strings so the filler placeholders those cells, not with
    anything guessed."""
    from modules import resourcing, returnable_schedules

    state = _State(
        resource_plan=[
            resourcing.ResourceAssignment(
                slot="Geotechnical", slot_kind="discipline", person_name="Sam Lee",
            )
        ],
        team_members=[],
    )
    person = returnable_schedules.build_fill_data(state)["personnel"][0]
    if person["quals"] or person["years"]:
        failures.append(f"[1a] blank credentials were not left blank: {person!r}")


def test_sender_address_is_wired_and_saved(failures: list[str]) -> None:
    """1(b): the schedule filler read letter_sender_address, but no widget
    ever set it and it wasn't in the saved-project key list, so the address
    labels on client forms could never be filled by anyone."""
    from modules import project_store, returnable_schedules

    if "letter_sender_address" not in project_store.PLAIN_KEYS:
        failures.append("[1b] letter_sender_address is not saved with the project")

    state = _State(
        letter_sender_address="Level 3, 100 Example St, Brisbane QLD 4000",
        resource_plan=[], team_members=[],
    )
    data = returnable_schedules.build_fill_data(state)
    if data.get("contact_address") != "Level 3, 100 Example St, Brisbane QLD 4000":
        failures.append(f"[1b] contact_address not built: {data.get('contact_address')!r}")

    # Round-trip it through a real save/load so a registration that exists but
    # doesn't survive the zip still fails here.
    saved = project_store.save_project(_State(letter_sender_address="12 Test Rd"))
    if project_store.load_project(saved).get("letter_sender_address") != "12 Test Rd":
        failures.append("[1b] letter_sender_address did not survive save/load")

    # And it must actually match an address label on a form.
    field_key, value = returnable_schedules.match_label("Registered Office:", data)
    if field_key != "contact_address" or not value:
        failures.append(f"[1b] a 'Registered Office' label didn't resolve: {field_key!r} {value!r}")


def test_program_start_date(failures: list[str]) -> None:
    """1(f): week headers can carry real calendar dates, and the date has to
    survive save/load (a date is not JSON-safe, so it needs its own handling
    rather than a PLAIN_KEYS entry)."""
    import datetime

    from modules import program_schedule, project_store

    plain = program_schedule.week_labels(3, None)
    if plain != ["Wk 1", "Wk 2", "Wk 3"]:
        failures.append(f"[1f] unchanged behaviour without a start date broke: {plain}")

    dated = program_schedule.week_labels(3, datetime.date(2026, 10, 6))
    if dated != ["Wk 1 - 6 Oct", "Wk 2 - 13 Oct", "Wk 3 - 20 Oct"]:
        failures.append(f"[1f] dated week labels are wrong: {dated}")

    if "program_start_date" not in project_store.DATE_KEYS:
        failures.append("[1f] program_start_date is not saved with the project")
    saved = project_store.save_project(_State(program_start_date=datetime.date(2026, 10, 6)))
    if project_store.load_project(saved).get("program_start_date") != datetime.date(2026, 10, 6):
        failures.append("[1f] the start date did not survive save/load")
    # Unset must round-trip as None, not as today's date or a crash.
    if project_store.load_project(project_store.save_project(_State())).get("program_start_date") is not None:
        failures.append("[1f] an unset start date did not round-trip as None")


def test_graphics_recommendations_are_current(failures: list[str]) -> None:
    """1(g): the graphics list is printed into the exported pack, so a stale
    entry tells the user to hand-build something the app already made."""
    from modules import graphics_engine, proposal_structure, tender_analyser

    sections = proposal_structure.build_proposal_structure(
        tender_analyser.TenderAnalysis(), [], [], "formal",
    )
    recs = {r.graphic_title: r for r in graphics_engine.recommend_graphics(
        sections, project_type="Coastal & Ocean Engineering")}

    for title in ("Organisation chart", "Methodology process diagram", "Programme timeline"):
        rec = recs.get(title)
        if rec is None:
            continue  # this brief's skeleton doesn't recommend it -- fine
        if rec.status != "Generated":
            failures.append(f"[1g] '{title}' is still reported as needing user input")
        if rec.placeholder_text:
            failures.append(f"[1g] '{title}' still carries a red placeholder instruction")

    # Genuinely ungenerated graphics must still be placeholdered.
    risk = recs.get("Key risk diagram")
    if risk is not None and risk.status == "Generated":
        failures.append("[1g] an ungenerated graphic was wrongly marked Generated")

    # The divider hint follows the project type instead of always saying bridge/road.
    divider = next((r for r in recs.values() if "divider" in r.graphic_title.lower()), None)
    if divider is None or "COASTAL" not in divider.placeholder_text:
        failures.append(f"[1g] the divider hint ignores project_type: "
                        f"{getattr(divider, 'placeholder_text', None)!r}")
    if divider is not None and "BRIDGE / ROAD" in divider.placeholder_text:
        failures.append("[1g] the divider hint is still hardcoded to bridge/road")


def _sheet_rows(blob: bytes) -> list[list]:
    import io as _io

    import openpyxl
    ws = openpyxl.load_workbook(_io.BytesIO(blob)).active
    return [list(row) for row in ws.iter_rows(values_only=True)]


def test_excel_fee_exports(failures: list[str]) -> None:
    """1(h): an unpriced discipline exported as a literal $0, which reads as
    an offer to do that work for nothing; the workbook carried nothing
    identifying the project; and the fee-% total cell was left empty."""
    from modules import fee_estimation_engine, resourcing

    info = {"project_name": "Example Creek Bridge", "client_name": "Example Shire Council",
            "tender_name": "RFT 2026-014", "bidder_name": "Test Engineering Pty Ltd"}

    lines = [
        resourcing.DisciplineFeeLine(discipline="Structural", total_hours=800, rate_per_hour=210),
        resourcing.DisciplineFeeLine(discipline="Hydraulics", total_hours=0, rate_per_hour=0),
    ]
    rows = _sheet_rows(resourcing.discipline_fee_lines_to_excel(lines, "Government", info))
    flat = [str(cell) for row in rows for cell in row if cell is not None]

    unpriced = next((r for r in rows if r and r[0] == "Hydraulics"), None)
    if unpriced is None:
        failures.append("[1h] the unpriced discipline row is missing entirely")
    else:
        if any(cell == 0 for cell in unpriced[1:4]):
            failures.append(f"[1h] an unpriced discipline still exports as 0: {unpriced}")
        if unpriced[4] != resourcing.UNPRICED_NOTE:
            failures.append(f"[1h] the unpriced row carries no explanatory note: {unpriced}")
    if not any("NOT been priced" in cell for cell in flat):
        failures.append("[1h] the workbook doesn't explain what a blank fee cell means")
    for expected in ("Example Creek Bridge", "Example Shire Council", "RFT 2026-014"):
        if expected not in flat:
            failures.append(f"[1h] the fee build-up workbook doesn't identify the project ({expected})")
    priced = next((r for r in rows if r and r[0] == "Structural"), None)
    if priced is None or priced[3] != 168000:
        failures.append(f"[1h] a priced discipline no longer exports its total: {priced}")

    estimates = [
        fee_estimation_engine.DisciplineFeeEstimate(
            discipline="Structural", fee_percentage=58.3, confidence="Medium", source="Benchmark"),
        fee_estimation_engine.DisciplineFeeEstimate(
            discipline="Geotechnical", fee_percentage=41.7, confidence="Medium", source="Benchmark"),
    ]
    rows = _sheet_rows(fee_estimation_engine.fee_estimates_to_excel(estimates, project_info=info))
    total_row = next((r for r in rows if r and r[0] == "Total"), None)
    if total_row is None:
        failures.append("[1h] the fee-split workbook has no Total row")
    elif total_row[1] is None:
        failures.append("[1h] the fee-% total cell is still empty")
    elif abs(total_row[1] - 100.0) > 0.05:
        failures.append(f"[1h] the fee-% total is wrong: {total_row[1]}")
    if "Example Creek Bridge" not in [str(c) for r in rows for c in r if c is not None]:
        failures.append("[1h] the fee-split workbook doesn't identify the project")


def test_fee_percentages_add_up(failures: list[str]) -> None:
    """1(i): independently rounding each share printed fee tables that
    totalled 99% or 101%."""
    from modules.export_docx import format_fee_percentages

    for values in ([33.4, 33.3, 33.3], [58.3, 25.4, 16.3], [16.7] * 6):
        printed = format_fee_percentages(values)
        total = sum(int(p.rstrip("%")) for p in printed)
        if total != 100:
            failures.append(f"[1i] {values} printed as {printed}, totalling {total}%")
        for value, shown in zip(values, printed):
            if abs(int(shown.rstrip("%")) - value) > 1:
                failures.append(f"[1i] {value} was printed as {shown} -- more than a point out")

    # A split that genuinely isn't 100% must keep showing that.
    if format_fee_percentages([40.0, 30.0]) == ["40%", "60%"]:
        failures.append("[1i] a 70% split was silently normalised to 100%")


def test_firm_profile(failures: list[str]) -> None:
    """Batch 3: firm-level facts fill placeholders that nobody could fill
    before -- and an EMPTY profile must leave every one of them exactly as
    it is today."""
    import json

    from modules import firm_profile, returnable_schedules

    class _Profile:
        company_name = "Example Engineering Pty Ltd"
        abn = "12 345 678 901"
        acn = ""
        registered_address = "Level 3, 100 Example St\nBrisbane QLD 4000"
        logo_bytes = b"not-a-real-png"
        signatory_name = "Jane Smith"
        signatory_title = "Project Director"
        signatory_phone = "07 3000 0000"
        signatory_email = "jane@example.com"
        insurances_json = json.dumps([
            {"type": "Professional Indemnity", "insurer": "Example Insurance Ltd",
             "policy_no": "PI-123456", "cover": "$10,000,000", "expiry": "30 June 2027"},
            {"type": "Public Liability", "insurer": "", "policy_no": "", "cover": "", "expiry": ""},
        ])
        certifications_json = json.dumps(["ISO 9001:2015"])
        rate_card_json = json.dumps({"Structural": 210.0})
        offices_text = "Offices in Brisbane and Townsville since 2004."
        community_text = "We fund two regional engineering scholarships each year."
        leadership_text = "Jane Smith (Managing Director) oversees delivery."
        terms_of_engagement_text = "AS 4122-2010 General Conditions."
        qa_statement = "All deliverables issued with WVRs."

    class _Empty:
        company_name = abn = acn = registered_address = ""
        logo_bytes = None
        signatory_name = signatory_title = signatory_phone = signatory_email = ""
        insurances_json = certifications_json = rate_card_json = ""
        offices_text = community_text = leadership_text = ""
        terms_of_engagement_text = qa_statement = ""

    filled, empty = _Profile(), _Empty()

    if firm_profile.is_empty(filled) or not firm_profile.is_empty(empty):
        failures.append("[3] is_empty() doesn't distinguish a filled profile from a blank one")

    # A half-filled insurance row must not export as a row of blanks.
    rows = firm_profile.insurances(filled)
    if len(rows) != 1 or rows[0]["insurer"] != "Example Insurance Ltd":
        failures.append(f"[3d] blank insurance rows are not being dropped: {rows}")

    # Footer: real line when complete, the original red placeholder when not.
    text, complete = firm_profile.footer_line(filled)
    if not complete or "ABN 12 345 678 901" not in text:
        failures.append(f"[3d] the footer line isn't built from the profile: {text!r}")
    text, complete = firm_profile.footer_line(empty, bidder_name="Test Pty Ltd")
    if complete or "[REGISTERED ADDRESS]" not in text:
        failures.append(f"[3d] an empty profile changed the footer placeholder: {text!r}")

    # Schedule filler: knowable now, still never-known when blank.
    data = returnable_schedules.build_fill_data(
        _State(resource_plan=[], team_members=[]), firm_profile.schedule_fill_data(filled))
    for label, expected in (("ABN:", "12 345 678 901"),
                            ("Professional Indemnity Insurer:", "Example Insurance Ltd"),
                            ("Certifications:", "ISO 9001:2015")):
        key, value = returnable_schedules.match_label(label, data)
        if value != expected:
            failures.append(f"[3d] '{label}' resolved to {value!r}, expected {expected!r}")
    # A signature can never come from a saved setting.
    if returnable_schedules.match_label("Signature:", data)[1] is not None:
        failures.append("[3d] a signature was filled from the firm profile")
    # An insurance the profile does NOT hold stays placeholdered.
    if returnable_schedules.match_label("Public Liability Insurer:", data)[1] is not None:
        failures.append("[3d] an unheld insurance was answered anyway")

    empty_data = returnable_schedules.build_fill_data(
        _State(resource_plan=[], team_members=[]), firm_profile.schedule_fill_data(empty))
    for label in ("ABN:", "Professional Indemnity Insurer:"):
        key, value = returnable_schedules.match_label(label, empty_data)
        if key != "never_known" or value is not None:
            failures.append(f"[3e] with a blank profile '{label}' no longer placeholders: {key!r}")

    # Seeding fills blanks only.
    seed = firm_profile.project_seed(filled)
    if seed.get("bidder_name") != "Example Engineering Pty Ltd":
        failures.append("[3c] the project seed doesn't carry the company name")
    if firm_profile.project_seed(empty):
        failures.append("[3c] an empty profile still produced a seed")


def test_failure_honesty(failures: list[str]) -> None:
    """Batch 6: a step that didn't happen must not report as one that found
    nothing."""
    from modules import ai_interface, reference_projects, team_bios

    # A repaired (truncated) parse is flagged rather than looking clean.
    if not ai_interface.was_repaired({ai_interface.REPAIRED_FLAG: True}):
        failures.append("[6d] a repaired response isn't detectable by callers")
    if ai_interface.was_repaired({"ok": 1}):
        failures.append("[6d] an ordinary response is being reported as repaired")

    # Truncation is reported, not silent.
    long_text = "x " * 40000
    members, warns = team_bios.draft_team_bios_from_cv("", max_chars=100)
    if members:
        failures.append("[6c] empty CV text produced members")
    # The truncation branch must produce a warning before any AI call is
    # attempted, so check the message construction directly.
    import inspect
    source = inspect.getsource(team_bios.draft_team_bios_from_cv)
    if "characters of the CV material were read" not in source:
        failures.append("[6c] CV truncation is still silent")
    source = inspect.getsource(reference_projects.draft_reference_projects)
    if "characters of the reference material were read" not in source:
        failures.append("[6c] reference-material truncation is still silent")


def test_fee_presentation_ticks(failures: list[str]) -> None:
    """Batch 7A: which fee presentations reach the proposal is the user's
    choice, the defaults reproduce today's packs, and nothing ticked gives a
    visible placeholder rather than a silently fee-less proposal."""
    from modules import export_docx, project_store

    if "fee_sections_included" not in project_store.PLAIN_KEYS:
        failures.append("[7A] the fee ticks aren't saved with the project")
    saved = project_store.save_project(_State(fee_sections_included={
        "pct_split": False, "discipline_buildup": True, "scope_buildup": True}))
    loaded = project_store.load_project(saved).get("fee_sections_included")
    if loaded != {"pct_split": False, "discipline_buildup": True, "scope_buildup": True}:
        failures.append(f"[7A] the ticks didn't survive save/load: {loaded}")

    # Defaults must reproduce what both packs exported before the choice
    # existed: % split + discipline build-up in, scope-item build-up out.
    defaults = export_docx.fee_sections(None)
    if defaults != {"pct_split": True, "discipline_buildup": True, "scope_buildup": False}:
        failures.append(f"[7A] the defaults changed today's packs: {defaults}")

    import test_exports
    project = test_exports.build_sample_project()

    def letter_text(included):
        return test_exports.docx_text(export_docx.build_letter_docx(
            project_info=project["project_info"], sender=project["sender"],
            analysis=project["analysis"], understanding_text="", methodology_text="",
            resource_plan=[], personnel_photos={},
            program_schedule=project["program_schedule"],
            program_week_labels=project["program_week_labels"],
            terms_of_engagement_text="", fee_estimates=project["fee_estimates"],
            discipline_fee_lines=project["discipline_fee_lines"],
            fee_sections_included=included,
        ).getvalue())

    combos = {
        "default": (None, True, True),
        "pct only": ({"pct_split": True, "discipline_buildup": False, "scope_buildup": False}, True, False),
        "buildup only": ({"pct_split": False, "discipline_buildup": True, "scope_buildup": False}, False, True),
        "none": ({"pct_split": False, "discipline_buildup": False, "scope_buildup": False}, False, False),
    }
    for name, (included, want_split, want_buildup) in combos.items():
        text = letter_text(included)
        has_split = "Indicative fee split" in text
        has_buildup = "Discipline fee build-up" in text
        if has_split != want_split or has_buildup != want_buildup:
            failures.append(
                f"[7A] '{name}' exported split={has_split} buildup={has_buildup}, "
                f"expected {want_split}/{want_buildup}")
        # Nothing ticked must say so rather than silently omitting fees.
        wants_placeholder = not (want_split or want_buildup)
        if (export_docx.FEE_NOTHING_SELECTED in text) != wants_placeholder:
            failures.append(f"[7A] '{name}' placeholder handling is wrong")


def test_program_styles(failures: list[str]) -> None:
    """Batch 7B: the delivery program is drawn in the style the user picked,
    the same style in every output, and every style survives the awkward
    cases -- no stages, no program at all, and a program far bigger than a
    slide."""
    import datetime
    import io as _io

    from docx import Document
    from pptx import Presentation

    from modules import export_docx, program_pptx, program_render, program_schedule, project_store

    if "program_style" not in project_store.PLAIN_KEYS:
        failures.append("[7B] the chosen program style isn't saved with the project")
    loaded = project_store.load_project(
        project_store.save_project(_State(program_style="timeline"))).get("program_style")
    if loaded != "timeline":
        failures.append(f"[7B] the style didn't survive save/load: {loaded}")

    weeks = 14
    start = datetime.date(2026, 9, 7)
    span = lambda a, b: [a <= i + 1 <= b for i in range(weeks)]  # noqa: E731
    schedule = {
        "Project initiation & site establishment": span(1, 3),
        "Topographic & utility survey": span(2, 5),
        "Concept design development": span(5, 8),
        "Detailed design (civil, drainage, pavement)": span(8, 12),
        "Issue for construction documentation": span(13, 14),
    }
    labels = program_schedule.week_labels(weeks, start)
    stages = [
        methodology_stage("Stage 1 - Investigation", 1, 6),
        methodology_stage("Stage 2 - Concept design", 5, 10),
        methodology_stage("Stage 3 - Detailed design", 8, 12),
        methodology_stage("Stage 4 - Documentation", 11, 14),
    ]

    model = program_render.build_model(schedule, labels, stages, start, None, "P", "C")

    # Stage assignment must take the containing stage that begins LATEST --
    # first-match filed everything under whichever stage merely began first.
    by_label = {item.label: model.stages[item.stage_index] for item in model.items}
    if by_label.get("Detailed design (civil, drainage, pavement)") != "Stage 3 - Detailed design":
        failures.append(f"[7B] stage assignment picked the wrong lane: {by_label}")

    for style in program_render.STYLES:
        if not program_render.render_png(model, style, "#1D4ED8"):
            failures.append(f"[7B] the {style} style didn't render")

    # Swimlanes has nothing to group by without stages, and must SAY so by
    # resolving to a style the caller can name -- not quietly draw one lane.
    no_stages = program_render.build_model(schedule, labels, [], start, None, "P", "C")
    if program_render.effective_style(no_stages, "swimlanes") != "gantt":
        failures.append("[7B] swimlanes without stages doesn't fall back to the Gantt")
    if program_render.effective_style(model, "swimlanes") != "swimlanes":
        failures.append("[7B] swimlanes falls back even when stages exist")

    # An empty program keeps its red placeholder in every style, rather than
    # exporting a plausible-looking empty grid.
    empty = program_render.build_model({}, labels, stages, start, None, "P", "C")
    for style in program_render.STYLES:
        if not program_render.render_png(empty, style, "#1D4ED8"):
            failures.append(f"[7B] the {style} style lost its empty-program placeholder")
    empty_deck = Presentation(_io.BytesIO(program_pptx.populate_program({}, [], style="table")))
    deck_text = " ".join(
        shape.text_frame.text for shape in empty_deck.slides[0].shapes
        if shape.has_text_frame)
    if "NO PROGRAM ENTERED" not in deck_text:
        failures.append("[7B] an empty program exports a deck with no placeholder")

    # 1(i)'s overflow fix has to hold in all four styles, not just the one it
    # was written for: PowerPoint stores shapes outside the slide and simply
    # doesn't show them.
    big = {f"Scope item {i}": [True] * weeks for i in range(24)}
    for style in program_render.STYLES:
        blob = program_pptx.populate_program(
            big, labels, project_name="P", style=style, methodology_stages=stages,
            start_date=start)
        prs = Presentation(_io.BytesIO(blob))
        lowest = max(shape.top + shape.height for shape in prs.slides[0].shapes)
        if lowest > prs.slide_height:
            failures.append(
                f"[7B] 24 scope items overflow the {style} slide by "
                f"{(lowest - prs.slide_height) / 914400:.2f} in")

    # The formal table must arrive as a real, editable PowerPoint table --
    # the whole reason that style exists.
    table_deck = Presentation(_io.BytesIO(program_pptx.populate_program(
        schedule, labels, project_name="P", style="table", start_date=start)))
    if not any(shape.has_table for shape in table_deck.slides[0].shapes):
        failures.append("[7B] the formal-table deck isn't a real PowerPoint table")

    # ... and as a real Word table in the letter pack, while the other three
    # arrive as a picture of exactly what the preview drew.
    import test_exports
    project = test_exports.build_sample_project()

    def letter(style):
        return export_docx.build_letter_docx(
            project_info=project["project_info"], sender=project["sender"],
            analysis=project["analysis"], understanding_text="", methodology_text="",
            resource_plan=[], personnel_photos={},
            program_schedule=schedule, program_week_labels=labels,
            terms_of_engagement_text="", fee_estimates=project["fee_estimates"],
            discipline_fee_lines=project["discipline_fee_lines"],
            program_style=style, methodology_stages=stages, program_start_date=start,
        ).getvalue()

    table_doc = Document(_io.BytesIO(letter("table")))
    if not any("Commence" in cell.text for table in table_doc.tables for row in table.rows
               for cell in row.cells):
        failures.append("[7B] the letter pack's formal table isn't a native Word table")
    for style in ("gantt", "swimlanes", "timeline"):
        doc = Document(_io.BytesIO(letter(style)))
        if not doc.inline_shapes:
            failures.append(f"[7B] the letter pack's {style} program has no image")

    # A program that can't be drawn must still export a program the reader
    # can follow, not an empty section.
    broken = _State()
    original = program_render.render_png
    try:
        program_render.render_png = lambda *a, **k: None
        fallback = Document(_io.BytesIO(letter("timeline")))
    finally:
        program_render.render_png = original
    grid_headers = [cell.text for table in fallback.tables for cell in table.rows[0].cells]
    if not any(str(labels[0]) == header for header in grid_headers):
        failures.append("[7B] a failed program image doesn't fall back to the week grid")
    del broken


def methodology_stage(name: str, week_start: int, week_end: int):
    from modules.methodology_stages import MethodologyStage

    return MethodologyStage(name=name, week_start=week_start, week_end=week_end)


def test_optional_design_manager(failures: list[str]) -> None:
    """Batch 11: Design Manager is optional. A commission without one must
    render with the role simply absent -- never as a red TBC, because an
    intentional removal is not a gap."""
    import io as _io

    from pptx import Presentation

    from modules import org_chart, org_chart_pptx, project_store, resourcing

    if "removed_management_roles" not in project_store.PLAIN_KEYS:
        failures.append("[11] the removal isn't saved with the project")
    loaded = project_store.load_project(project_store.save_project(
        _State(removed_management_roles=["Design Manager"]))).get("removed_management_roles")
    if loaded != ["Design Manager"]:
        failures.append(f"[11] the removal didn't survive save/load: {loaded}")

    # Project Director and Project Manager stay mandatory today.
    if resourcing.OPTIONAL_MANAGEMENT_ROLES != {"Design Manager"}:
        failures.append(
            f"[11] unexpected optional roles: {resourcing.OPTIONAL_MANAGEMENT_ROLES}")
    for role in ("Project Director", "Project Manager", resourcing.CLIENT_ROLE):
        if resourcing.is_removable_management_role(role):
            failures.append(f"[11] {role} became removable")
    if not resourcing.is_removable_management_role("Design Manager"):
        failures.append("[11] Design Manager isn't removable")

    # A rebuild -- which is what re-running Tender Analysis does -- must not
    # quietly put a removed role back.
    rebuilt = resourcing.build_resource_plan(["Structural"], ["Design Manager"])
    if any(a.slot == "Design Manager" for a in rebuilt):
        failures.append("[11] rebuilding the plan resurrected the removed Design Manager")
    for role in ("Project Director", "Project Manager", resourcing.CLIENT_ROLE):
        if not any(a.slot == role for a in rebuilt):
            failures.append(f"[11] rebuilding the plan dropped the mandatory {role}")

    # Naming a mandatory role in the removal list must be ignored, not obeyed.
    if "Project Director" not in resourcing.management_roles_for_plan(["Project Director"]):
        failures.append("[11] a hand-edited project file can delete the Project Director")

    plan = resourcing.build_resource_plan(["Structural", "Geotechnical"], ["Design Manager"])
    for a in plan:
        if a.slot == "Project Director":
            a.person_name = "Jane Citizen"
        if a.slot == "Structural":
            a.person_name = "Tom Sample"

    deck = Presentation(_io.BytesIO(org_chart_pptx.populate_org_chart(
        plan, client_name="Example Council", project_name="P", tender_name="RFT-1")))
    slide_text = " ".join(
        shape.text_frame.text for shape in deck.slides[0].shapes if shape.has_text_frame)
    if "Design Manager" in slide_text:
        failures.append("[11] the removed Design Manager is still drawn on the org chart deck")
    for role in ("Project Director", "Project Manager"):
        if role not in slide_text:
            failures.append(f"[11] the org chart deck lost its {role}")

    # The in-app / DOCX-embedded PNG must render too -- it takes a different
    # code path from the deck.
    if not org_chart.render_org_chart(plan, project_title="P"):
        failures.append("[11] the org chart PNG didn't render without a Design Manager")

    # Nothing downstream may emit a phantom line for the missing role.
    from modules import draft_generator

    if "Design Manager" in draft_generator.format_team_context(plan):
        failures.append("[11] the AI team context names a Design Manager that isn't staffed")
    if any(e["role_label"] == "Design Manager" for e in resourcing.letter_team_entries(plan)):
        failures.append("[11] the letter pack's team list still has a Design Manager")
    if any(r == "Design Manager"
           for entry in resourcing.personnel_profiles_deduped(plan) for r in entry["roles"]):
        failures.append("[11] a Key Personnel profile is still reserved for the removed role")


def test_org_chart_styles(failures: list[str]) -> None:
    """Batch 9: the org chart is drawn in the style the user picked, the same
    style in the preview and the deck, and every style survives the awkward
    cases -- an unfilled lead, no reviewer, and far more disciplines than fit
    across one row."""
    import io as _io

    from pptx import Presentation

    from modules import org_chart_pptx, org_chart_render, project_store, resourcing

    if "org_chart_style" not in project_store.PLAIN_KEYS:
        failures.append("[9] the chosen org chart style isn't saved with the project")
    loaded = project_store.load_project(project_store.save_project(
        _State(org_chart_style="bands"))).get("org_chart_style")
    if loaded != "bands":
        failures.append(f"[9] the style didn't survive save/load: {loaded}")

    def slot(name, kind, person="", lead=True, title="", qual=""):
        return resourcing.ResourceAssignment(
            slot=name, slot_kind=kind, person_name=person, is_lead=lead,
            custom_title=title, qualification=qual)

    plan = [
        slot("Client Project Manager", "management", "Dana Client"),
        slot("Project Director", "management", "Jane Citizen", qual="BE(Civil)"),
        slot("Project Manager", "management", "Alex Demo"),
        slot("Independent Review", "discipline", "Chris Invented"),
        slot("Structural", "discipline", "Tom Sample"),
        slot("Structural", "discipline", "Ryan Example", lead=False, title="Bridge Engineer"),
        slot("Hydraulics", "discipline"),                      # unfilled lead -> TBC
        slot("Survey", "discipline", "Pat Mockup"),
        slot("Survey", "discipline", "Lee Fictional", lead=False),   # no title -> placeholder
    ]
    model = org_chart_render.build_model(plan, "Coastal Council", "Bridge Duplication", "RFT-1")

    if [p.role for p in model.leadership] != ["Project Director", "Project Manager"]:
        failures.append(f"[9] leadership came out wrong: {[p.role for p in model.leadership]}")
    if [p.name for p in model.assurance] != ["Chris Invented"]:
        failures.append(f"[9] the reviewer wasn't picked up: {model.assurance}")
    # The client's own PM names the client box; it is not one of OUR leaders.
    if any(p.role == resourcing.CLIENT_ROLE for p in model.leadership):
        failures.append("[9] the client's PM was filed as firm leadership")
    hydraulics = [g for g in model.disciplines if g.name == "Hydraulics"][0]
    if not hydraulics.lead.is_tbc:
        failures.append("[9] an unfilled lead isn't marked TBC")
    untitled = [g for g in model.disciplines if g.name == "Survey"][0].supports[0]
    if untitled.role != org_chart_render.CONFIRM_TITLE or not untitled.role_is_placeholder:
        failures.append(f"[9] an untitled support member isn't flagged: {untitled.role!r}")

    for style in org_chart_render.STYLES:
        if not org_chart_render.render_png(model, style, "#1D4ED8"):
            failures.append(f"[9] the {style} style didn't render")

    # A project with no reviewer must render with no assurance element at all
    # -- an empty ASSURANCE band reads as a missing answer, not an absent role.
    lean_plan = [a for a in plan if a.slot != "Independent Review"]
    if org_chart_render.build_model(lean_plan, "C", "P", "T").has_assurance:
        failures.append("[9] a reviewer was invented for a project that has none")
    bands = Presentation(_io.BytesIO(org_chart_pptx.populate_org_chart(
        lean_plan, style="bands")))
    band_text = " ".join(sh.text_frame.text for sh in bands.slides[0].shapes if sh.has_text_frame)
    if "ASSURANCE" in band_text:
        failures.append("[9] the bands deck shows an assurance band with nobody in it")

    # An empty plan keeps its placeholder rather than exporting a bare chart.
    empty_model = org_chart_render.build_model([], "", "", "")
    for style in org_chart_render.STYLES:
        if not org_chart_render.render_png(empty_model, style, "#1D4ED8"):
            failures.append(f"[9] the {style} style lost its empty-plan placeholder")
    empty_deck = Presentation(_io.BytesIO(org_chart_pptx.populate_org_chart([], style="cards")))
    if "NO TEAM ASSIGNED" not in " ".join(
            sh.text_frame.text for sh in empty_deck.slides[0].shapes if sh.has_text_frame):
        failures.append("[9] an empty plan exports a deck with no placeholder")

    # Eight disciplines must not run off the bottom of the slide in any style:
    # PowerPoint stores shapes outside the slide and simply doesn't show them.
    big = [slot("Client Project Manager", "management", "Dana Client"),
           slot("Project Director", "management", "Jane Citizen"),
           slot("Project Manager", "management", "Alex Demo")]
    for index in range(8):
        big.append(slot(f"Discipline {index + 1}", "discipline", f"Person {index + 1}"))
        big.append(slot(f"Discipline {index + 1}", "discipline", f"Support {index + 1}",
                        lead=False, title="Design Engineer"))
    big_model = org_chart_render.build_model(big, "C", "P", "T")
    for style in org_chart_render.STYLES:
        deck = Presentation(_io.BytesIO(org_chart_pptx.populate_org_chart(
            big, client_name="C", project_name="P", style=style)))
        lowest = max(shape.top + shape.height for shape in deck.slides[0].shapes)
        if lowest > deck.slide_height:
            failures.append(
                f"[9] 8 disciplines overflow the {style} slide by "
                f"{(lowest - deck.slide_height) / 914400:.2f} in")
        if not org_chart_render.render_png(big_model, style, "#1D4ED8"):
            failures.append(f"[9] the {style} preview failed with 8 disciplines")


def test_fee_prepopulation(failures: list[str]) -> None:
    """Batch 10: three labelled tiers of fee prepopulation -- the firm's own
    history, the bundled rule-of-thumb, and AI-modelled benchmarks -- with
    history ranked first, snapshots isolated per user, and a failed AI call
    that refuses to impersonate a successful one."""
    from modules import db, fee_estimation_engine, fee_history, resourcing

    db.init_db()

    def line(discipline, hours, rate):
        return resourcing.DisciplineFeeLine(discipline=discipline, total_hours=hours,
                                            rate_per_hour=rate)

    user_a, user_b = "test-fee-user-a", "test-fee-user-b"
    try:
        with db.get_session() as session:
            session.query(db.FeeSnapshot).filter(
                db.FeeSnapshot.user_id.in_([user_a, user_b])).delete(synchronize_session=False)
            session.commit()
    except Exception as exc:  # noqa: BLE001
        failures.append(f"[10] couldn't reach the fee_snapshots table: {exc}")
        return

    # One bid is not a median -- nothing is offered until there are two.
    fee_history.record_snapshot(user_a, "proj-1", "Bridge",
                                [line("Structural", 100, 200), line("Geotechnical", 50, 200)])
    if fee_history.fee_history_benchmarks(user_a, "Bridge")["disciplines"]:
        failures.append("[10] a single past bid is presented as a benchmark")

    fee_history.record_snapshot(user_a, "proj-2", "Bridge",
                                [line("Structural", 60, 200), line("Geotechnical", 140, 200)])
    history = fee_history.fee_history_benchmarks(user_a, "Bridge")
    if history["bids"] != 2:
        failures.append(f"[10] wrong bid count: {history['bids']}")
    medians = {e["discipline"]: e["median_pct"] for e in history["disciplines"]}
    # Structural was 66.7% then 30%; Geotechnical 33.3% then 70%.
    if abs(medians.get("Structural", 0) - 48.4) > 0.6:
        failures.append(f"[10] the Structural median is wrong: {medians}")
    ranges = {e["discipline"]: (e["low_pct"], e["high_pct"]) for e in history["disciplines"]}
    if abs(ranges["Structural"][0] - 30.0) > 0.2 or abs(ranges["Structural"][1] - 66.7) > 0.2:
        failures.append(f"[10] the Structural range is wrong: {ranges}")

    # Re-exporting the same project must refresh its snapshot, not add a bid.
    fee_history.record_snapshot(user_a, "proj-2", "Bridge",
                                [line("Structural", 60, 200), line("Geotechnical", 140, 200)])
    if fee_history.fee_history_benchmarks(user_a, "Bridge")["bids"] != 2:
        failures.append("[10] regenerating a pack counted as another bid")

    # Per-user isolation: one firm's pricing must never reach another's.
    if fee_history.fee_history_benchmarks(user_b, "Bridge")["disciplines"]:
        failures.append("[10] one user's fee history leaked into another user's benchmark")
    if fee_history.fee_history_benchmarks(None, "Bridge")["disciplines"]:
        failures.append("[10] fee history is served without a user to scope it to")

    # An unpriced row is dropped, not averaged in as a 0% discipline.
    fee_history.record_snapshot(user_b, "proj-3", "Road",
                                [line("Structural", 10, 100), line("Survey", 0, 0)])
    fee_history.record_snapshot(user_b, "proj-4", "Road",
                                [line("Structural", 10, 100), line("Survey", 0, 0)])
    road = fee_history.fee_history_benchmarks(user_b, "Road")
    if any(e["discipline"] == "Survey" for e in road["disciplines"]):
        failures.append("[10] an unpriced discipline was stored as a 0% benchmark")

    # A project with nothing priced isn't a bid at all.
    if fee_history.record_snapshot(user_b, "proj-5", "Road", [line("Structural", 0, 0)]):
        failures.append("[10] an unpriced project was recorded as a bid")

    # Tier ranking: history wins where it covers the disciplines, bundled
    # otherwise, and both say which they are.
    split, source = fee_history.best_available_split(
        user_a, "Bridge", ["Structural", "Geotechnical"])
    if source != fee_history.SOURCE_HISTORY:
        failures.append(f"[10] history isn't ranked above the bundled table: {source}")
    if abs(sum(split.values()) - 100) > 0.5:
        failures.append(f"[10] the applied split doesn't total 100%: {split}")
    _, source = fee_history.best_available_split(user_b, "Bridge", ["Structural", "Survey"])
    if source != fee_history.SOURCE_BUNDLED:
        failures.append(f"[10] a user with no history isn't falling back to bundled: {source}")
    split, _ = fee_history.best_available_split(user_b, "Bridge", ["Nonexistent Discipline"])
    if abs(sum(split.values()) - 100) > 0.5:
        failures.append(f"[10] an unrecognised discipline list doesn't total 100%: {split}")

    # B4: the AI call must actually receive this brief's context. Asserted
    # against the REAL prompt -- re-deriving it in the test would pass even if
    # the context stopped reaching the model.
    class _Analysis:
        project_scope = "Duplicate the existing creek bridge and realign 400 m of approach road."
        scope_items = [type("Item", (), {"title": "Bridge superstructure design", "tasks": []})()]

    captured: list = []
    estimates, error = fee_estimation_engine.refresh_estimate_from_ai(
        "Bridge", ["Structural", "Geotechnical"], "$450,000",
        config={"provider": "Anthropic Claude", "api_key": ""},
        scope_summary=_Analysis.project_scope, analysis=_Analysis(),
        capture_prompt=captured,
    )
    if not captured:
        failures.append("[10] the AI benchmark prompt was never built")
    else:
        prompt = captured[0]
        for expected in ("Duplicate the existing creek bridge", "Bridge superstructure design",
                         "Structural", "Geotechnical", "$450,000", "Australian civil engineering"):
            if expected not in prompt:
                failures.append(f"[10] the AI prompt is missing context: {expected!r}")
    # No API key configured, so the call fails -- and a failure must NOT hand
    # back the bundled table dressed as a fresh estimate.
    if estimates:
        failures.append("[10] a failed AI call still returned estimates")
    if not error or fee_estimation_engine.AI_BENCHMARK_LABEL not in error:
        failures.append(f"[10] a failed AI call didn't surface an honest error: {error!r}")

    # The label no longer claims to browse the web.
    if "web" in fee_estimation_engine.AI_BENCHMARK_LABEL.lower():
        failures.append("[10] the AI benchmark button still claims to fetch from the web")
    if hasattr(fee_estimation_engine, "refresh_estimate_from_web"):
        failures.append("[10] the old refresh_estimate_from_web name is still exported")

    # Ranges, not single numbers, once the source gives one.
    ranged = fee_estimation_engine.DisciplineFeeEstimate(
        discipline="Structural", fee_percentage=15.0, pct_low=12.0, pct_high=18.0,
        source="x", confidence="Low")
    if ranged.range_text != "12.0-18.0%":
        failures.append(f"[10] a ranged estimate doesn't render its range: {ranged.range_text}")
    point = fee_estimation_engine.DisciplineFeeEstimate(
        discipline="Structural", fee_percentage=15.0, source="x", confidence="Low")
    if point.range_text != "15.0%":
        failures.append(f"[10] a point estimate invented a range: {point.range_text}")

    try:
        with db.get_session() as session:
            session.query(db.FeeSnapshot).filter(
                db.FeeSnapshot.user_id.in_([user_a, user_b])).delete(synchronize_session=False)
            session.commit()
    except Exception:  # noqa: BLE001
        pass


def main() -> int:
    sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
    failures: list[str] = []

    test_schedule_filler_reads_real_quals(failures)
    test_blank_person_still_placeholders(failures)
    test_sender_address_is_wired_and_saved(failures)
    test_program_start_date(failures)
    test_graphics_recommendations_are_current(failures)
    test_excel_fee_exports(failures)
    test_fee_percentages_add_up(failures)
    test_firm_profile(failures)
    test_failure_honesty(failures)
    test_fee_presentation_ticks(failures)
    test_program_styles(failures)
    test_optional_design_manager(failures)
    test_org_chart_styles(failures)
    test_fee_prepopulation(failures)

    if failures:
        print("BATCH 1 WIRING TESTS FAILED:")
        for failure in failures:
            print("  -", failure)
        return 1
    print("BATCH 1 WIRING TESTS PASSED")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
