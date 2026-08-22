"""
test_exports.py -- builds a full sample project and generates every export.

This is the open-check harness for output changes: it constructs one
realistic project (analysis, structure, drafts, resourcing plan, fees,
program, references) and runs it through all six exports the app produces:

    Proposal (Large Scope) DOCX      Org Chart PPTX
    Tender Summary DOCX              Methodology Table PPTX
    Proposal (Small Scope) DOCX      Delivery Program PPTX

Then it asserts the things a human would otherwise have to open six files
to check. Every assertion here corresponds to a specific behaviour that was
wrong before, so a regression shows up as a named failure rather than as a
document someone has to notice is subtly worse.

    python test_exports.py               # assert only
    python test_exports.py --out DIR     # also write the files out for review

The sample project is deliberately hand-written rather than AI-generated:
these tests must be deterministic and must never need an API key.
"""

from __future__ import annotations

import io
import os
import sys


# ---------------------------------------------------------------------------
# The sample project
# ---------------------------------------------------------------------------

def build_sample_project() -> dict:
    """One realistic in-progress project, as a plain dict of the session-state
    values the exporters read."""
    from modules import (
        draft_generator, fee_estimation_engine, page_allocation, proposal_structure,
        reference_projects, resourcing, tender_analyser, weighting_engine,
    )

    analysis = tender_analyser.TenderAnalysis(
        project_scope=(
            "Detailed design of the replacement Example Creek bridge and 400 m of "
            "approach road, including geotechnical investigation and stormwater upgrades."
        ),
        client_objectives=[
            "Restore flood immunity to a 1% AEP event",
            "Minimise disruption to the adjoining residential street",
        ],
        submission_date="14 July 2026",
        mandatory_requirements=[
            "Professional Indemnity insurance of not less than $10,000,000",
            "RPEQ certification for all structural design",
        ],
        deliverables=[
            "Concept design report",
            "Detailed design drawings (IFC)",
            "Design certification and RPEQ sign-off",
        ],
        scope_items=[
            tender_analyser.ScopeItem(
                title="Project inception and site investigation",
                tasks=["Inception meeting", "Site inspection", "Geotechnical investigation"],
            ),
            tender_analyser.ScopeItem(
                title="Concept design",
                tasks=["Options assessment", "Concept drawings", "Cost estimate"],
            ),
            tender_analyser.ScopeItem(
                title="Detailed design",
                tasks=["Structural design", "Drainage design", "IFC drawing set"],
            ),
        ],
        risks=["Wet-season access to the creek bed", "Unknown service locations"],
        assumptions=["Survey is provided by the Principal"],
        fee_cap="$450,000 ex GST",
        disciplines_involved=["Structural", "Geotechnical", "Hydraulics & Hydrology"],
    )

    # The real Large Scope skeleton, built through the same entry point the
    # Proposal Structure tab uses, so the exported document has the shape a
    # user's would. Weighted criteria are supplied (rather than left empty)
    # because a bare skeleton is only the four fixed sections -- no Commercial
    # or Relationship Management section, and so no coverage of the parts of
    # the exporter that render them.
    weighted_criteria = [
        weighting_engine.WeightedCriterion(
            criterion_name="Methodology and program", criterion_code="SC1",
            applied_weighting=35.0, weighting_source="tender_provided",
            mapped_section="Methodology / Approach", priority_rank=1,
        ),
        weighting_engine.WeightedCriterion(
            criterion_name="Relevant experience", criterion_code="SC2",
            applied_weighting=25.0, weighting_source="tender_provided",
            mapped_section="Relevant Experience", priority_rank=2,
        ),
        weighting_engine.WeightedCriterion(
            criterion_name="Key personnel", criterion_code="SC3",
            applied_weighting=20.0, weighting_source="tender_provided",
            mapped_section="Key Personnel / Capability", priority_rank=3,
        ),
        weighting_engine.WeightedCriterion(
            criterion_name="Commercial and value for money", criterion_code="SC4",
            applied_weighting=20.0, weighting_source="tender_provided",
            mapped_section="Commercial", priority_rank=4,
        ),
    ]
    allocations = [
        page_allocation.PageAllocation(
            section_name=c.mapped_section, weighting=c.applied_weighting,
            page_limit_source="weighted_total_limit", allocated_pages=3,
            reason="Weighted share of the total page limit.",
        )
        for c in weighted_criteria
    ]
    sections = proposal_structure.build_proposal_structure(
        analysis, weighted_criteria=weighted_criteria, allocations=allocations,
        proposal_format="formal",
    )

    resource_plan = [
        resourcing.ResourceAssignment(
            slot="Project Director", slot_kind="management", person_name="Jane Smith",
            qualification="BEng (Civil) (Hons), UQ, 2003", rpeq_status="RPEQ 12345",
            years_experience="18 years",
            value_to_project="lead the design and hold RPEQ certification for the structure.",
            relevant_projects=["Example River bridge replacement, 2024"],
        ),
        resourcing.ResourceAssignment(
            slot="Structural", slot_kind="discipline", person_name="Mat Williams",
            qualification="MEng (Structural)", rpeq_status="RPEQ 22334",
            years_experience="12 years",
        ),
        # A support member with NO custom title -- must render a red placeholder,
        # never a silent "Team member".
        resourcing.ResourceAssignment(
            slot="Structural", slot_kind="discipline", person_name="Ryan Swagemakers",
            is_lead=False, years_experience="4 years",
        ),
        resourcing.ResourceAssignment(
            slot="Geotechnical", slot_kind="discipline", person_name="Sam Lee",
            qualification="BSc (Geology)", years_experience="9 years",
        ),
        # Deliberately unassigned -- exports must show this as TBC, not hide it.
        resourcing.ResourceAssignment(slot="Hydraulics & Hydrology", slot_kind="discipline"),
    ]

    discipline_fee_lines = [
        resourcing.DisciplineFeeLine(discipline="Structural", total_hours=800, rate_per_hour=210),
        resourcing.DisciplineFeeLine(discipline="Geotechnical", total_hours=260, rate_per_hour=195),
        # A priced-at-nothing row: must never export as a literal $0.
        resourcing.DisciplineFeeLine(discipline="Hydraulics & Hydrology", total_hours=0, rate_per_hour=0),
    ]

    fee_estimates = [
        fee_estimation_engine.DisciplineFeeEstimate(
            discipline="Structural", fee_percentage=58.3, confidence="Medium", source="Benchmark",
        ),
        fee_estimation_engine.DisciplineFeeEstimate(
            discipline="Geotechnical", fee_percentage=25.4, confidence="Medium", source="Benchmark",
        ),
        fee_estimation_engine.DisciplineFeeEstimate(
            discipline="Hydraulics & Hydrology", fee_percentage=16.3, confidence="Low", source="Benchmark",
        ),
    ]

    week_labels = [f"Wk {i}" for i in range(1, 13)]
    program_schedule = {
        "Project inception and site investigation": [i < 3 for i in range(12)],
        "Concept design": [2 <= i < 6 for i in range(12)],
        # Week 12 is deliberately idle: the derived cash flow must show a
        # gap there rather than smearing fee across a week with no work.
        "Detailed design": [5 <= i < 11 for i in range(12)],
    }

    drafts = {
        "Project Understanding": draft_generator.SectionDraft(
            section_title="Project Understanding",
            draft_heading="Understanding the brief",
            draft_text=(
                "**Our understanding**\n\nThe Principal requires the replacement of the "
                "Example Creek bridge together with 400 m of approach road.\n\n"
                "The controlling constraint is flood immunity to a 1% AEP event."
            ),
            required_user_inputs=[],
            recommended_graphic_placeholders=[],
        ),
        "Methodology and Deliverables": draft_generator.SectionDraft(
            section_title="Methodology and Deliverables",
            draft_heading="How we will deliver",
            draft_text=(
                "**Staged delivery**\n\nThe work is delivered in three stages: inception "
                "and investigation, concept design, and detailed design."
            ),
            required_user_inputs=[],
            recommended_graphic_placeholders=[],
        ),
    }

    refs = [
        reference_projects.ReferenceProject(
            title="Example River bridge replacement",
            client="Example Shire Council",
            description="Replacement of a 3-span bridge on a flood-affected rural road.",
            relevance_text="Same structure type and the same 1% AEP immunity requirement.",
        ),
    ]

    from modules import methodology_stages as ms
    stages = [
        ms.MethodologyStage(
            name="Project inception", week_start=1, week_end=3,
            key_tasks=["Inception meeting", "Site inspection", "Geotechnical investigation"],
            engagement_activities=["Inception meeting with the Principal"],
            outcome="Scope, program and governance confirmed.",
            deliverables=["Inception meeting minutes"],
        ),
        ms.MethodologyStage(
            name="Concept design", week_start=3, week_end=6,
            key_tasks=["Options assessment", "Concept drawings", "Cost estimate"],
            # The brief says nothing about engagement in this stage -- must
            # stay TBC, never be filled in to look complete.
            engagement_activities=["TBC"],
            outcome="Preferred option endorsed.",
            deliverables=["Concept design report"],
        ),
        ms.MethodologyStage(
            name="Detailed design", week_start=6, week_end=11,
            key_tasks=["Structural design", "Drainage design", "IFC drawing set"],
            engagement_activities=["TBC"], outcome="TBC",
            deliverables=["Detailed design drawings (IFC)", "Design certification and RPEQ sign-off"],
        ),
    ]

    return {
        "methodology_stages": stages,
        "project_info": {
            "project_name": "Example Creek Bridge Replacement",
            "client_name": "Example Shire Council",
            "tender_name": "RFT 2026-014",
            "bidder_name": "Test Engineering Pty Ltd",
            "submission_date_input": "14 July 2026",
            "proposal_theme": "Government",
            "project_type": "Bridge",
        },
        "analysis": analysis,
        "sections": sections,
        "drafts": drafts,
        "resource_plan": resource_plan,
        "discipline_fee_lines": discipline_fee_lines,
        "fee_estimates": fee_estimates,
        "program_schedule": program_schedule,
        "program_week_labels": week_labels,
        "program_start_date": __import__("datetime").date(2026, 9, 7),
        "program_style": "swimlanes",
        "org_chart_style": "cards",
        "methodology_style": "matrix",
        "reference_projects": refs,
        "sender": {
            "name": "Jane Smith", "title": "Project Director",
            "phone": "07 3000 0000", "email": "jane@example.com",
        },
        "terms_of_engagement_text": "AS 4122-2010 General Conditions of Contract.",
    }


# ---------------------------------------------------------------------------
# Generation
# ---------------------------------------------------------------------------

def generate_all(project: dict) -> dict:
    """Runs every exporter. Returns {filename: bytes}."""
    from modules import (
        export_docx, methodology_pptx, methodology_stages, org_chart, org_chart_pptx,
        program_pptx,
    )

    info = project["project_info"]
    from modules import org_chart_render

    org_png = org_chart_render.render_png(
        org_chart_render.build_model(project["resource_plan"], info["client_name"],
                                     info["project_name"], info["tender_name"]),
        project["org_chart_style"],
    ) or org_chart.render_org_chart(
        project["resource_plan"], theme_name=info["proposal_theme"],
        project_title=info["project_name"],
    )

    out: dict[str, bytes] = {}

    out["Proposal_LargeScope.docx"] = export_docx.build_docx(
        project_info=info,
        analysis=project["analysis"],
        weighted_criteria=[],
        allocations=[],
        sections=project["sections"],
        guidance_notes={},
        drafts=project["drafts"],
        compliance_items=[],
        gap_items=[],
        graphics=[],
        fee_estimates=project["fee_estimates"],
        resource_plan=project["resource_plan"],
        org_chart_png=org_png,
        reference_projects=project["reference_projects"],
        discipline_fee_lines=project["discipline_fee_lines"],
        program_schedule=project["program_schedule"],
        program_week_labels=project["program_week_labels"],
    ).getvalue()

    out["Tender_Summary.docx"] = export_docx.build_tender_summary_docx(
        project_info=info,
        analysis=project["analysis"],
        weighting_chart_png=None,
        compliance_items=[],
        gap_items=[],
        sections=project["sections"],
        drafts=project["drafts"],
    ).getvalue()

    out["Proposal_SmallScope.docx"] = export_docx.build_letter_docx(
        project_info=info,
        sender=project["sender"],
        analysis=project["analysis"],
        understanding_text=project["drafts"]["Project Understanding"].draft_text,
        methodology_text=project["drafts"]["Methodology and Deliverables"].draft_text,
        resource_plan=project["resource_plan"],
        personnel_photos={},
        program_schedule=project["program_schedule"],
        program_week_labels=project["program_week_labels"],
        terms_of_engagement_text=project["terms_of_engagement_text"],
        fee_estimates=project["fee_estimates"],
        discipline_fee_lines=project["discipline_fee_lines"],
        program_style=project["program_style"],
        methodology_stages=project["methodology_stages"],
        program_start_date=project["program_start_date"],
    ).getvalue()

    # One deck per presentation style: they share a model but not a line of
    # layout, so a regression in one is invisible in the others.
    for style in org_chart_render.STYLES:
        out[f"Org_Chart_{style}.pptx"] = org_chart_pptx.populate_org_chart(
            project["resource_plan"],
            client_name=info["client_name"],
            project_name=info["project_name"],
            tender_name=info["tender_name"],
            theme_name=info["proposal_theme"],
            style=style,
        )
    out["Org_Chart.pptx"] = out[f"Org_Chart_{project['org_chart_style']}.pptx"]

    # One deck per presentation style, same reasoning as the org chart and
    # program loops above: the four methodology styles share a data model
    # (methodology_render.build_columns) but not a line of native-pptx shape
    # code, so a regression in one style is invisible in the others.
    from modules import methodology_render

    for style in methodology_render.STYLES:
        out[f"Methodology_Table_{style}.pptx"] = methodology_pptx.populate_methodology(
            project["analysis"],
            client_name=info["client_name"],
            project_name=info["project_name"],
            theme_name=info["proposal_theme"],
            stages=project["methodology_stages"],
            week_labels=project["program_week_labels"],
            wvr_confirmed=False,
            style=style,
        )
    out["Methodology_Table.pptx"] = out[
        f"Methodology_Table_{project.get('methodology_style', methodology_render.DEFAULT_STYLE)}.pptx"
    ]
    out["Methodology_Table_no_stages.pptx"] = methodology_pptx.populate_methodology(
        project["analysis"],
        client_name=info["client_name"],
        project_name=info["project_name"],
        theme_name=info["proposal_theme"],
    )
    # One deck per presentation style: they share a model but not a single
    # line of layout, so a regression in one is invisible in the others.
    from modules import program_render

    for style in program_render.STYLES:
        out[f"Delivery_Program_{style}.pptx"] = program_pptx.populate_program(
            project["program_schedule"], project["program_week_labels"],
            client_name=info["client_name"],
            project_name=info["project_name"],
            theme_name=info["proposal_theme"],
            style=style,
            methodology_stages=project["methodology_stages"],
            start_date=project.get("program_start_date"),
            analysis=project["analysis"],
        )
    out["Delivery_Program.pptx"] = out["Delivery_Program_swimlanes.pptx"]
    if org_png:
        out["_org_chart_preview.png"] = org_png
    return out


# ---------------------------------------------------------------------------
# Reading the outputs back
# ---------------------------------------------------------------------------

def docx_text(blob: bytes) -> str:
    """All body + table text of a DOCX, as one string."""
    from docx import Document
    doc = Document(io.BytesIO(blob))
    parts = [p.text for p in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            parts.extend(c.text for c in row.cells)
    return "\n".join(parts)


def docx_image_count(blob: bytes) -> int:
    from docx import Document
    doc = Document(io.BytesIO(blob))
    return sum(1 for part in doc.part.package.parts if part.content_type.startswith("image/"))


def pptx_text(blob: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(blob))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    parts.extend(c.text for c in row.cells)
    return "\n".join(parts)


# ---------------------------------------------------------------------------
# Assertions
# ---------------------------------------------------------------------------

def check(files: dict, failures: list[str]) -> None:
    large = docx_text(files["Proposal_LargeScope.docx"])

    # 1(c): the generated org chart must actually be in the pack, and the
    # "replace me with the finished chart" note must survive alongside it.
    # (Was < 3 when the methodology grid image was also embedded here; Part C
    # removed that embed in favour of a red placeholder note, so the pack now
    # carries one fewer image -- the cover banner and the org chart.)
    if docx_image_count(files["Proposal_LargeScope.docx"]) < 2:
        failures.append("[1c] the Large Scope pack lost an image -- org chart not embedded")
    if "FIRST-PASS CHART ABOVE" not in large:
        failures.append("[1c] the org chart's replace-me note is missing")
    if "Project organisation chart" not in large:
        failures.append("[1c] the org chart heading disappeared")

    # 1(d): the cash-flow profile is derived when the fee build-up and the
    # program both exist, instead of telling the user to combine them by hand.
    if "[INSERT PROJECT CASH FLOW PROFILE" in large:
        failures.append("[1d] cash flow is still a placeholder despite a priced fee and a program")
    if "Indicative only, derived from your fee build-up and program" not in large:
        failures.append("[1d] the derived cash flow's honesty note is missing")

    from modules import export_docx
    project = build_sample_project()
    rows = export_docx.cash_flow_rows(
        project["discipline_fee_lines"], project["program_schedule"], project["program_week_labels"],
    )
    expected_total = sum(l.fee_amount for l in project["discipline_fee_lines"])
    if not rows:
        failures.append("[1d] cash_flow_rows returned nothing for a fully-specified project")
    else:
        if abs(rows[-1][2] - expected_total) > 0.01:
            failures.append(f"[1d] cash flow doesn't sum to the fee total: {rows[-1][2]} vs {expected_total}")
        # Week 12 has no scope item active in the fixture -- it must carry $0
        # rather than have money smeared into it.
        if abs(rows[-1][1]) > 0.01:
            failures.append("[1d] a week with no programmed work was given fee")

    # ...and it must NOT invent a profile when either input is missing.
    if export_docx.cash_flow_rows(project["discipline_fee_lines"], {}, []):
        failures.append("[1d] a cash flow was derived with no program")
    if export_docx.cash_flow_rows([], project["program_schedule"], project["program_week_labels"]):
        failures.append("[1d] a cash flow was derived with no priced fee")

    # 1(e): org chart PPTX -- themed, carrying the tender number, and marking
    # an untitled support member instead of calling them "Team member".
    org = pptx_text(files["Org_Chart.pptx"])
    if "RFT 2026-014" not in org:
        failures.append("[1e] the tender number never reaches the org chart title")
    if "[Project Number]" in org:
        failures.append("[1e] the org chart still shows the [Project Number] placeholder")
    if "Team member" in org:
        failures.append("[1e] an untitled support member still renders as 'Team member'")
    if "[CONFIRM TITLE]" not in org:
        failures.append("[1e] an untitled support member has no red placeholder")

    from modules import org_chart_pptx
    if org_chart_pptx._resolve_palette("Government")["header"] == org_chart_pptx._CYAN_HEADER:
        failures.append("[1e] the org chart is still hardcoded to the cyan palette")
    if org_chart_pptx._resolve_palette(None)["header"] != org_chart_pptx._CYAN_HEADER:
        failures.append("[1e] an unthemed project no longer gets the original palette")


def check_methodology_stages(failures: list[str], files: dict) -> None:
    """Batch 2: every column of the methodology table comes from the reviewed
    stage grid, TBC survives as TBC, and the WVR line is not asserted."""
    from modules import export_docx, methodology_pptx, methodology_stages

    meth = pptx_text(files["Methodology_Table.pptx"])
    for expected in ("Project inception", "Concept design", "Detailed design",
                     "Options assessment", "Design certification and RPEQ sign-off",
                     "Preferred option endorsed."):
        if expected not in meth:
            failures.append(f"[2c] the methodology table is missing real content: {expected!r}")
    if "15% design stage" in meth:
        failures.append("[2c] the hardcoded stage headers are still being used")
    if "TBC" not in meth:
        failures.append("[2a] a TBC cell was filled in rather than left as TBC")
    if "Example Creek Bridge Replacement" not in meth:
        failures.append("[2c] project_name is still not rendered on the methodology table")
    if "Wk 1 - Wk 3" not in meth:
        failures.append("[2c] the date chevrons were not filled from the program")

    # 2(e): the WVR claim must not be asserted unless confirmed.
    if methodology_pptx.WVR_STATEMENT in meth:
        failures.append("[2e] the WVR statement is asserted without confirmation")
    if methodology_pptx.WVR_CONFIRM_PLACEHOLDER not in meth:
        failures.append("[2e] no red placeholder where the WVR statement used to be")
    confirmed = pptx_text(methodology_pptx.populate_methodology(
        build_sample_project()["analysis"], stages=build_sample_project()["methodology_stages"],
        wvr_confirmed=True))
    if methodology_pptx.WVR_STATEMENT not in confirmed:
        failures.append("[2e] confirming the WVR statement doesn't bring it back")

    # Falling back with no stages must reproduce the old behaviour.
    legacy = pptx_text(files["Methodology_Table_no_stages.pptx"])
    if "15% design stage" not in legacy:
        failures.append("[2c] the no-stages fallback lost its original layout")

    # 2(d): the DOCX no longer embeds the methodology grid image -- it carries
    # exactly the red "paste the finished table here" placeholder note instead
    # (Part C: the methodology summary lives only in the PPTX export now).
    large = docx_text(files["Proposal_LargeScope.docx"])
    if "FIRST-PASS TABLE ABOVE" in large:
        failures.append("[2d] the old methodology grid image note is still in the DOCX")
    # Round 3, Part 1a: the note is now resolved from export_i18n rather than
    # a module-level English constant (removed).
    from modules import export_i18n
    if export_i18n.export_t("export_methodology_table_placeholder", "en") not in large:
        failures.append("[2d] the methodology red placeholder note is missing from the DOCX")

    # Overflow: a stage with many long tasks must stay inside the slide.
    import io as _io

    from pptx import Presentation
    fat = [methodology_stages.MethodologyStage(
        name="Detailed design",
        key_tasks=[f"A fairly long scope task description number {i} for overflow testing" for i in range(18)],
        engagement_activities=["TBC"], outcome="TBC", deliverables=["TBC"],
    )]
    prs = Presentation(_io.BytesIO(methodology_pptx.populate_methodology(
        build_sample_project()["analysis"], stages=fat)))
    lowest = max(shape.top + shape.height for shape in prs.slides[0].shapes)
    if lowest > prs.slide_height:
        failures.append("[2c] a stage with many tasks overflows the methodology slide")


def check_methodology_styles(failures: list[str], files: dict) -> None:
    """Part D: the four methodology presentation styles share one data model
    (methodology_render.build_columns) but not a line of layout code, so
    each one is checked independently for the same things -- real content
    from the reviewed grid, TBC/placeholder text staying red, no overflow
    at a heavy stage/task load -- plus the Programme style's hold-point
    diamonds coming only from a stage's own text, and the chosen style
    surviving a save/load round trip."""
    import io as _io

    from pptx import Presentation

    from modules import methodology_pptx, methodology_render, methodology_stages, project_store

    def run_colors(blob: bytes) -> list[tuple[str, str | None]]:
        prs = Presentation(_io.BytesIO(blob))
        out = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        rgb = None
                        try:
                            if run.font.color and run.font.color.type is not None:
                                rgb = str(run.font.color.rgb)
                        except Exception:  # noqa: BLE001
                            rgb = None
                        out.append((run.text, rgb))
        return out

    project = build_sample_project()

    # A heavy load, mirroring the existing "18 long tasks" overflow test but
    # exercised across all four styles: 6 stages, 8 tasks each, mixed TBC and
    # real content, a hold-point mention on one stage only.
    fat_stages = [
        methodology_stages.MethodologyStage(
            name=f"Stage {i + 1}", week_start=i * 3 + 1, week_end=i * 3 + 3,
            key_tasks=[f"Task {j + 1} for stage {i + 1} covering a fairly long scope "
                       "description for overflow testing" for j in range(8)],
            engagement_activities=(["Client workshop", "Hold point sign-off review"]
                                    if i == 1 else ["TBC"]),
            outcome="Outcome achieved." if i % 2 == 0 else "TBC",
            deliverables=[f"Deliverable {j + 1}" for j in range(6)],
        )
        for i in range(6)
    ]

    for style in methodology_render.STYLES:
        name = f"Methodology_Table_{style}.pptx"
        if name not in files:
            failures.append(f"[2f] {style} was not generated by generate_all()")
            continue

        text = pptx_text(files[name])
        if "Project inception" not in text or "Concept design" not in text:
            failures.append(f"[2f] {style} is missing real stage content from the reviewed grid")
        if "TBC" not in text:
            failures.append(f"[2f] {style} lost the TBC placeholder from the sample project")

        tbc_runs = [c for t, c in run_colors(files[name])
                    if "TBC" in t.strip().upper() or "[" in t.strip()]
        if not any(c == "C00000" for c in tbc_runs if c):
            failures.append(f"[2f] {style} doesn't render TBC/placeholder text in red")

        try:
            blob = methodology_pptx.populate_methodology(
                project["analysis"], client_name="Client", project_name="Project",
                stages=fat_stages, week_labels=[f"Wk {i + 1}" for i in range(20)], style=style)
        except Exception as exc:  # noqa: BLE001
            failures.append(f"[2f] {style} raised at 6 stages / 8 tasks: {exc}")
            continue
        prs = Presentation(_io.BytesIO(blob))
        lowest = max(shape.top + shape.height for shape in prs.slides[0].shapes)
        rightmost = max(shape.left + shape.width for shape in prs.slides[0].shapes)
        if lowest > prs.slide_height:
            failures.append(f"[2f] {style} overflows the bottom of the slide at 6 stages / 8 tasks")
        if rightmost > prs.slide_width:
            failures.append(f"[2f] {style} overflows the right edge of the slide at 6 stages / 8 tasks")

    # The Programme style's hold-point diamonds are derived from a stage's
    # own engagement/outcome text (methodology_render.stage_carries_hold_point)
    # -- never invented for a stage that doesn't mention one, never dropped
    # for a stage that does.
    with_mention = [
        methodology_stages.MethodologyStage(
            name="Stage A", key_tasks=["Task"], engagement_activities=["Kickoff"],
            outcome="Done", deliverables=["Report"]),
        methodology_stages.MethodologyStage(
            name="Stage B", key_tasks=["Task"],
            engagement_activities=["Hold point sign-off review"],
            outcome="Done", deliverables=["Report"]),
        methodology_stages.MethodologyStage(
            name="Stage C", key_tasks=["Task"], engagement_activities=["Wrap-up"],
            outcome="Done", deliverables=["Report"]),
    ]
    without_mention = [
        methodology_stages.MethodologyStage(
            name=s.name, key_tasks=s.key_tasks, engagement_activities=["Kickoff"],
            outcome=s.outcome, deliverables=s.deliverables)
        for s in with_mention
    ]
    got_holdpoint = pptx_text(methodology_pptx.populate_methodology(
        project["analysis"], stages=with_mention, style="programme"))
    no_holdpoint = pptx_text(methodology_pptx.populate_methodology(
        project["analysis"], stages=without_mention, style="programme"))
    if "HOLD" not in got_holdpoint:
        failures.append("[2f] a stage mentioning a hold point produced no hold-point diamond")
    if "HOLD" in no_holdpoint:
        failures.append("[2f] a hold-point diamond appeared despite no stage mentioning one")

    # methodology_style is a plain presentation-style key, same convention as
    # program_style/org_chart_style -- it must round-trip through a saved
    # project file untouched.
    saved = project_store.save_project({"methodology_style": "spine"})
    loaded = project_store.load_project(saved)
    if loaded.get("methodology_style") != "spine":
        failures.append(
            f"[2f] methodology_style didn't round-trip through save/load: {loaded.get('methodology_style')!r}")


def check_empty_letter_sections(failures: list[str]) -> None:
    """1(i): a Small Scope pack with nothing priced rendered a numbered
    "5. Fees" heading with nothing at all beneath it."""
    from modules import export_docx

    project = build_sample_project()
    blob = export_docx.build_letter_docx(
        project_info=project["project_info"], sender=project["sender"],
        analysis=project["analysis"], understanding_text="", methodology_text="",
        resource_plan=[], personnel_photos={}, program_schedule={}, program_week_labels=[],
        terms_of_engagement_text="", fee_estimates=None, discipline_fee_lines=None,
    ).getvalue()
    text = docx_text(blob)
    fees_index = text.find("5. Fees")
    if fees_index < 0:
        failures.append("[1i] the Small Scope pack lost its Fees section")
    else:
        between = text[fees_index + len("5. Fees"):text.find("6. Program")].strip()
        if not between:
            failures.append("[1i] '5. Fees' still renders with nothing beneath it")


def check_program_overflow(failures: list[str]) -> None:
    """1(i): beyond ~15 scope items the program's rows ran off the bottom of
    the slide, and PowerPoint shows nothing rather than complaining."""
    import io as _io

    from pptx import Presentation

    from modules import program_pptx

    for count in (3, 15, 30):
        schedule = {f"Scope item {i}": [True] * 12 for i in range(count)}
        blob = program_pptx.populate_program(
            schedule, [f"Wk {i + 1}" for i in range(12)], project_name="P")
        prs = Presentation(_io.BytesIO(blob))
        lowest = max(shape.top + shape.height for shape in prs.slides[0].shapes)
        if lowest > prs.slide_height:
            failures.append(
                f"[1i] {count} scope items overflow the program slide by "
                f"{(lowest - prs.slide_height) / 914400:.2f} in")


def check_spanish_placeholders(failures: list[str]) -> None:
    """Audit Round 2, Part 4: a Spanish-output project's placeholders must
    use the canonical Spanish markers (never a free translation), and
    collect_placeholders()'s sweep must find BOTH Spanish and English
    markers in the same document -- a pack can mix the two, since the
    returnable-schedule filler already respects output_language but the
    exporters' own hardcoded scaffolding strings are Part 5's job, not
    Part 4's."""
    from docx import Document

    from modules import export_docx, export_i18n, returnable_schedules

    es_marker = returnable_schedules.make_placeholder("ABN", "es")
    if "POR COMPLETAR" not in es_marker:
        failures.append(f"[Part4] make_placeholder(..., 'es') isn't using the canonical Spanish marker: {es_marker!r}")
    en_marker = returnable_schedules.make_placeholder("ABN", "en")
    if "TO BE COMPLETED" not in en_marker:
        failures.append(f"[Part4] make_placeholder(..., 'en') regressed the English marker: {en_marker!r}")
    default_marker = returnable_schedules.make_placeholder("ABN")
    if "TO BE COMPLETED" not in default_marker:
        failures.append(f"[Part4] make_placeholder() with no language argument should default to English: {default_marker!r}")

    doc = Document()
    doc.add_paragraph(es_marker)
    doc.add_paragraph(en_marker)
    doc.add_paragraph("[NO FEES ENTERED -- price the discipline fee build-up, or generate one]")
    found = export_docx.collect_placeholders(doc)
    if not any("POR COMPLETAR" in f for f in found):
        failures.append("[Part4] collect_placeholders() missed a Spanish [POR COMPLETAR ...] marker")
    if not any("TO BE COMPLETED" in f for f in found):
        failures.append("[Part4] collect_placeholders() missed an English [TO BE COMPLETED ...] marker")
    if not any("NO FEES ENTERED" in f for f in found):
        failures.append("[Part4] collect_placeholders() missed an existing English [NO ...] marker")

    # A Spanish-output returnable schedule (a real fill, not just the
    # marker builder in isolation) must itself carry only canonical
    # Spanish placeholders for anything it couldn't fill.
    fill_data = returnable_schedules.build_fill_data(
        {"bidder_name": "Acme Consulting", "output_language": "es"})
    if fill_data.get("output_language") != "es":
        failures.append(
            f"[Part4] build_fill_data() didn't carry output_language through: {fill_data.get('output_language')!r}")
    schedule_placeholder = returnable_schedules.make_placeholder(
        "ABN/ACN", fill_data.get("output_language"))
    if "POR COMPLETAR" not in schedule_placeholder:
        failures.append(
            f"[Part4] a Spanish-project returnable schedule wrote a non-Spanish placeholder: {schedule_placeholder!r}")

    # A generated Spanish pack must itself sweep non-empty: the placeholder
    # checklist must never silently report nothing when placeholders exist.
    project = build_sample_project()
    es_blob = export_docx.build_letter_docx(
        project_info=project["project_info"], sender=project["sender"], analysis=project["analysis"],
        understanding_text="", methodology_text="", resource_plan=[], personnel_photos={},
        program_schedule={}, program_week_labels=[], terms_of_engagement_text="",
        fee_estimates=None, discipline_fee_lines=None, output_language="es",
    ).getvalue()
    from docx import Document as _Document
    es_doc = _Document(io.BytesIO(es_blob))
    es_found = export_docx.collect_placeholders(es_doc)
    if not es_found:
        failures.append("[Part4] a Spanish-output pack with unfilled sections swept to an empty placeholder list")
    # Round 3, Part 5a: collect_placeholders() matches against
    # ALL_PLACEHOLDER_PREFIXES, the flattened set across BOTH languages (see
    # its own docstring on why -- a mid-project language switch can leave a
    # document genuinely mixed) -- so the check above passes even if every
    # single placeholder it found were an EN-prefixed leak into a Spanish
    # document, which is exactly the class of bug Part 1 fixed and this
    # entire test file exists to catch. Confirming es_found is non-empty
    # alone never proves the SPANISH markers specifically were recognised;
    # this asserts at least one match actually uses an ES-only prefix.
    es_only_prefixes = tuple(export_i18n.PLACEHOLDER_PREFIXES["es"].values())
    if not any(prefix in f.upper() for f in es_found for prefix in es_only_prefixes):
        failures.append(
            "[Part5a] a Spanish-output pack's placeholder sweep found matches, but none used an "
            f"ES-only prefix -- every match may have been an English-prefix accident: {es_found[:5]!r}")


def check_spanish_pptx(failures: list[str]) -> None:
    """Audit Round 2, Part 5 + Round 3, Part 2/5b: the three PPTX companions
    (org chart, methodology, delivery program) ignored export_i18n entirely
    -- a Spanish project got fully English PowerPoints. Round 2 fixed this
    for one non-default style each and only checked slide TITLES; Round 3
    found the actual DEFAULT styles (cards / matrix / swimlanes) were still
    untouched, plus body text (row labels, legends, badges) that title-only
    checks never exercised. This now loops every style of all three decks,
    with BOTH an empty plan/schedule (empty-state notes) and a real one
    (row/label/badge text), and asserts real translated body content -- not
    just that a title localised."""
    from pptx import Presentation

    from modules import methodology_pptx, org_chart_pptx, program_pptx, resourcing

    def all_text(blob: bytes) -> str:
        prs = Presentation(io.BytesIO(blob))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            parts.append(cell.text_frame.text)
        return "\n".join(parts)

    # A small real resourcing plan: one management row (so the client band
    # renders), a titled discipline lead, and one deliberately-untitled
    # support member (exercises the "[CONFIRMAR CARGO]" / "Líder de ..."
    # body text every style renders, not just the title).
    org_plan = [
        resourcing.ResourceAssignment(slot="Project Director", slot_kind="management",
                                      person_name="Jane Smith"),
        resourcing.ResourceAssignment(slot="Structural", slot_kind="discipline",
                                      person_name="Mat Williams", is_lead=True),
        resourcing.ResourceAssignment(slot="Structural", slot_kind="discipline",
                                      person_name="Ryan Swagemakers", is_lead=False),
    ]

    for style in ("cards", "columns", "bands", "tree"):
        org_es = all_text(org_chart_pptx.populate_org_chart(
            org_plan, client_name="Cliente Real", project_name="P", tender_name="T",
            style=style, output_language="es"))
        if "Organización del proyecto" not in org_es:
            failures.append(f"[Part2] org_chart_pptx (style={style}) title didn't localise to Spanish: {org_es[:120]!r}")
        if "CONFIRMAR CARGO" not in org_es:
            failures.append(f"[Part2] org_chart_pptx (style={style}) untitled-member body text didn't localise to Spanish")
        if "Líder de" not in org_es:
            failures.append(f"[Part2] org_chart_pptx (style={style}) discipline-lead body text didn't localise to Spanish")

        org_es_empty = all_text(org_chart_pptx.populate_org_chart(
            [], client_name="Client", project_name="P", tender_name="T",
            style=style, output_language="es"))
        if "SIN EQUIPO ASIGNADO" not in org_es_empty:
            failures.append(f"[Part2] org_chart_pptx (style={style}) empty-team note didn't localise to Spanish")

    program_schedule = {"Item A": [True, True, False, True], "Item B": [False, True, True, True]}
    for style in ("gantt", "swimlanes", "table", "timeline"):
        prog_es = all_text(program_pptx.populate_program(
            program_schedule, ["Wk 1", "Wk 2", "Wk 3", "Wk 4"],
            project_name="P", client_name="Cliente Real", style=style, output_language="es"))
        if "Programa de ejecución" not in prog_es:
            failures.append(f"[Part2] program_pptx (style={style}) title didn't localise to Spanish: {prog_es[:120]!r}")
        # The timeline style draws item labels on its bars but never a
        # separate duration string (no "X wk" text to localise there) --
        # every other style does.
        if style != "timeline" and "sem" not in prog_es.lower():
            failures.append(f"[Part2] program_pptx (style={style}) duration body text didn't localise to Spanish")

        prog_es_empty = all_text(program_pptx.populate_program(
            {}, [], project_name="P", client_name="Cliente Real", style=style, output_language="es"))
        if "SIN PROGRAMA INGRESADO" not in prog_es_empty:
            failures.append(f"[Part2] program_pptx (style={style}) empty-schedule note didn't localise to Spanish")

    for style in ("matrix", "chevrons", "programme", "spine"):
        # No stages yet -- exercises methodology_render.py's "legacy"
        # boilerplate columns (Round 3, Part 2 follow-up), the exact case
        # that made this loop fail before that fix.
        meth_es = all_text(methodology_pptx.populate_methodology(
            None, project_name="P", client_name="Cliente Real", style=style, output_language="es"))
        if "metodología propuesta" not in meth_es:
            failures.append(f"[Part2] methodology_pptx (style={style}) title didn't localise to Spanish: {meth_es[:120]!r}")
        if "TAREAS" not in meth_es.upper() and "QUÉ HACEMOS" not in meth_es.upper():
            failures.append(f"[Part2] methodology_pptx (style={style}) legacy-boilerplate body text didn't localise to Spanish")
        if not any(marker in meth_es.upper() for marker in ("ENTREGABLES", "LO QUE USTED RECIBE", "USTED RECIBE")):
            failures.append(f"[Part2] methodology_pptx (style={style}) deliverables body text didn't localise to Spanish")

    # English output must be completely unaffected by any of this threading.
    org_en = all_text(org_chart_pptx.populate_org_chart(
        [], client_name="Client", project_name="P", tender_name="T", style="bands"))
    if "Project organisation" not in org_en:
        failures.append("[Part5] org_chart_pptx default (no output_language passed) regressed away from English")


def check_spanish_png_previews(failures: list[str]) -> None:
    """Round 3, Part 4b: the live PNG previews next to the Team &
    Resourcing / Fees & Program / Draft Responses tabs
    (org_chart_render.py, program_render.py, methodology_render.py) drew
    with plain English scaffolding regardless of output_language -- only
    the PPTX companion decks (Part 2, see check_spanish_pptx() above) had
    language support. Fixed by threading the same `language: str | None =
    None` opt-in pattern through render_png() and each module's four
    per-style drawing functions, reusing the identical export_i18n keys
    Part 2 already added wherever the wording matches exactly.

    Text is read directly off the matplotlib Text artists these renderers
    draw with (fig.axes[0].texts), not OCR'd from the rendered pixels --
    exact and deterministic, unlike pixel OCR, and nothing else in this
    suite depends on tesseract being installed."""
    from modules import methodology_render, org_chart_render, program_render, resourcing
    from modules.methodology_stages import MethodologyStage

    def all_text(fig) -> str:
        return "\n".join(t.get_text() for ax in fig.axes for t in ax.texts)

    # --- org_chart_render --------------------------------------------------
    org_plan = [
        resourcing.ResourceAssignment(slot="Project Director", slot_kind="management",
                                      person_name="Jane Smith"),
        resourcing.ResourceAssignment(slot="Structural", slot_kind="discipline",
                                      person_name="Mat Williams", is_lead=True),
        resourcing.ResourceAssignment(slot="Structural", slot_kind="discipline",
                                      person_name="Ryan Swagemakers", is_lead=False),
    ]
    org_model = org_chart_render.build_model(org_plan, "Cliente Real", "P", "T", language="es")
    for style in org_chart_render.STYLES:
        png_bytes = org_chart_render.render_png(org_model, style, "#1D4ED8", language="es")
        if png_bytes is None:
            failures.append(f"[Part4b] org_chart_render.render_png (style={style}, es) failed/crashed")
            continue
        renderer = {"cards": org_chart_render._render_cards, "columns": org_chart_render._render_columns,
                   "bands": org_chart_render._render_bands, "tree": org_chart_render._render_tree}[style]
        text = all_text(renderer(org_model, "#1D4ED8", language="es"))
        if "Organización" not in text:
            failures.append(f"[Part4b] org_chart_render (style={style}) title didn't localise to Spanish: {text[:120]!r}")
        if "CONFIRMAR CARGO" not in text:
            failures.append(f"[Part4b] org_chart_render (style={style}) untitled-member body text didn't localise to Spanish")
        if "Líder de" not in text:
            failures.append(f"[Part4b] org_chart_render (style={style}) discipline-lead body text didn't localise to Spanish")

    org_model_empty = org_chart_render.build_model([], "Client", "P", "T", language="es")
    for style in org_chart_render.STYLES:
        renderer = {"cards": org_chart_render._render_cards, "columns": org_chart_render._render_columns,
                   "bands": org_chart_render._render_bands, "tree": org_chart_render._render_tree}[style]
        text = all_text(renderer(org_model_empty, "#1D4ED8", language="es"))
        if "SIN EQUIPO ASIGNADO" not in text:
            failures.append(f"[Part4b] org_chart_render (style={style}) empty-team note didn't localise to Spanish")

    # --- program_render ------------------------------------------------------
    from datetime import date
    program_schedule = {"Item A": [True, True, False, True], "Item B": [False, True, True, True]}
    prog_stages = [MethodologyStage(name="Etapa 1", week_start=1, week_end=4,
                                    engagement_activities=["Client hold point"])]
    prog_model = program_render.build_model(
        program_schedule, ["Wk 1", "Wk 2", "Wk 3", "Wk 4"], prog_stages, date(2026, 8, 1), None,
        "P", "Cliente Real", language="es")
    for style in program_render.STYLES:
        renderer = {"gantt": program_render._render_gantt, "swimlanes": program_render._render_swimlanes,
                   "table": program_render._render_table, "timeline": program_render._render_timeline}[style]
        text = all_text(renderer(prog_model, "#1D4ED8", language="es"))
        if "Programa de ejecución" not in text:
            failures.append(f"[Part4b] program_render (style={style}) title didn't localise to Spanish: {text[:120]!r}")
        if style != "timeline" and "sem" not in text.lower():
            failures.append(f"[Part4b] program_render (style={style}) duration body text didn't localise to Spanish")
    # Month bands are only drawn by the timeline style.
    tl_text = all_text(program_render._render_timeline(prog_model, "#1D4ED8", language="es"))
    if "AGOSTO" not in tl_text.upper():
        failures.append("[Part4b] program_render (style=timeline) month band didn't localise to Spanish")
    if "Punto de espera del cliente" not in tl_text and "Punto de espera del cliente" not in all_text(
            program_render._render_gantt(prog_model, "#1D4ED8", language="es")):
        failures.append("[Part4b] program_render client-hold-point milestone label didn't localise to Spanish")

    prog_model_empty = program_render.build_model({}, [], [], None, None, "P", "Cliente Real", language="es")
    for style in program_render.STYLES:
        renderer = {"gantt": program_render._render_gantt, "swimlanes": program_render._render_swimlanes,
                   "table": program_render._render_table, "timeline": program_render._render_timeline}[style]
        text = all_text(renderer(prog_model_empty, "#1D4ED8", language="es"))
        if "SIN PROGRAMA INGRESADO" not in text:
            failures.append(f"[Part4b] program_render (style={style}) empty-schedule note didn't localise to Spanish")

    # --- methodology_render ---------------------------------------------------
    for style in methodology_render.STYLES:
        renderer = {"matrix": methodology_render._render_matrix, "chevrons": methodology_render._render_chevrons,
                   "programme": methodology_render._render_programme, "spine": methodology_render._render_spine}[style]
        # No stages yet -- exercises the legacy boilerplate columns, mirroring
        # check_spanish_pptx()'s identical methodology_pptx.py check.
        cols = methodology_render.build_columns(None, None, [], "es")
        text = all_text(renderer(cols, "P", "Cliente Real", language="es"))
        if "metodología propuesta" not in text:
            failures.append(f"[Part4b] methodology_render (style={style}) title didn't localise to Spanish: {text[:120]!r}")
        if "TAREAS" not in text.upper() and "QUÉ HACEMOS" not in text.upper():
            failures.append(f"[Part4b] methodology_render (style={style}) legacy-boilerplate body text didn't localise to Spanish")
        if not any(marker in text.upper() for marker in ("ENTREGABLES", "LO QUE USTED RECIBE", "USTED RECIBE")):
            failures.append(f"[Part4b] methodology_render (style={style}) deliverables body text didn't localise to Spanish")

    # English output must be completely unaffected by any of this threading --
    # this is the un-migrated PNG-preview caller's exact call shape (no
    # `language` kwarg at all), which must still resolve to plain English via
    # the `language: str | None = None` opt-in default.
    org_en = all_text(org_chart_render._render_bands(
        org_chart_render.build_model([], "Client", "P", "T"), "#1D4ED8"))
    if "Project organisation" not in org_en:
        failures.append("[Part4b] org_chart_render default (no language passed) regressed away from English")


def check_spanish_hold_point_detection(failures: list[str]) -> None:
    """Round 3, Part 4c: methodology_render.stage_carries_hold_point() (which
    methodology_pptx.py also imports and uses directly) and
    program_render.build_model()'s milestone derivation both used to detect
    a hold point by checking a stage's own engagement/outcome text for the
    literal English phrase "hold point" -- which never matches once that
    text is legitimately in Spanish, silently dropping the Programme
    style's gate diamond and the delivery-program's milestone for every
    Spanish project. Fixed via export_i18n.mentions_hold_point(), which
    checks both languages' canonical phrase."""
    from datetime import date

    from modules import export_i18n, methodology_render, program_render
    from modules.methodology_stages import MethodologyStage

    if not methodology_render.stage_carries_hold_point(
            {"engagement": ["Punto de espera del cliente antes de continuar"], "outcome": ""}):
        failures.append("[Part4c] methodology_render.stage_carries_hold_point() missed a Spanish hold point")
    if not methodology_render.stage_carries_hold_point(
            {"engagement": ["Client hold point before proceeding"], "outcome": ""}):
        failures.append("[Part4c] methodology_render.stage_carries_hold_point() regressed on an English hold point")
    if methodology_render.stage_carries_hold_point({"engagement": ["Site inspection"], "outcome": ""}):
        failures.append("[Part4c] methodology_render.stage_carries_hold_point() false-positived with no hold point mentioned")

    stages = [MethodologyStage(name="Etapa", week_start=1, week_end=3,
                               engagement_activities=["Punto de espera del cliente"])]
    model = program_render.build_model(
        {"A": [True, True, True]}, ["Sem 1", "Sem 2", "Sem 3"], stages, date(2026, 8, 1), None,
        "P", "C", language="es")
    if not any("espera" in m.label.lower() for m in model.milestones):
        failures.append("[Part4c] program_render.build_model() missed a Spanish hold-point milestone")

    if not export_i18n.mentions_hold_point("hay un HOLD POINT aquí"):
        failures.append("[Part4c] export_i18n.mentions_hold_point() is case-sensitive and shouldn't be")


def check_spanish_resourcing_reasons(failures: list[str]) -> None:
    """Audit Round 2, Part 5: resourcing.suggest_proposal_inclusion()'s
    "reason" strings are user-visible in the Team & Resourcing tab
    (pages/60_team.py) but were English-only regardless of the project's
    output_language. Covers the two paths that don't need an AI call: the
    fixed firm-leadership reason, and the "no tender analysis yet" fallback."""
    from modules import resourcing

    class _FakeAssignment:
        def __init__(self, slot, slot_kind):
            self.slot = slot
            self.slot_kind = slot_kind

    if "Liderazgo" not in resourcing.firm_leadership_reason("es"):
        failures.append(
            f"[Part5] resourcing.firm_leadership_reason('es') didn't localise: "
            f"{resourcing.firm_leadership_reason('es')!r}")
    if resourcing.firm_leadership_reason("en") != resourcing.FIRM_LEADERSHIP_REASON:
        failures.append("[Part5] resourcing.firm_leadership_reason('en') regressed away from the English default")

    plan = [_FakeAssignment("Project Director", "management"), _FakeAssignment("Bridges", "discipline")]
    result_es = resourcing.suggest_proposal_inclusion(plan, analysis=None, output_language="es")
    if "Liderazgo" not in result_es.get("Project Director", {}).get("reason", ""):
        failures.append("[Part5] suggest_proposal_inclusion(..., output_language='es') management reason not localised")
    if "análisis" not in result_es.get("Bridges", {}).get("reason", "").lower():
        failures.append(
            "[Part5] suggest_proposal_inclusion(..., output_language='es') "
            "no-analysis-yet fallback reason not localised")

    result_en = resourcing.suggest_proposal_inclusion(plan, analysis=None)
    if result_en.get("Project Director", {}).get("reason") != resourcing.FIRM_LEADERSHIP_REASON:
        failures.append("[Part5] suggest_proposal_inclusion() default (no output_language) regressed away from English")


def check_generated_language_stale_notice(failures: list[str]) -> None:
    """Audit Round 2, Part 5's last piece: the stale-language notice.
    generated_language is a plain companion key to output_language (same
    project_store.PLAIN_KEYS category), stamped with the language a
    "Generate/Regenerate All Drafts" run actually used -- the drafting/
    export tabs (10_state_helpers.py's _generated_language_stale()) compare
    it against the project's CURRENT output_language to decide whether to
    show a non-blocking "these drafts were generated in X -- regenerate to
    get them in Y" notice. This check covers the two pieces that don't
    require the Streamlit page-script machinery those tabs run in: the
    key's save/load round trip, and the i18n notice string itself actually
    localising and substituting both language names."""
    from modules import i18n, project_store

    saved = project_store.save_project({"output_language": "es", "generated_language": "en"})
    loaded = project_store.load_project(saved)
    if loaded.get("generated_language") != "en":
        failures.append(
            f"[Part5] generated_language didn't round-trip through save/load: {loaded.get('generated_language')!r}")
    if loaded.get("output_language") != "es":
        failures.append(
            f"[Part5] output_language didn't round-trip alongside generated_language: {loaded.get('output_language')!r}")

    en_notice = i18n.t("generated_language_stale_notice", from_lang=i18n.LANGUAGES["en"], to_lang=i18n.LANGUAGES["es"])
    if "English" not in en_notice or "Español" not in en_notice:
        failures.append(f"[Part5] English stale-language notice didn't substitute both language names: {en_notice!r}")

    # i18n.t() always resolves against st.session_state["_lang"], which isn't
    # set up outside a running app -- fetch the Spanish catalog entry
    # directly instead, so this check exercises the real Spanish string
    # rather than guessing at a language-selection call convention.
    from modules.translations import es as _es_catalog
    es_template = _es_catalog.STRINGS["generated_language_stale_notice"]
    if "{from_lang}" not in es_template or "{to_lang}" not in es_template:
        failures.append(f"[Part5] Spanish stale-language notice template lost a format placeholder: {es_template!r}")
    es_filled = es_template.format(from_lang=i18n.LANGUAGES["en"], to_lang=i18n.LANGUAGES["es"])
    if "English" not in es_filled or "Español" not in es_filled:
        failures.append(f"[Part5] Spanish stale-language notice didn't substitute both language names: {es_filled!r}")


def check_swimlane_bar_colours_match_stage(failures: list[str]) -> None:
    """Audit Round 2, Part 8: program_render.py:883's draw_row() closure
    captured `colour` by reference from the enclosing lane loop (only
    `item=item` was bound as a default arg) -- since every draw_row callback
    across every lane actually runs later, inside flow.render(), by the
    time any of them ran `colour` had already advanced to the LAST lane's
    value, so every swimlane bar drew in that one colour regardless of
    which stage it belonged to, contradicting the legend right next to it.
    The PPTX twin (program_pptx.py:668) already bound correctly. Renders
    four stages' worth of bars and asserts each one's facecolor is its OWN
    stage's colour, not all four collapsed to the last stage's."""
    import matplotlib
    matplotlib.use("Agg")
    from matplotlib.patches import FancyBboxPatch
    import matplotlib.colors as mcolors

    from modules.program_render import ProgramModel, ProgramItem, STAGE_COLOURS, _render_swimlanes

    model = ProgramModel(
        items=[
            ProgramItem(label="Task A1", start_week=1, end_week=3, stage_index=0),
            ProgramItem(label="Task B1", start_week=3, end_week=5, stage_index=1),
            ProgramItem(label="Task C1", start_week=5, end_week=7, stage_index=2),
            ProgramItem(label="Task D1", start_week=6, end_week=8, stage_index=3),
        ],
        week_labels=[f"Wk {i}" for i in range(1, 10)],
        stages=["Stage A", "Stage B", "Stage C", "Stage D"],
    )
    fig = _render_swimlanes(model, "#1D4ED8")
    axes = fig.axes[0]
    bar_colours = [p.get_facecolor() for p in axes.patches if isinstance(p, FancyBboxPatch)]
    expected = [mcolors.to_rgba(STAGE_COLOURS[i]) for i in (0, 1, 2, 3)]
    if len(bar_colours) != len(expected):
        failures.append(
            f"[Part8] check_swimlane_bar_colours_match_stage: expected {len(expected)} bars, found {len(bar_colours)}")
        return
    for i, (got, want) in enumerate(zip(bar_colours, expected)):
        if got != want:
            failures.append(
                f"[Part8] check_swimlane_bar_colours_match_stage: bar {i} (stage {i}) drew in "
                f"{got} instead of its own stage's colour {want} -- the closure-capture bug is back")


def main() -> int:
    sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
    out_dir = None
    if "--out" in sys.argv:
        out_dir = sys.argv[sys.argv.index("--out") + 1]
        os.makedirs(out_dir, exist_ok=True)

    project = build_sample_project()
    files = generate_all(project)

    if out_dir:
        for name, blob in files.items():
            with open(os.path.join(out_dir, name), "wb") as fh:
                fh.write(blob)
        print(f"Wrote {len(files)} files to {out_dir}")

    failures: list[str] = []
    check(files, failures)
    check_program_overflow(failures)
    check_methodology_stages(failures, files)
    check_methodology_styles(failures, files)
    check_empty_letter_sections(failures)
    check_spanish_placeholders(failures)
    check_spanish_pptx(failures)
    check_spanish_png_previews(failures)
    check_spanish_hold_point_detection(failures)
    check_spanish_resourcing_reasons(failures)
    check_generated_language_stale_notice(failures)
    check_swimlane_bar_colours_match_stage(failures)

    if failures:
        print("EXPORT TESTS FAILED:")
        for failure in failures:
            print("  -", failure)
        return 1
    print(f"EXPORT TESTS PASSED ({len(files)} artefacts generated and checked)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
